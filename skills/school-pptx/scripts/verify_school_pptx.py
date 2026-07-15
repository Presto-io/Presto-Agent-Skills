#!/usr/bin/env python3
"""Public Phase 44 verification orchestrator for school-pptx."""

from __future__ import annotations

import argparse
import ast
import contextlib
import errno
import hashlib
import importlib
import io
import json
import os
import platform
import re
import secrets
import shutil
import stat
import subprocess
import sys
import tempfile
import xml.etree.ElementTree as ET
from zipfile import ZipFile
from datetime import datetime, timezone
from dataclasses import dataclass
from pathlib import Path
from typing import Callable


PUBLIC_OUTPUT_LIMIT = 8 * 1024
MAX_EVIDENCE_BYTES = 256 * 1024
MAX_COLLECTION_ITEMS = 64
MAX_STRING_LENGTH = 512
SCRIPT_PATH = Path(__file__).resolve()
DEFAULT_SKILL_ROOT = SCRIPT_PATH.parent.parent
VERIFY_GATE_ORDER = (
    "dependency-readiness",
    "example-generation",
    "template-validation",
    "canonical-render",
    "structural-inspection",
    "phase-43-regression",
    "negative-cases",
    "evidence-integrity",
)
VERIFY_REQUIRED_GATES = frozenset({
    "dependency-readiness",
    "example-generation",
    "template-validation",
    "canonical-render",
    "structural-inspection",
    "phase-43-regression",
    "negative-cases",
    "evidence-integrity",
})
NEGATIVE_CASE_ORDER = (
    "unknown-theme",
    "unknown-layout",
    "missing-media",
    "unsupported-styling",
    "unresolved-review-marker",
    "template-manifest-mismatch",
)
NEGATIVE_EXPECTED_CODES = {
    "unknown-theme": "THEME_UNKNOWN",
    "unknown-layout": "LAYOUT_UNKNOWN",
    "missing-media": "MEDIA_MISSING",
    "unsupported-styling": "UNSUPPORTED_STYLE",
    "unresolved-review-marker": "REVIEW_MARKER_UNRESOLVED",
    "template-manifest-mismatch": "TEMPLATE_MANIFEST_MISMATCH",
}
PHASE_43_GATE_ORDER = (
    "contract-model", "pagination", "frozen-slot-content", "frozen-numbering-row-heights",
    "ooxml-bootstrap", "editable-objects", "code-literal-roundtrip", "mixed-fragment-capacity",
    "frame-capacity-consistency", "emit-structure", "frozen-plan-emission", "cli-publication",
    "best-effort", "publication-safety", "publication-descriptor-race", "table-header-only",
    "object-error-bounded", "media-descriptor-binding", "template-reader-security", "determinism",
    "phase_41_42_regression",
)
RUN_ID_PATTERN = re.compile(r"^[a-f0-9]{16}$")
SHA256_PATTERN = re.compile(r"^[a-f0-9]{64}$")
UAT_ITEM_ORDER = (
    *(f"UAT-V{index:02d}" for index in range(1, 7)),
    *(f"UAT-I{index:02d}" for index in range(1, 5)),
)
UAT_RESULT_VALUES = frozenset({"PENDING", "PASS", "FAIL", "BLOCKED"})
UAT_VIEWERS = frozenset({"Microsoft PowerPoint", "WPS Presentation"})
REQUIRED_DEPENDENCIES = (
    ("python-pptx", "pptx"),
    ("Pillow", "PIL"),
    ("lxml", "lxml"),
    ("PyYAML", "yaml"),
)
STABLE_REMEDIATIONS = {
    "VERIFY_ARGUMENT_INVALID": "Use exactly: verify --workdir <directory>.",
    "VERIFY_WORKDIR_UNSAFE": "Choose a real, writable directory that is not a symlink.",
    "VERIFY_NOT_IMPLEMENTED": "Complete the remaining verification gates before accepting this run.",
    "VERIFY_EVIDENCE_INVALID": "Inspect the bounded current evidence and correct the verification implementation.",
    "VERIFY_DEPENDENCY_BLOCKED": "Restore required dependencies, then run verification again.",
    "VERIFY_DOCUMENTATION_INVALID": "Restore the canonical skill, references, runtime rows, and repository discovery links.",
}


class VerifyFailure(RuntimeError):
    def __init__(self, code: str, remediation: str | None = None) -> None:
        super().__init__(code)
        self.code = code
        self.remediation = remediation or STABLE_REMEDIATIONS[code]


class UATLintFailure(RuntimeError):
    def __init__(self, code: str) -> None:
        super().__init__(code)
        self.code = code


@dataclass
class HeldWorkdir:
    path: Path
    fd: int

    def close(self) -> None:
        if self.fd >= 0:
            os.close(self.fd)
            self.fd = -1

    def __enter__(self) -> "HeldWorkdir":
        return self

    def __exit__(self, exc_type: object, exc: object, traceback: object) -> None:
        del exc_type, exc, traceback
        self.close()


@dataclass
class OwnedTree:
    root: HeldWorkdir
    delivery_fd: int
    evidence_fd: int
    work_fd: int
    run_fd: int
    run_id: str

    def close(self) -> None:
        for attribute in ("run_fd", "work_fd", "evidence_fd", "delivery_fd"):
            descriptor = getattr(self, attribute)
            if descriptor >= 0:
                os.close(descriptor)
                setattr(self, attribute, -1)

    def __enter__(self) -> "OwnedTree":
        return self

    def __exit__(self, exc_type: object, exc: object, traceback: object) -> None:
        del exc_type, exc, traceback
        self.close()


def bounded_text(value: object, limit: int = 240) -> str:
    text = " ".join(str(value).split())
    return text[:limit]


def emit_failure(failure: VerifyFailure) -> int:
    message = f"{failure.code}: {bounded_text(failure.remediation)}\n"
    encoded = message.encode("utf-8")[: PUBLIC_OUTPUT_LIMIT - 1]
    sys.stderr.buffer.write(encoded + (b"\n" if not encoded.endswith(b"\n") else b""))
    return 2


def _split_internal_skill_root(argv: list[str]) -> tuple[Path, list[str]]:
    if argv and not argv[0].startswith("-"):
        supplied = Path(argv[0])
        try:
            if supplied.resolve(strict=True) != DEFAULT_SKILL_ROOT.resolve(strict=True):
                raise VerifyFailure("VERIFY_ARGUMENT_INVALID")
        except (OSError, RuntimeError) as exc:
            raise VerifyFailure("VERIFY_ARGUMENT_INVALID") from exc
        return supplied, argv[1:]
    return DEFAULT_SKILL_ROOT, argv


def parse_public_args(argv: list[str]) -> tuple[Path, Path]:
    skill_root, public = _split_internal_skill_root(argv)
    if len(public) != 2 or public[0] != "--workdir" or not public[1]:
        raise VerifyFailure("VERIFY_ARGUMENT_INVALID")
    if public[1].startswith("-"):
        raise VerifyFailure("VERIFY_ARGUMENT_INVALID")
    return skill_root, Path(public[1])


def _open_directory(path: Path) -> int:
    flags = os.O_RDONLY | getattr(os, "O_DIRECTORY", 0) | getattr(os, "O_NOFOLLOW", 0)
    return os.open(path, flags)


def _open_directory_at(parent_fd: int, name: str) -> int:
    flags = os.O_RDONLY | getattr(os, "O_DIRECTORY", 0) | getattr(os, "O_NOFOLLOW", 0)
    return os.open(name, flags, dir_fd=parent_fd)


def _ensure_directory_at(parent_fd: int, name: str) -> int:
    if not name or name in {".", ".."} or "/" in name or "\x00" in name:
        raise VerifyFailure("VERIFY_WORKDIR_UNSAFE")
    try:
        os.mkdir(name, mode=0o700, dir_fd=parent_fd)
    except FileExistsError:
        pass
    except OSError as exc:
        raise VerifyFailure("VERIFY_WORKDIR_UNSAFE") from exc
    try:
        descriptor = _open_directory_at(parent_fd, name)
        details = os.stat(name, dir_fd=parent_fd, follow_symlinks=False)
        held = os.fstat(descriptor)
        if not stat.S_ISDIR(details.st_mode) or (details.st_dev, details.st_ino) != (held.st_dev, held.st_ino):
            raise VerifyFailure("VERIFY_WORKDIR_UNSAFE")
        return descriptor
    except (OSError, VerifyFailure) as exc:
        if "descriptor" in locals():
            os.close(descriptor)
        if isinstance(exc, VerifyFailure):
            raise
        raise VerifyFailure("VERIFY_WORKDIR_UNSAFE") from exc


def hold_workdir(path: Path) -> HeldWorkdir:
    if path == Path(""):
        raise VerifyFailure("VERIFY_ARGUMENT_INVALID")
    try:
        details = os.lstat(path)
    except FileNotFoundError:
        parent = path.parent if path.parent != Path("") else Path(".")
        try:
            parent_fd = _open_directory(parent)
        except OSError as exc:
            raise VerifyFailure("VERIFY_WORKDIR_UNSAFE") from exc
        try:
            os.mkdir(path.name, mode=0o700, dir_fd=parent_fd)
        except OSError as exc:
            raise VerifyFailure("VERIFY_WORKDIR_UNSAFE") from exc
        finally:
            os.close(parent_fd)
        details = os.lstat(path)
    except OSError as exc:
        raise VerifyFailure("VERIFY_WORKDIR_UNSAFE") from exc
    if stat.S_ISLNK(details.st_mode) or not stat.S_ISDIR(details.st_mode):
        raise VerifyFailure("VERIFY_WORKDIR_UNSAFE")
    try:
        fd = _open_directory(path)
        held = os.fstat(fd)
        current = os.lstat(path)
        if (held.st_dev, held.st_ino) != (current.st_dev, current.st_ino):
            raise VerifyFailure("VERIFY_WORKDIR_UNSAFE")
        if not os.access(path, os.W_OK | os.X_OK):
            raise VerifyFailure("VERIFY_WORKDIR_UNSAFE")
    except (OSError, VerifyFailure) as exc:
        if "fd" in locals():
            os.close(fd)
        if isinstance(exc, VerifyFailure):
            raise
        if getattr(exc, "errno", None) in {errno.ELOOP, errno.ENOTDIR, errno.EACCES}:
            raise VerifyFailure("VERIFY_WORKDIR_UNSAFE") from exc
        raise VerifyFailure("VERIFY_WORKDIR_UNSAFE") from exc
    return HeldWorkdir(path=path, fd=fd)


def build_owned_tree(root: HeldWorkdir, run_id: str | None = None) -> OwnedTree:
    bounded_run_id = run_id or secrets.token_hex(8)
    if not RUN_ID_PATTERN.fullmatch(bounded_run_id):
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    opened: list[int] = []
    try:
        delivery_fd = _ensure_directory_at(root.fd, "delivery")
        opened.append(delivery_fd)
        evidence_fd = _ensure_directory_at(root.fd, "evidence")
        opened.append(evidence_fd)
        work_fd = _ensure_directory_at(root.fd, "work")
        opened.append(work_fd)
        run_fd = _ensure_directory_at(work_fd, bounded_run_id)
        opened.append(run_fd)
        return OwnedTree(root, delivery_fd, evidence_fd, work_fd, run_fd, bounded_run_id)
    except Exception:
        for descriptor in reversed(opened):
            os.close(descriptor)
        raise


def iso_now() -> str:
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat().replace("+00:00", "Z")


def gate_result(status_value: str, code: str, remediation: str) -> dict[str, object]:
    if status_value not in {"passed", "failed"}:
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    return {
        "status": status_value,
        "code": bounded_text(code, 64),
        "remediation": bounded_text(remediation, 240),
    }


def not_implemented_result(name: str) -> dict[str, object]:
    return gate_result(
        "failed",
        "NOT_IMPLEMENTED",
        f"Gate {name} is reserved for Phase 44 Plan 02 and is not accepted yet.",
    )


def blocked_result() -> dict[str, object]:
    return gate_result("failed", "VERIFY_DEPENDENCY_BLOCKED", STABLE_REMEDIATIONS["VERIFY_DEPENDENCY_BLOCKED"])


def dependency_readiness_placeholder() -> dict[str, object]:
    return dependency_readiness_gate()


def _module_version(module: object) -> str:
    value = getattr(module, "__version__", None)
    if value is None:
        value = getattr(module, "VERSION", None)
    if isinstance(value, (tuple, list)):
        value = ".".join(str(item) for item in value)
    return bounded_text(value if value is not None else "unknown", 80)


def _pandoc_probe() -> dict[str, object]:
    executable = shutil.which("pandoc")
    if executable is None:
        return {"required": False, "available": False, "version": None}
    try:
        completed = subprocess.run(
            [executable, "--version"],
            stdin=subprocess.DEVNULL,
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
            text=True,
            timeout=2,
            check=False,
        )
        first_line = completed.stdout.splitlines()[0] if completed.stdout.splitlines() else ""
        version = bounded_text(first_line, 120) if completed.returncode == 0 and first_line else None
        return {"required": False, "available": completed.returncode == 0, "version": version}
    except (OSError, subprocess.SubprocessError):
        return {"required": False, "available": False, "version": None}


def observe_dependencies(
    importer: Callable[[str], object] = importlib.import_module,
    pandoc_probe: Callable[[], dict[str, object]] = _pandoc_probe,
) -> tuple[dict[str, object], list[str]]:
    selected = "SCHOOL_PPTX_PYTHON" if os.environ.get("SCHOOL_PPTX_PYTHON") else "default"
    executable_basename = bounded_text(Path(sys.executable).name, 80)
    python_observation = {
        "selection_source": selected,
        "executable_basename": executable_basename,
        "executable_identity": f"<external>/{executable_basename}",
        "implementation": bounded_text(platform.python_implementation(), 40),
        "version": bounded_text(platform.python_version(), 40),
    }
    packages: list[dict[str, object]] = []
    missing: list[str] = []
    for distribution, import_name in REQUIRED_DEPENDENCIES:
        try:
            module = importer(import_name)
            error = None
            version = _module_version(module)
        except (ImportError, ModuleNotFoundError):
            error = "missing"
            version = None
            missing.append(distribution)
        packages.append({
            "distribution": distribution,
            "import_name": import_name,
            "required": True,
            "ready": error is None,
            "version": version,
        })
    pandoc = pandoc_probe()
    if set(pandoc) != {"required", "available", "version"} or pandoc.get("required") is not False:
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    evidence = {"python": python_observation, "packages": packages, "pandoc": pandoc}
    _walk_bounded(evidence)
    return evidence, missing


def dependency_readiness_gate(
    importer: Callable[[str], object] = importlib.import_module,
    pandoc_probe: Callable[[], dict[str, object]] = _pandoc_probe,
) -> dict[str, object]:
    observed, missing = observe_dependencies(importer, pandoc_probe)
    if missing:
        result = gate_result(
            "failed",
            "VERIFY_DEPENDENCY_MISSING",
            "Install the required school-pptx Python packages in the selected interpreter.",
        )
    else:
        result = gate_result("passed", "OK", "All required Python imports were observed in this run.")
    result["observed"] = observed
    return result


def _relative_artifact_path(value: str) -> bool:
    candidate = Path(value)
    return bool(value) and not candidate.is_absolute() and ".." not in candidate.parts and "\x00" not in value


def _walk_bounded(value: object, *, key: str = "", depth: int = 0) -> None:
    if depth > 8:
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    if isinstance(value, dict):
        if len(value) > MAX_COLLECTION_ITEMS:
            raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
        for child_key, child in value.items():
            if not isinstance(child_key, str) or len(child_key) > 80:
                raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
            _walk_bounded(child, key=child_key, depth=depth + 1)
        return
    if isinstance(value, (list, tuple)):
        if len(value) > MAX_COLLECTION_ITEMS:
            raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
        for child in value:
            _walk_bounded(child, key=key, depth=depth + 1)
        return
    if isinstance(value, str):
        if len(value) > MAX_STRING_LENGTH or "Traceback" in value or "-----BEGIN" in value:
            raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
        if key in {"path", "evidence_path"} and not _relative_artifact_path(value):
            raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
        if key == "sha256" and not SHA256_PATTERN.fullmatch(value):
            raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
        return
    if value is not None and not isinstance(value, (bool, int, float)):
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")


def validate_candidate(candidate: dict[str, object], *, include_integrity: bool) -> None:
    if candidate.get("schema_version") != 1:
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    registry = candidate.get("registry")
    gates = candidate.get("gates")
    if not isinstance(registry, dict) or not isinstance(gates, dict):
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    required = tuple(registry.get("required", ()))
    called = tuple(registry.get("called", ()))
    if len(VERIFY_GATE_ORDER) != len(set(VERIFY_GATE_ORDER)):
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    if set(VERIFY_GATE_ORDER) != VERIFY_REQUIRED_GATES:
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    if required != VERIFY_GATE_ORDER or len(required) != len(set(required)):
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    expected_called = VERIFY_GATE_ORDER if include_integrity else VERIFY_GATE_ORDER[:-1]
    if called != expected_called or len(called) != len(set(called)):
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    if registry.get("dynamic_skips") != 0 or registry.get("unique") is not True:
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    if set(gates) != set(expected_called):
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    for name in expected_called:
        result = gates.get(name)
        if not isinstance(result, dict) or result.get("status") not in {"passed", "failed"}:
            raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    if include_integrity and candidate.get("status") != _aggregate_status(gates, list(called)):
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    phase_43 = gates.get("phase-43-regression")
    if isinstance(phase_43, dict) and phase_43.get("status") == "passed":
        phase_evidence = phase_43.get("evidence")
        if not isinstance(phase_evidence, dict):
            raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
        if tuple(phase_evidence.get("required", ())) != PHASE_43_GATE_ORDER:
            raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
        if tuple(phase_evidence.get("called", ())) != PHASE_43_GATE_ORDER:
            raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
        if phase_evidence.get("unique") is not True or phase_evidence.get("dynamic_skips") != 0:
            raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    negative = gates.get("negative-cases")
    if isinstance(negative, dict) and negative.get("status") == "passed":
        evidence = negative.get("evidence")
        if not isinstance(evidence, dict) or not isinstance(evidence.get("registry"), dict):
            raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
        negative_registry = evidence["registry"]
        if tuple(negative_registry.get("required", ())) != NEGATIVE_CASE_ORDER:
            raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
        if tuple(negative_registry.get("called", ())) != NEGATIVE_CASE_ORDER:
            raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
        if negative_registry.get("unique") is not True or negative_registry.get("dynamic_skips") != 0:
            raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
        cases = evidence.get("cases")
        if not isinstance(cases, list) or [item.get("name") for item in cases] != list(NEGATIVE_CASE_ORDER):
            raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
        if any(item.get("status") != "failed-as-expected" for item in cases):
            raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    structure = gates.get("structural-inspection")
    if isinstance(structure, dict) and structure.get("status") == "passed":
        evidence = structure.get("evidence")
        if not isinstance(evidence, dict) or evidence.get("transition_mode") not in {"none", "preserved", "generated"}:
            raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
        if evidence.get("logical_count") != len(evidence.get("logical_to_physical", [])):
            raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
        if evidence.get("physical_count") != len(evidence.get("layout_parts", [])):
            raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
        mapped = [index for item in evidence["logical_to_physical"] for index in item["physical_indices"]]
        if mapped != list(range(1, evidence["physical_count"])):
            raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
        contents = evidence.get("contents", {})
        notes = evidence.get("notes", {})
        if contents.get("source") != contents.get("reopen") or contents.get("source_count") != contents.get("reopen_count"):
            raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
        if notes.get("expected_indices") != notes.get("observed_indices") or notes.get("expected_count") != notes.get("observed_count"):
            raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    _walk_bounded(candidate)
    encoded = json.dumps(candidate, ensure_ascii=False, sort_keys=True).encode("utf-8")
    if len(encoded) > MAX_EVIDENCE_BYTES:
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")


def _aggregate_status(gates: dict[str, dict[str, object]], called: list[str]) -> str:
    complete = tuple(called) == VERIFY_GATE_ORDER
    all_passed = complete and all(gates[name]["status"] == "passed" for name in VERIFY_GATE_ORDER)
    return "passed" if all_passed else "failed"


def markdown_from_candidate(candidate: dict[str, object]) -> str:
    registry = candidate["registry"]
    gates = candidate["gates"]
    lines = [
        "# School PPTX Verification",
        "",
        f"Status: {str(candidate['status']).upper()}",
        f"Run: {candidate['run']['id']}",
        f"Gates: {len(registry['called'])}/{len(registry['required'])}",
        f"Dynamic skips: {registry['dynamic_skips']}",
        "",
        "| Gate | Status | Code |",
        "|---|---|---|",
    ]
    for name in registry["called"]:
        result = gates[name]
        lines.append(f"| {name} | {result['status']} | {result['code']} |")
    lines.extend(("", "Generated from the same in-memory observed evidence as verification.json.", ""))
    return "\n".join(lines)


def _reject_unsafe_target(directory_fd: int, name: str) -> None:
    try:
        details = os.stat(name, dir_fd=directory_fd, follow_symlinks=False)
    except FileNotFoundError:
        return
    except OSError as exc:
        raise VerifyFailure("VERIFY_WORKDIR_UNSAFE") from exc
    if not stat.S_ISREG(details.st_mode):
        raise VerifyFailure("VERIFY_WORKDIR_UNSAFE")


def _atomic_publish(directory_fd: int, name: str, payload: bytes, run_id: str) -> None:
    _reject_unsafe_target(directory_fd, name)
    temporary = f".{name}.{run_id}.tmp"
    flags = os.O_WRONLY | os.O_CREAT | os.O_EXCL | getattr(os, "O_NOFOLLOW", 0)
    descriptor = -1
    try:
        descriptor = os.open(temporary, flags, 0o600, dir_fd=directory_fd)
        written = 0
        while written < len(payload):
            written += os.write(descriptor, payload[written:])
        os.fsync(descriptor)
        details = os.fstat(descriptor)
        current = os.stat(temporary, dir_fd=directory_fd, follow_symlinks=False)
        if not stat.S_ISREG(current.st_mode) or (details.st_dev, details.st_ino) != (current.st_dev, current.st_ino):
            raise VerifyFailure("VERIFY_WORKDIR_UNSAFE")
        os.close(descriptor)
        descriptor = -1
        _reject_unsafe_target(directory_fd, name)
        os.replace(temporary, name, src_dir_fd=directory_fd, dst_dir_fd=directory_fd)
        os.fsync(directory_fd)
    except OSError as exc:
        raise VerifyFailure("VERIFY_WORKDIR_UNSAFE") from exc
    finally:
        if descriptor >= 0:
            os.close(descriptor)
        try:
            os.unlink(temporary, dir_fd=directory_fd)
        except FileNotFoundError:
            pass


def publish_evidence(tree: OwnedTree, candidate: dict[str, object]) -> None:
    validate_candidate(candidate, include_integrity=True)
    json_payload = (json.dumps(candidate, ensure_ascii=False, indent=2, sort_keys=True) + "\n").encode("utf-8")
    markdown_payload = markdown_from_candidate(candidate).encode("utf-8")
    if max(len(json_payload), len(markdown_payload)) > MAX_EVIDENCE_BYTES:
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    _atomic_publish(tree.evidence_fd, "verification.json", json_payload, tree.run_id)
    _atomic_publish(tree.evidence_fd, "verification.md", markdown_payload, tree.run_id)


def sha256_path(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(64 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def run_bounded(command: list[str], *, cwd: Path, timeout: int = 180) -> dict[str, object]:
    try:
        completed = subprocess.run(
            command,
            cwd=cwd,
            stdin=subprocess.DEVNULL,
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
            timeout=timeout,
            check=False,
        )
        captured = completed.stdout[:PUBLIC_OUTPUT_LIMIT]
        overflow = len(completed.stdout) >= PUBLIC_OUTPUT_LIMIT
        return {
            "exit_code": completed.returncode,
            "capture_bytes": len(completed.stdout),
            "bounded": not overflow,
            "traceback": b"Traceback" in captured,
            "text": captured.decode("utf-8", errors="replace"),
        }
    except (OSError, subprocess.SubprocessError) as exc:
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID", bounded_text(exc)) from exc


def require_command_ok(result: dict[str, object], code: str) -> None:
    if result["exit_code"] != 0 or not result["bounded"] or result["traceback"]:
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID", code)


def public_cli(skill_root: Path) -> str:
    return str(skill_root / "scripts" / "school-pptx.sh")


def canonical_flow(tree: OwnedTree, skill_root: Path) -> tuple[dict[str, dict[str, object]], dict[str, object]]:
    run_root = tree.root.path / "work" / tree.run_id
    example_root = run_root / "canonical-example"
    report_root = run_root / "template-report"
    example_root.mkdir(mode=0o700)
    report_root.mkdir(mode=0o700)
    example = run_bounded([public_cli(skill_root), "example", "--out-dir", "canonical-example"], cwd=run_root)
    require_command_ok(example, "Public example generation failed.")
    source = example_root / "school-pptx-full.md"
    expected_example = {"school-pptx-full.md", "media"}
    if set(path.name for path in example_root.iterdir()) != expected_example:
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID", "Example inventory changed.")
    example_gate = gate_result("passed", "OK", "Fresh public example generation completed.")
    example_gate["evidence"] = {
        "command": "example",
        "exit_code": example["exit_code"],
        "capture_bytes": example["capture_bytes"],
        "markdown_sha256": sha256_path(source),
        "media_count": len(list((example_root / "media").iterdir())),
    }

    report_json = report_root / "template-report.json"
    report_md = report_root / "template-report.md"
    template = run_bounded([
        public_cli(skill_root), "template-report", "--theme", "standard-school",
        "--out-md", "template-report/template-report.md",
        "--out-json", "template-report/template-report.json",
    ], cwd=run_root)
    require_command_ok(template, "Canonical template validation failed.")
    report = json.loads(report_json.read_text(encoding="utf-8"))
    if report.get("failures") != [] or not report_md.is_file():
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID", "Canonical template report contains failures.")
    template_gate = gate_result("passed", "OK", "Canonical template and manifest agree.")
    template_gate["evidence"] = {
        "command": "template-report",
        "exit_code": template["exit_code"],
        "capture_bytes": template["capture_bytes"],
        "failure_count": len(report["failures"]),
        "template_sha256": sha256_path(skill_root / "templates" / "standard-school.pptx"),
        "manifest_sha256": sha256_path(skill_root / "templates" / "standard-school.manifest.yaml"),
    }

    delivery_root = tree.root.path / "delivery"
    existing = {path.name for path in delivery_root.iterdir()}
    if existing - {"canonical.md", "canonical.pptx"}:
        raise VerifyFailure("VERIFY_WORKDIR_UNSAFE", "Delivery contains caller-owned files.")
    render = run_bounded([
        public_cli(skill_root), "render", "--input", "canonical-example/school-pptx-full.md",
        "--out-dir", str(delivery_root), "--stem", "canonical",
    ], cwd=run_root)
    require_command_ok(render, "Canonical public render failed.")
    markdown = delivery_root / "canonical.md"
    pptx = delivery_root / "canonical.pptx"
    if {path.name for path in delivery_root.iterdir()} != {"canonical.md", "canonical.pptx"}:
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID", "Delivery is not the exact same-stem pair.")
    if markdown.read_bytes() != source.read_bytes() or pptx.stat().st_size <= 0:
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID", "Fresh delivery bytes do not match the public source.")
    render_gate = gate_result("passed", "OK", "Fresh public render published one same-stem pair.")
    render_gate["evidence"] = {
        "command": "render", "exit_code": render["exit_code"], "capture_bytes": render["capture_bytes"],
        "markdown_sha256": sha256_path(markdown), "pptx_sha256": sha256_path(pptx),
        "pptx_bytes": pptx.stat().st_size, "delivery": ["delivery/canonical.md", "delivery/canonical.pptx"],
    }
    return {
        "example-generation": example_gate,
        "template-validation": template_gate,
        "canonical-render": render_gate,
    }, {"source": source, "markdown": markdown, "pptx": pptx}


def _slide_texts(slide: object) -> list[str]:
    values: list[str] = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            values.append(shape.text.strip())
    return values


def structural_inspection(artifacts: dict[str, object]) -> dict[str, object]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from markdown_contract import load_manifest, parse_document

    source = artifacts["source"]
    pptx_path = artifacts["pptx"]
    document = parse_document(source, load_manifest(DEFAULT_SKILL_ROOT))
    if document["errors"]:
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID", "Canonical source failed reopen parsing.")
    presentation = Presentation(str(pptx_path))
    slides = list(presentation.slides)
    layout_parts = [slide.slide_layout.part.partname.lstrip("/") for slide in slides]
    slide_texts = [_slide_texts(slide) for slide in slides]
    logical = document["logical_slides"]
    cursor = 0
    mappings: list[dict[str, object]] = []
    manifest = load_manifest(DEFAULT_SKILL_ROOT)
    for logical_index, item in enumerate(logical, 1):
        layout = item["layout"]
        expected_part = manifest["layouts"][layout]["pptx_layout"]
        expected_title = document["metadata"].get("title") if layout == "cover" else "目录" if layout == "contents" else item["title"]
        indices: list[int] = []
        while cursor < len(slides) - 1:
            first = slide_texts[cursor][0] if slide_texts[cursor] else ""
            if layout_parts[cursor] != expected_part or (layout != "contents" and first != expected_title):
                break
            indices.append(cursor + 1)
            cursor += 1
        mappings.append({"logical": logical_index, "layout": layout, "physical_indices": indices})
    source_contents = list(document["contents_entries"])
    content_indices = mappings[1]["physical_indices"]
    reopen_contents: list[str] = []
    for index in content_indices:
        body = slide_texts[index - 1][0] if slide_texts[index - 1] else ""
        reopen_contents.extend(re.sub(r"^\d+\.\s*", "", line) for line in body.splitlines() if line.strip())
    expected_notes = list(range(1, len(slides) + 1))
    observed_notes = [
        index for index, slide in enumerate(slides, 1)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip()
    ]
    media: list[dict[str, object]] = []
    tables: list[dict[str, object]] = []
    groups: list[dict[str, object]] = []
    whole_slide_picture_count = 0
    for page, slide in enumerate(slides, 1):
        page_groups = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                blob = shape.image.blob
                crops = [shape.crop_left, shape.crop_top, shape.crop_right, shape.crop_bottom]
                media.append({"page": page, "sha256": hashlib.sha256(blob).hexdigest(), "crop": crops})
                if shape.left <= 0 and shape.top <= 0 and shape.width >= presentation.slide_width and shape.height >= presentation.slide_height:
                    whole_slide_picture_count += 1
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                page_groups += 1
            if getattr(shape, "has_table", False):
                table = shape.table
                tables.append({
                    "page": page, "rows": len(table.rows), "columns": len(table.columns),
                    "header": [cell.text for cell in table.rows[0].cells],
                })
        if page_groups:
            groups.append({"page": page, "count": page_groups})
    code_source = "\n".join(
        block["text"] for item in logical for block in item["blocks"] if block["kind"] == "code"
    )
    code_pages = [
        index for mapping in mappings if mapping["layout"] == "code" for index in mapping["physical_indices"]
    ]
    code_frames: list[str] = []
    code_run_properties: list[dict[str, object]] = []
    for page in code_pages:
        text_shapes = [shape for shape in slides[page - 1].shapes if hasattr(shape, "text") and shape.text.strip()]
        for shape in text_shapes[1:]:
            code_frames.append(shape.text)
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    code_run_properties.append({
                        "page": page,
                        "size": run.font.size.pt if run.font.size is not None else None,
                        "font": run.font.name,
                    })
    code_reopen = "\n".join(code_frames)
    with ZipFile(pptx_path) as package:
        transition_count = 0
        embedded_hashes = []
        for name in package.namelist():
            if re.fullmatch(r"ppt/slides/slide\d+\.xml", name):
                transition_count += len(ET.fromstring(package.read(name)).findall("{http://schemas.openxmlformats.org/presentationml/2006/main}transition"))
            if name.startswith("ppt/media/") and not name.endswith("/"):
                embedded_hashes.append(hashlib.sha256(package.read(name)).hexdigest())
    transition_mode = "preserved" if transition_count else "none"
    evidence = {
        "artifact": {"path": "delivery/canonical.pptx", "sha256": sha256_path(pptx_path), "bytes": pptx_path.stat().st_size},
        "logical_count": len(logical), "physical_count": len(slides), "layout_parts": layout_parts,
        "logical_to_physical": mappings,
        "contents": {"source": source_contents, "reopen": reopen_contents, "source_count": len(source_contents), "reopen_count": len(reopen_contents)},
        "notes": {"expected_indices": expected_notes, "observed_indices": observed_notes, "expected_count": len(expected_notes), "observed_count": len(observed_notes)},
        "media": {"pictures": media, "embedded_hashes": sorted(embedded_hashes), "picture_count": len(media)},
        "tables": tables,
        "code": {
            "source_sha256": hashlib.sha256(code_source.encode()).hexdigest(),
            "reopen_sha256": hashlib.sha256(code_reopen.encode()).hexdigest(),
            "pages": code_pages, "run_properties": code_run_properties,
        },
        "groups": groups, "whole_slide_picture_count": whole_slide_picture_count,
        "transition_mode": transition_mode,
    }
    valid = (
        evidence["artifact"]["sha256"] == sha256_path(pptx_path)
        and evidence["logical_count"] == 13 and evidence["physical_count"] == 26
        and cursor == len(slides) - 1 and all(item["physical_indices"] for item in mappings)
        and source_contents == reopen_contents and expected_notes == observed_notes
        and bool(media) and bool(embedded_hashes) and bool(tables) and bool(code_pages)
        and evidence["code"]["source_sha256"] == evidence["code"]["reopen_sha256"]
        and all(item["crop"] == [0.0, 0.0, 0.0, 0.0] for item in media)
        and whole_slide_picture_count == 0 and transition_mode in {"none", "preserved", "generated"}
    )
    result = gate_result("passed" if valid else "failed", "OK" if valid else "STRUCTURE_MISMATCH", "Fresh PPTX raw structure was independently reopened and compared.")
    result["evidence"] = evidence
    return result


def phase_43_regression() -> dict[str, object]:
    import verify_pptx_renderer

    capture = io.StringIO()
    with contextlib.redirect_stdout(capture):
        evidence = verify_pptx_renderer.run_phase_43()
    registry = evidence.get("registry", {})
    required = tuple(registry.get("required", ()))
    called = tuple(registry.get("called", ()))
    valid = (
        required == PHASE_43_GATE_ORDER == called
        and len(called) == len(set(called))
        and registry.get("unique") is True
        and registry.get("dynamic_skips") == 0
        and len(capture.getvalue().encode()) < PUBLIC_OUTPUT_LIMIT
    )
    result = gate_result("passed" if valid else "failed", "OK" if valid else "PHASE_43_REGRESSION_FAILED", "Phase 43 fixed registry was executed and revalidated.")
    result["evidence"] = {"required": list(required), "called": list(called), "unique": len(called) == len(set(called)), "dynamic_skips": registry.get("dynamic_skips"), "capture_bytes": len(capture.getvalue().encode())}
    return result


def _extract_codes(text: str) -> list[str]:
    codes = re.findall(r"\[([A-Z][A-Z0-9_]+)\]", text)
    for stable in NEGATIVE_EXPECTED_CODES.values():
        if stable in text and stable not in codes:
            codes.append(stable)
    return codes


def _copy_example(source_root: Path, case_root: Path) -> Path:
    target = case_root / "input"
    shutil.copytree(source_root, target)
    return target / "school-pptx-full.md"


def negative_cases(tree: OwnedTree, skill_root: Path, artifacts: dict[str, object]) -> dict[str, object]:
    run_root = tree.root.path / "work" / tree.run_id
    negative_root = run_root / "negatives"
    negative_root.mkdir(mode=0o700)
    source_root = artifacts["source"].parent
    protected = [
        skill_root / "templates" / "standard-school.pptx",
        skill_root / "templates" / "standard-school.manifest.yaml",
        skill_root / "fixtures" / "school-pptx-full.md",
        artifacts["source"],
    ]
    before = [sha256_path(path) for path in protected]
    caller_sentinel = tree.root.path / ".verify-caller-sentinel"
    caller_sentinel.write_bytes(b"caller-owned")
    cases: list[dict[str, object]] = []
    called: list[str] = []
    for name in NEGATIVE_CASE_ORDER:
        case_root = negative_root / name
        case_root.mkdir(mode=0o700)
        expected = NEGATIVE_EXPECTED_CODES[name]
        if name == "template-manifest-mismatch":
            import yaml

            shutil.copy2(skill_root / "templates" / "standard-school.pptx", case_root / "template.pptx")
            manifest = yaml.safe_load(
                (skill_root / "templates" / "standard-school.manifest.yaml").read_text(encoding="utf-8")
            )
            manifest["layouts"]["cover"]["pptx_layout"] = "ppt/slideLayouts/missing-layout.xml"
            (case_root / "manifest.yaml").write_text(
                yaml.safe_dump(manifest, allow_unicode=True, sort_keys=False),
                encoding="utf-8",
            )
            command = [public_cli(skill_root), "template-report", "--theme", "standard-school", "--manifest", "manifest.yaml", "--template", "template.pptx", "--out-md", "report.md", "--out-json", "report.json"]
        else:
            markdown = _copy_example(source_root, case_root)
            text_value = markdown.read_text(encoding="utf-8")
            if name == "unknown-theme":
                text_value = text_value.replace('theme: "standard-school"', 'theme: "unknown-school"', 1)
            elif name == "unknown-layout":
                text_value = text_value.replace('layout="title-content"', 'layout="unknown-layout"', 1)
            elif name == "missing-media":
                text_value = text_value.replace("media/curriculum-map.png", "media/not-found.png", 1)
            elif name == "unsupported-styling":
                text_value = text_value.replace("课程以", "style=color:red 课程以", 1)
            elif name == "unresolved-review-marker":
                text_value = text_value.replace("课程以", "{{待补充: 教师确认}} 课程以", 1)
            markdown.write_text(text_value, encoding="utf-8")
            if name in {"missing-media", "unresolved-review-marker"}:
                command = [public_cli(skill_root), "render", "--input", "input/school-pptx-full.md", "--out-dir", "delivery", "--stem", "negative"]
            else:
                command = [public_cli(skill_root), "validate", "--input", "input/school-pptx-full.md"]
        observed = run_bounded(command, cwd=case_root)
        codes = _extract_codes(observed["text"])
        writes = sorted(path.relative_to(case_root).as_posix() for path in case_root.rglob("*") if path.is_file())
        false_success = bool(re.search(r"\b(?:PASS|PASSED|SUCCESS)\b", observed["text"], re.I))
        passed = (
            observed["exit_code"] != 0 and expected in codes and observed["bounded"]
            and not observed["traceback"] and not false_success
            and caller_sentinel.read_bytes() == b"caller-owned"
        )
        cases.append({
            "name": name, "status": "failed-as-expected" if passed else "unexpected-result",
            "exit_code": observed["exit_code"], "expected_code": expected,
            "observed_code": expected if expected in codes else (codes[0] if codes else "NONE"),
            "capture_bytes": observed["capture_bytes"], "bounded": observed["bounded"],
            "traceback": observed["traceback"], "false_success": false_success, "writes": writes,
        })
        called.append(name)
    after = [sha256_path(path) for path in protected]
    registry = {"required": list(NEGATIVE_CASE_ORDER), "called": called, "unique": len(called) == len(set(called)), "dynamic_skips": len(NEGATIVE_CASE_ORDER) - len(called)}
    valid = before == after and tuple(called) == NEGATIVE_CASE_ORDER and all(item["status"] == "failed-as-expected" for item in cases)
    result = gate_result("passed" if valid else "failed", "OK" if valid else "NEGATIVE_CASE_FAILED", "All fixed negative cases failed safely with exact stable codes.")
    result["evidence"] = {"registry": registry, "cases": cases, "protected_hashes_unchanged": before == after}
    return result


def build_candidate(
    tree: OwnedTree,
    skill_root: Path = DEFAULT_SKILL_ROOT,
    dependency_gate: Callable[[], dict[str, object]] = dependency_readiness_placeholder,
) -> dict[str, object]:
    started = iso_now()
    called: list[str] = []
    gates: dict[str, dict[str, object]] = {}
    artifact_evidence: list[dict[str, object]] = []
    dependency = dependency_gate()
    gates["dependency-readiness"] = dependency
    called.append("dependency-readiness")
    if dependency.get("status") == "passed":
        flow, artifacts = canonical_flow(tree, skill_root)
        for name in ("example-generation", "template-validation", "canonical-render"):
            gates[name] = flow[name]
            called.append(name)
        gates["structural-inspection"] = structural_inspection(artifacts)
        called.append("structural-inspection")
        gates["phase-43-regression"] = phase_43_regression()
        called.append("phase-43-regression")
        gates["negative-cases"] = negative_cases(tree, skill_root, artifacts)
        called.append("negative-cases")
        artifact_evidence = [
            {"path": "delivery/canonical.md", "sha256": sha256_path(artifacts["markdown"]), "bytes": artifacts["markdown"].stat().st_size},
            {"path": "delivery/canonical.pptx", "sha256": sha256_path(artifacts["pptx"]), "bytes": artifacts["pptx"].stat().st_size},
        ]
    else:
        for name in VERIFY_GATE_ORDER[1:-1]:
            gates[name] = blocked_result()
            called.append(name)
    dependency_observation = gates["dependency-readiness"].get("observed", {})
    candidate: dict[str, object] = {
        "schema_version": 1,
        "status": "failed",
        "run": {"id": tree.run_id, "started_at": started, "finished_at": iso_now()},
        "registry": {
            "required": list(VERIFY_GATE_ORDER),
            "called": list(called),
            "unique": len(VERIFY_GATE_ORDER) == len(set(VERIFY_GATE_ORDER)),
            "dynamic_skips": 0,
        },
        "gates": gates,
        "dependencies": dependency_observation,
        "artifacts": artifact_evidence,
        "diagnostics": [],
    }
    try:
        documentation_integrity_check()
        validate_candidate(candidate, include_integrity=False)
        integrity = gate_result("passed", "OK", "Evidence authority and bounded projection are intact.")
        integrity["evidence"] = {"documentation_integrity": "passed"}
    except VerifyFailure as failure:
        code = failure.code if failure.code == "VERIFY_DOCUMENTATION_INVALID" else "VERIFY_EVIDENCE_INVALID"
        integrity = gate_result("failed", code, STABLE_REMEDIATIONS[code])
    gates["evidence-integrity"] = integrity
    called.append("evidence-integrity")
    candidate["registry"]["called"] = list(called)
    candidate["status"] = _aggregate_status(gates, called)
    negative_evidence = gates["negative-cases"].get("evidence", {})
    negative_registry = negative_evidence.get("registry", {}) if isinstance(negative_evidence, dict) else {}
    negative_items = negative_evidence.get("cases", []) if isinstance(negative_evidence, dict) else []
    phase_43_evidence = gates["phase-43-regression"].get("evidence", {})
    candidate["canonical"] = {"artifacts": artifact_evidence}
    candidate["negative_registry"] = negative_registry
    candidate["negatives"] = [
        {"id": item.get("name"), **{key: value for key, value in item.items() if key != "name"}}
        for item in negative_items
        if isinstance(item, dict)
    ]
    candidate["phase_43"] = {
        "status": gates["phase-43-regression"].get("status"),
        "registry": phase_43_evidence if isinstance(phase_43_evidence, dict) else {},
    }
    validate_candidate(candidate, include_integrity=True)
    return candidate


def source_hashes(skill_root: Path) -> dict[str, str]:
    paths = sorted(
        path for path in skill_root.rglob("*")
        if path.is_file() and not path.is_symlink() and "__pycache__" not in path.parts
    )
    return {
        path.relative_to(skill_root).as_posix(): hashlib.sha256(path.read_bytes()).hexdigest()
        for path in paths
    }


def semantic_projection(candidate: dict[str, object]) -> dict[str, object]:
    return {
        "status": candidate["status"],
        "required": candidate["registry"]["required"],
        "called": candidate["registry"]["called"],
        "gates": candidate["gates"],
        "artifacts": candidate["artifacts"],
    }


def _read_at(directory_fd: int, name: str) -> bytes:
    flags = os.O_RDONLY | getattr(os, "O_NOFOLLOW", 0)
    descriptor = os.open(name, flags, dir_fd=directory_fd)
    try:
        return os.read(descriptor, MAX_EVIDENCE_BYTES + 1)
    finally:
        os.close(descriptor)


def _assert_registry_mutations_rejected(candidate: dict[str, object]) -> None:
    mutations = []
    missing = json.loads(json.dumps(candidate))
    del missing["gates"]["canonical-render"]
    mutations.append(missing)
    duplicate = json.loads(json.dumps(candidate))
    duplicate["registry"]["called"][2] = duplicate["registry"]["called"][1]
    mutations.append(duplicate)
    reordered = json.loads(json.dumps(candidate))
    reordered["registry"]["called"][1:3] = reversed(reordered["registry"]["called"][1:3])
    mutations.append(reordered)
    skipped = json.loads(json.dumps(candidate))
    skipped["registry"]["called"].pop()
    mutations.append(skipped)
    missing_status = json.loads(json.dumps(candidate))
    del missing_status["gates"]["negative-cases"]["status"]
    mutations.append(missing_status)
    absolute_path = json.loads(json.dumps(candidate))
    absolute_path["artifacts"] = [{"path": "/private/secret.pptx", "sha256": "0" * 64, "bytes": 1}]
    mutations.append(absolute_path)
    hardcoded_pass = json.loads(json.dumps(candidate))
    hardcoded_pass["gates"]["canonical-render"]["status"] = "failed"
    hardcoded_pass["status"] = "passed"
    mutations.append(hardcoded_pass)
    constant_count = json.loads(json.dumps(candidate))
    if "evidence" in constant_count["gates"]["structural-inspection"]:
        constant_count["gates"]["structural-inspection"]["evidence"]["physical_count"] = 31
        mutations.append(constant_count)
    oversized = json.loads(json.dumps(candidate))
    oversized["diagnostics"] = ["x" * (MAX_STRING_LENGTH + 1)]
    mutations.append(oversized)
    copied_phase = json.loads(json.dumps(candidate))
    if "evidence" in copied_phase["gates"]["phase-43-regression"]:
        copied_phase["gates"]["phase-43-regression"]["evidence"]["called"] = list(PHASE_43_GATE_ORDER[:-1])
        mutations.append(copied_phase)
    for mutation in mutations:
        try:
            validate_candidate(mutation, include_integrity=True)
        except VerifyFailure:
            continue
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID", "A registry or evidence mutation was accepted.")


def foundation_self_test() -> None:
    before_sources = source_hashes(DEFAULT_SKILL_ROOT)
    projections: list[dict[str, object]] = []
    with tempfile.TemporaryDirectory(prefix="school-pptx-foundation-") as temporary:
        base = Path(temporary)
        for index in range(2):
            root_path = base / f"fresh-{index}"
            with hold_workdir(root_path) as root, build_owned_tree(root) as tree:
                candidate = build_candidate(tree, dependency_gate=lambda: gate_result(
                    "failed", "VERIFY_DEPENDENCY_MISSING", "Synthetic dependency fault."
                ))
                publish_evidence(tree, candidate)
                projections.append(semantic_projection(candidate))
                stored = json.loads(_read_at(tree.evidence_fd, "verification.json"))
                markdown = _read_at(tree.evidence_fd, "verification.md").decode("utf-8")
                if stored["status"] != "failed" or "Status: FAILED" not in markdown:
                    raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
                if f"Gates: {len(VERIFY_GATE_ORDER)}/{len(VERIFY_GATE_ORDER)}" not in markdown:
                    raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
                if os.listdir(tree.delivery_fd):
                    raise VerifyFailure("VERIFY_EVIDENCE_INVALID", "Delivery was polluted by foundation verification.")
                _assert_registry_mutations_rejected(candidate)
        if projections[0] != projections[1]:
            raise VerifyFailure("VERIFY_EVIDENCE_INVALID", "Fresh roots produced different semantic evidence.")

        rerun_path = base / "rerun"
        sentinel = rerun_path / "caller-sentinel.bin"
        with hold_workdir(rerun_path) as root:
            sentinel.write_bytes(b"caller-owned")
            with build_owned_tree(root, "0" * 16) as first_tree:
                passed = build_candidate(first_tree)
                if passed["status"] != "passed":
                    raise VerifyFailure("VERIFY_EVIDENCE_INVALID", "Foundation PASS fixture did not pass.")
                publish_evidence(first_tree, passed)
            with build_owned_tree(root, "1" * 16) as failed_tree:
                failed = build_candidate(failed_tree, dependency_gate=lambda: gate_result(
                    "failed", "VERIFY_DEPENDENCY_MISSING", "Synthetic dependency fault."
                ))
                publish_evidence(failed_tree, failed)
                current = json.loads(_read_at(failed_tree.evidence_fd, "verification.json"))
                if current["status"] != "failed" or current["run"]["id"] != "1" * 16:
                    raise VerifyFailure("VERIFY_EVIDENCE_INVALID", "Stale PASS remained current after failure.")
            if sentinel.read_bytes() != b"caller-owned":
                raise VerifyFailure("VERIFY_WORKDIR_UNSAFE")
    if source_hashes(DEFAULT_SKILL_ROOT) != before_sources:
        raise VerifyFailure("VERIFY_WORKDIR_UNSAFE", "Repository source changed during foundation self-test.")


def dependencies_self_test() -> None:
    normal = dependency_readiness_gate(pandoc_probe=lambda: {
        "required": False, "available": False, "version": None,
    })
    if normal["status"] != "passed":
        raise VerifyFailure("VERIFY_DEPENDENCY_MISSING")
    packages = normal["observed"]["packages"]
    if len(packages) != 4 or not all(item["ready"] and item["version"] for item in packages):
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    python_observation = normal["observed"]["python"]
    if "/" in python_observation["executable_basename"] or not python_observation["executable_identity"].startswith("<external>/"):
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    if normal["observed"]["pandoc"] != {"required": False, "available": False, "version": None}:
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")

    missing_import = REQUIRED_DEPENDENCIES[2][1]

    def fault_importer(name: str) -> object:
        if name == missing_import:
            raise ModuleNotFoundError(name)
        return importlib.import_module(name)

    fault = dependency_readiness_gate(
        importer=fault_importer,
        pandoc_probe=lambda: {"required": False, "available": True, "version": "pandoc test"},
    )
    if fault["status"] != "failed" or fault["code"] != "VERIFY_DEPENDENCY_MISSING":
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    fault_packages = fault["observed"]["packages"]
    if [item["distribution"] for item in fault_packages if not item["ready"]] != ["lxml"]:
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")

    with tempfile.TemporaryDirectory(prefix="school-pptx-dependency-fault-") as temporary:
        sentinel = Path(temporary) / "sentinel"
        with hold_workdir(Path(temporary)) as root:
            sentinel.write_bytes(b"dependency-sentinel")
            with build_owned_tree(root) as tree:
                candidate = build_candidate(tree, dependency_gate=lambda: fault)
                publish_evidence(tree, candidate)
                current = json.loads(_read_at(tree.evidence_fd, "verification.json"))
                if current["status"] != "failed":
                    raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
                if current["gates"]["dependency-readiness"]["code"] != "VERIFY_DEPENDENCY_MISSING":
                    raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
                blocked = [current["gates"][name]["code"] for name in VERIFY_GATE_ORDER[1:-1]]
                if blocked != ["VERIFY_DEPENDENCY_BLOCKED"] * 6:
                    raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
            if sentinel.read_bytes() != b"dependency-sentinel":
                raise VerifyFailure("VERIFY_WORKDIR_UNSAFE")

    source = SCRIPT_PATH.read_text(encoding="utf-8")
    forbidden = (
        "pip " + "install",
        "uv " + "pip",
        "conda " + "install",
        "curl" + " ",
        "wget" + " ",
    )
    if any(value in source for value in forbidden):
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")


def canonical_self_test() -> None:
    with tempfile.TemporaryDirectory(prefix="school-pptx-canonical-") as temporary:
        with hold_workdir(Path(temporary)) as root, build_owned_tree(root) as tree:
            gates, artifacts = canonical_flow(tree, DEFAULT_SKILL_ROOT)
            structure = structural_inspection(artifacts)
            if any(result["status"] != "passed" for result in gates.values()) or structure["status"] != "passed":
                raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
            evidence = structure["evidence"]
            if evidence["logical_count"] != 13 or evidence["physical_count"] != 26:
                raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
            if evidence["transition_mode"] != "none" or evidence["whole_slide_picture_count"] != 0:
                raise VerifyFailure("VERIFY_EVIDENCE_INVALID")


def stable_diagnostics_self_test() -> None:
    from markdown_contract import load_manifest, parse_document

    source = (DEFAULT_SKILL_ROOT / "fixtures" / "school-pptx-full.md").read_text(encoding="utf-8")
    with tempfile.TemporaryDirectory(prefix="school-pptx-diagnostics-") as temporary:
        root = Path(temporary)
        variants = {
            "pending": source.replace("课程以", "{{待补充: 教师确认}} 课程以", 1),
            "draft": source.replace("课程以", "{{AI草稿: 教师确认}} 课程以", 1),
            "ordinary": source.replace("课程以", "{普通花括号} 课程以", 1),
            "fenced": source.replace("<section", "{{待补充: code}}\n<section", 1),
        }
        observed: dict[str, list[str]] = {}
        shutil.copytree(DEFAULT_SKILL_ROOT / "fixtures" / "media", root / "media")
        for name, text_value in variants.items():
            path = root / f"{name}.md"
            path.write_text(text_value, encoding="utf-8")
            document = parse_document(path, load_manifest(DEFAULT_SKILL_ROOT))
            observed[name] = [item["code"] for item in document["errors"]]
        if "REVIEW_MARKER_UNRESOLVED" not in observed["pending"] or "REVIEW_MARKER_UNRESOLVED" not in observed["draft"]:
            raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
        if "REVIEW_MARKER_UNRESOLVED" in observed["ordinary"] or "REVIEW_MARKER_UNRESOLVED" in observed["fenced"]:
            raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
        render_root = root / "render"
        render_root.mkdir()
        (render_root / "input.md").write_text(variants["pending"], encoding="utf-8")
        shutil.copytree(DEFAULT_SKILL_ROOT / "fixtures" / "media", render_root / "media")
        result = run_bounded([public_cli(DEFAULT_SKILL_ROOT), "render", "--input", "input.md", "--out-dir", "delivery", "--stem", "review"], cwd=render_root)
        if result["exit_code"] == 0 or "REVIEW_MARKER_UNRESOLVED" not in _extract_codes(result["text"]):
            raise VerifyFailure("VERIFY_EVIDENCE_INVALID")


def negatives_integrity_self_test() -> None:
    with tempfile.TemporaryDirectory(prefix="school-pptx-negative-integrity-") as temporary:
        with hold_workdir(Path(temporary)) as root, build_owned_tree(root) as tree:
            flow, artifacts = canonical_flow(tree, DEFAULT_SKILL_ROOT)
            structure = structural_inspection(artifacts)
            negative = negative_cases(tree, DEFAULT_SKILL_ROOT, artifacts)
            gates = {name: flow.get(name, gate_result("passed", "OK", "self-test")) for name in VERIFY_GATE_ORDER[:-1]}
            gates["dependency-readiness"] = gate_result("passed", "OK", "self-test")
            gates["structural-inspection"] = structure
            gates["phase-43-regression"] = gate_result("passed", "OK", "self-test")
            gates["phase-43-regression"]["evidence"] = {
                "required": list(PHASE_43_GATE_ORDER), "called": list(PHASE_43_GATE_ORDER),
                "unique": True, "dynamic_skips": 0,
            }
            gates["negative-cases"] = negative
            gates["evidence-integrity"] = gate_result("passed", "OK", "self-test")
            candidate = {
                "schema_version": 1, "status": "passed",
                "run": {"id": tree.run_id, "started_at": iso_now(), "finished_at": iso_now()},
                "registry": {"required": list(VERIFY_GATE_ORDER), "called": list(VERIFY_GATE_ORDER), "unique": True, "dynamic_skips": 0},
                "gates": gates, "dependencies": {}, "artifacts": [], "diagnostics": [],
            }
            validate_candidate(candidate, include_integrity=True)
            _assert_registry_mutations_rejected(candidate)


def _skill_document() -> str:
    path = DEFAULT_SKILL_ROOT / "SKILL.md"
    if not path.is_file():
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    return path.read_text(encoding="utf-8")


def docs_entry_self_test() -> None:
    source = _skill_document()
    frontmatter, separator, body = source.removeprefix("---\n").partition("\n---\n")
    if not separator or 'description: "Use when ' not in frontmatter:
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    runtimes = re.findall(r"^    - (Codex|Claude Code|Gemini CLI|OpenCode|OpenClaw|Hermes Agent)$", frontmatter, re.MULTILINE)
    expected_runtimes = ["Codex", "Claude Code", "Gemini CLI", "OpenCode", "OpenClaw", "Hermes Agent"]
    if runtimes != expected_runtimes or 'portability: "canonical"' not in frontmatter:
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    required_sections = (
        "## Objective and Use When", "## Required Inputs and References", "## Review Before Render",
        "## Missing Information Questions", "## Public Commands", "## Outputs and Verification",
        "## Runtime Adapter Notes", "## Safety",
    )
    if any(section not in body for section in required_sections):
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    required_commands = (" example ", " validate ", " template-report ", " render ", " verify ", "--workdir")
    normalized = f" {body.replace(chr(10), ' ')} "
    if any(command not in normalized for command in required_commands):
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    clarification_terms = (
        "学校/课程标题", "subtitle/term/author metadata", "controlled theme", "required media",
        "unresolved review markers", "Markdown finalized state", "禁止 render",
    )
    if any(term not in body for term in clarification_terms):
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    reference_names = (
        "markdown-contract.md", "template-contract.md", "template-editing.md",
        "renderer-and-pagination.md", "verification-contract.md", "visual-uat.md",
    )
    if any(not (DEFAULT_SKILL_ROOT / "references" / name).is_file() for name in reference_names):
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    main_workflow = body.split("## Runtime Adapter Notes", 1)[0]
    forbidden_long_form = ("UAT-V01", "dependency-readiness", "cover.subtitle", "slideLayout7.xml")
    if any(term in main_workflow for term in forbidden_long_form):
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")
    private_syntax = (".claude/skills", "GEMINI.md", "Claude-compatible", "automatic script discovery")
    if any(term in main_workflow for term in private_syntax):
        raise VerifyFailure("VERIFY_EVIDENCE_INVALID")


def _require_documentation(condition: bool) -> None:
    if not condition:
        raise VerifyFailure("VERIFY_DOCUMENTATION_INVALID")


def documentation_integrity_check() -> None:
    docs_entry_self_test()
    source = _skill_document()
    adapter = source.split("## Runtime Adapter Notes", 1)[1].split("## Safety", 1)[0]
    rows = [line for line in adapter.splitlines() if line.startswith("| ")][1:]
    expected_runtimes = ["Codex", "Claude Code", "Gemini CLI", "OpenCode", "OpenClaw", "Hermes Agent"]
    parsed_rows: list[list[str]] = []
    for row in rows:
        cells = [cell.strip() for cell in row.strip("|").split("|")]
        if cells and cells[0] in expected_runtimes:
            parsed_rows.append(cells)
    _require_documentation([row[0] for row in parsed_rows] == expected_runtimes)
    _require_documentation(all(len(row) == 8 and all(row[1:]) for row in parsed_rows))
    for runtime in ("OpenClaw", "Hermes Agent"):
        row_text = next(" ".join(row[1:]) for row in parsed_rows if row[0] == runtime)
        _require_documentation("installation-time verification" in row_text)
        _require_documentation("fallback" in row_text and "automatic script discovery" in row_text)

    verification = (DEFAULT_SKILL_ROOT / "references" / "verification-contract.md").read_text(encoding="utf-8")
    expected_top = list(VERIFY_GATE_ORDER)
    observed_top = re.findall(r"^\d+\. `([^`]+)`$", verification, re.MULTILINE)[: len(expected_top)]
    _require_documentation(observed_top == expected_top)
    for case_name, code in NEGATIVE_EXPECTED_CODES.items():
        _require_documentation(f"| `{case_name}` | `{code}` |" in verification)
    authority = verification.split("## Phase 43 Nested Authority", 1)[1].split("## Dependency Semantics", 1)[0]
    observed_phase = re.findall(r"^\d+\. `([^`]+)`$", authority, re.MULTILINE)
    _require_documentation(observed_phase == list(PHASE_43_GATE_ORDER))
    required_verification_terms = (
        "--workdir <caller-workdir>", "python-pptx", "Pillow", "lxml", "PyYAML", "Pandoc",
        "none|preserved|generated", "13 logical slides → 26 physical slides", "relative",
        "documentation-integrity", "status: failed",
    )
    _require_documentation(all(term in verification for term in required_verification_terms))

    uat = (DEFAULT_SKILL_ROOT / "references" / "visual-uat.md").read_text(encoding="utf-8")
    expected_uat = [*(f"UAT-V{index:02d}" for index in range(1, 7)), *(f"UAT-I{index:02d}" for index in range(1, 5))]
    observed_uat = re.findall(r"^\| `(UAT-[VI]\d{2})` \|", uat, re.MULTILINE)
    _require_documentation(observed_uat == expected_uat)
    _require_documentation(all(term in uat for term in ("Microsoft PowerPoint", "WPS Presentation", "fixture_sha256", "pptx_sha256")))
    _require_documentation(all(f"`{status}`" in uat for status in ("PENDING", "PASS", "FAIL", "BLOCKED")))
    _require_documentation("YOLO/auto" in uat and "不得填写 tester" in uat)
    _require_documentation("Current approval: PASS" not in uat and "status: PASS" not in uat)

    repository_root = DEFAULT_SKILL_ROOT.parent.parent
    forbidden_directories = (
        repository_root / "skills" / "school-pptx-codex",
        DEFAULT_SKILL_ROOT / "adapters",
        DEFAULT_SKILL_ROOT / "runtime-adapters",
    )
    _require_documentation(not any(path.exists() for path in forbidden_directories))
    indexes = {
        "README.md": ("skills/school-pptx/SKILL.md", "verify --workdir", "Microsoft PowerPoint", "WPS Presentation"),
        "skills/README.md": (
            "`school-pptx`", "fixtures/school-pptx-full.md", "references/markdown-contract.md",
            "references/template-contract.md", "references/verification-contract.md", "references/visual-uat.md",
            "scripts/school-pptx.sh", "templates/standard-school.pptx", "templates/standard-school.manifest.yaml",
        ),
        "docs/compatibility-matrix.md": (
            "### School PPTX", "skills/school-pptx/SKILL.md", "whole-folder", "python-pptx", "Pillow",
            "lxml", "PyYAML", "Pandoc", "external command", "execute", "read", "write", "sandbox/allowlist",
            "explicit invocation fallback", "verify --workdir", "installation-time verification",
        ),
    }
    for relative, terms in indexes.items():
        index_source = (repository_root / relative).read_text(encoding="utf-8")
        _require_documentation(all(term in index_source for term in terms))
        school_section = index_source[index_source.find("school-pptx") :]
        _require_documentation("UAT 已通过" not in school_section and "UAT PASS" not in school_section)

    directory_spec = (repository_root / "docs" / "directory-spec.md").read_text(encoding="utf-8")
    _require_documentation(all(term in directory_spec for term in (
        "skills/school-pptx/fixtures/", "skills/school-pptx/templates/", "references/visual-uat.md",
        ".planning/phases/<phase>/", "不得进入公开 delivery",
    )))


def _parse_uat_frontmatter(source: str) -> tuple[dict[str, object], str]:
    import yaml

    if len(source.encode("utf-8")) > 64 * 1024 or not source.startswith("---\n"):
        raise UATLintFailure("UAT_INVALID")
    frontmatter_text, separator, body = source[4:].partition("\n---\n")
    if not separator:
        raise UATLintFailure("UAT_INVALID")
    try:
        frontmatter = yaml.safe_load(frontmatter_text)
    except yaml.YAMLError as exc:
        raise UATLintFailure("UAT_INVALID") from exc
    if not isinstance(frontmatter, dict):
        raise UATLintFailure("UAT_INVALID")
    return frontmatter, body


def _parse_uat_rows(body: str) -> list[dict[str, str]]:
    rows: list[dict[str, str]] = []
    for line in body.splitlines():
        cells = [cell.strip() for cell in line.strip().strip("|").split("|")]
        if not cells or not re.fullmatch(r"UAT-[VI]\d{2}", cells[0]):
            continue
        if len(cells) != 6:
            raise UATLintFailure("UAT_INVALID")
        rows.append({
            "id": cells[0],
            "check": cells[1],
            "result": cells[2],
            "objects": cells[3],
            "observation": cells[4],
            "reference": cells[5],
        })
    if tuple(row["id"] for row in rows) != UAT_ITEM_ORDER:
        raise UATLintFailure("UAT_INVALID")
    if any(row["result"] not in UAT_RESULT_VALUES or not row["check"] for row in rows):
        raise UATLintFailure("UAT_INVALID")
    return rows


def _uat_timestamp_valid(value: str) -> bool:
    if not value or not re.search(r"(?:Z|[+-]\d{2}:\d{2})$", value):
        return False
    try:
        datetime.fromisoformat(value.replace("Z", "+00:00"))
    except ValueError:
        return False
    return True


def _uat_source_guard() -> None:
    tree = ast.parse(SCRIPT_PATH.read_text(encoding="utf-8"))
    guarded_names = {
        "_parse_uat_frontmatter", "_parse_uat_rows", "_uat_timestamp_valid",
        "_validate_uat_record", "check_uat_evidence",
    }
    guarded = [
        node for node in tree.body
        if isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef)) and node.name in guarded_names
    ]
    if {node.name for node in guarded} != guarded_names:
        raise UATLintFailure("UAT_GUARD_FAILED")
    forbidden_calls = {"write_text", "write_bytes", "unlink", "iso_now"}
    forbidden_names = {"getpass", "LOGNAME", "USER", "structural_inspection", "phase_43_regression"}
    for function in guarded:
        for node in ast.walk(function):
            if isinstance(node, ast.Call):
                name = node.func.attr if isinstance(node.func, ast.Attribute) else node.func.id if isinstance(node.func, ast.Name) else ""
                if name in forbidden_calls:
                    raise UATLintFailure("UAT_GUARD_FAILED")
            if isinstance(node, ast.Name) and node.id in forbidden_names:
                raise UATLintFailure("UAT_GUARD_FAILED")
            if isinstance(node, ast.Subscript) and isinstance(node.ctx, ast.Store):
                raise UATLintFailure("UAT_GUARD_FAILED")


def _validate_uat_record(
    frontmatter: dict[str, object],
    rows: list[dict[str, str]],
    fixture_sha256: str,
    pptx_sha256: str,
    evidence_source: str,
    delivery_pptx: str,
) -> str:
    expected_fields = {
        "status", "viewer", "viewer_version", "operating_system", "fixture_sha256",
        "pptx_sha256", "tested_at", "tester", "evidence_source", "delivery_pptx", "signoff",
    }
    if set(frontmatter) != expected_fields:
        raise UATLintFailure("UAT_INVALID")
    values = {key: frontmatter.get(key) for key in expected_fields}
    if any(not isinstance(value, str) for value in values.values()):
        raise UATLintFailure("UAT_INVALID")
    if values["fixture_sha256"] != fixture_sha256 or values["pptx_sha256"] != pptx_sha256:
        raise UATLintFailure("UAT_BLOCKED")
    if values["evidence_source"] != evidence_source or values["delivery_pptx"] != delivery_pptx:
        raise UATLintFailure("UAT_BLOCKED")
    identity_fields = ("viewer", "viewer_version", "operating_system", "tested_at", "tester")
    results = [row["result"] for row in rows]
    if values["status"] == "pending":
        if any(values[field] for field in identity_fields) or values["signoff"]:
            raise UATLintFailure("UAT_BLOCKED")
        if results != ["PENDING"] * len(UAT_ITEM_ORDER):
            raise UATLintFailure("UAT_BLOCKED")
        if any(row["objects"] or row["observation"] or row["reference"] for row in rows):
            raise UATLintFailure("UAT_BLOCKED")
        return "UAT_PENDING"
    if values["status"] not in {"passed", "failed", "blocked"}:
        raise UATLintFailure("UAT_INVALID")
    if values["viewer"] not in UAT_VIEWERS:
        raise UATLintFailure("UAT_BLOCKED")
    if any(not values[field] for field in identity_fields) or not _uat_timestamp_valid(values["tested_at"]):
        raise UATLintFailure("UAT_BLOCKED")
    if any(not row["objects"] or not row["observation"] for row in rows):
        raise UATLintFailure("UAT_BLOCKED")
    if "PENDING" in results:
        raise UATLintFailure("UAT_PENDING")
    if "FAIL" in results:
        raise UATLintFailure("UAT_FAILED")
    if "BLOCKED" in results:
        raise UATLintFailure("UAT_BLOCKED")
    if values["status"] != "passed" or not values["signoff"]:
        raise UATLintFailure("UAT_BLOCKED")
    return "UAT_PASSED"


def check_uat_evidence(uat_path: Path, evidence_path: Path) -> str:
    _uat_source_guard()
    try:
        uat_source = uat_path.read_text(encoding="utf-8")
        evidence = json.loads(evidence_path.read_text(encoding="utf-8"))
    except (OSError, UnicodeError, json.JSONDecodeError) as exc:
        raise UATLintFailure("UAT_INVALID") from exc
    if any(token in uat_source for token in ("Traceback", "/Users/", "/home/", "/tmp/", "C:\\")):
        raise UATLintFailure("UAT_INVALID")
    frontmatter, body = _parse_uat_frontmatter(uat_source)
    rows = _parse_uat_rows(body)
    try:
        validate_candidate(evidence, include_integrity=True)
    except VerifyFailure as exc:
        raise UATLintFailure("UAT_BLOCKED") from exc
    if evidence.get("status") != "passed":
        raise UATLintFailure("UAT_BLOCKED")
    artifacts = evidence.get("artifacts")
    if not isinstance(artifacts, list):
        raise UATLintFailure("UAT_BLOCKED")
    markdown = next((item for item in artifacts if item.get("path") == "delivery/canonical.md"), None)
    pptx = next((item for item in artifacts if item.get("path") == "delivery/canonical.pptx"), None)
    if not isinstance(markdown, dict) or not isinstance(pptx, dict):
        raise UATLintFailure("UAT_BLOCKED")
    workdir = evidence_path.parent.parent
    markdown_path = workdir / str(markdown["path"])
    pptx_path = workdir / str(pptx["path"])
    try:
        fixture_sha256 = sha256_path(markdown_path)
        pptx_sha256 = sha256_path(pptx_path)
    except OSError as exc:
        raise UATLintFailure("UAT_BLOCKED") from exc
    if fixture_sha256 != markdown.get("sha256") or pptx_sha256 != pptx.get("sha256"):
        raise UATLintFailure("UAT_BLOCKED")
    repository_root = DEFAULT_SKILL_ROOT.parents[1].resolve(strict=True)
    evidence_source = evidence_path.resolve(strict=True)
    delivery_value = frontmatter.get("delivery_pptx")
    if not isinstance(delivery_value, str) or Path(delivery_value).is_absolute():
        raise UATLintFailure("UAT_BLOCKED")
    try:
        evidence_relative = evidence_source.relative_to(repository_root).as_posix()
        delivery_path = (repository_root / delivery_value).resolve(strict=True)
        delivery_relative = delivery_path.relative_to(repository_root).as_posix()
        delivery_sha256 = sha256_path(delivery_path)
    except (OSError, ValueError) as exc:
        raise UATLintFailure("UAT_BLOCKED") from exc
    if delivery_relative != delivery_value or delivery_sha256 != pptx_sha256:
        raise UATLintFailure("UAT_BLOCKED")
    return _validate_uat_record(
        frontmatter,
        rows,
        fixture_sha256,
        pptx_sha256,
        evidence_relative,
        delivery_relative,
    )


def parse_uat_args(argv: list[str]) -> argparse.Namespace | None:
    if not argv or argv[0] != "--check-uat":
        return None
    parser = argparse.ArgumentParser(add_help=False)
    parser.add_argument("--check-uat", type=Path, required=True)
    parser.add_argument("--evidence", type=Path, required=True)
    try:
        return parser.parse_args(argv)
    except SystemExit as exc:
        raise UATLintFailure("UAT_INVALID") from exc


def parse_self_test(argv: list[str]) -> argparse.Namespace | None:
    if not argv or argv[0] != "--self-test":
        return None
    parser = argparse.ArgumentParser(add_help=False)
    parser.add_argument("--self-test", choices=("foundation", "dependencies", "canonical", "stable-diagnostics", "negatives-integrity", "docs-entry", "docs-portability"), required=True)
    try:
        return parser.parse_args(argv)
    except SystemExit as exc:
        raise VerifyFailure("VERIFY_ARGUMENT_INVALID") from exc


def run_self_test(name: str) -> None:
    if name == "foundation":
        foundation_self_test()
        return
    if name == "dependencies":
        dependencies_self_test()
        return
    if name == "canonical":
        canonical_self_test()
        return
    if name == "stable-diagnostics":
        stable_diagnostics_self_test()
        return
    if name == "docs-entry":
        docs_entry_self_test()
        return
    if name == "docs-portability":
        documentation_integrity_check()
        return
    negatives_integrity_self_test()


def run_public(skill_root: Path, workdir: HeldWorkdir) -> int:
    if skill_root.resolve(strict=True) != DEFAULT_SKILL_ROOT.resolve(strict=True):
        raise VerifyFailure("VERIFY_ARGUMENT_INVALID")
    with build_owned_tree(workdir) as tree:
        candidate = build_candidate(tree, skill_root)
        publish_evidence(tree, candidate)
    if candidate["status"] != "passed":
        raise VerifyFailure("VERIFY_NOT_IMPLEMENTED")
    print("School PPTX verification passed")
    return 0


def main(argv: list[str]) -> int:
    try:
        uat_args = parse_uat_args(argv)
        if uat_args is not None:
            result = check_uat_evidence(uat_args.check_uat, uat_args.evidence)
            print(result)
            return 0 if result == "UAT_PASSED" else 3
    except UATLintFailure as failure:
        print(failure.code)
        return 3
    try:
        self_test = parse_self_test(argv)
        if self_test is not None:
            run_self_test(self_test.self_test)
            return 0
        skill_root, requested_workdir = parse_public_args(argv)
        with hold_workdir(requested_workdir) as workdir:
            return run_public(skill_root, workdir)
    except VerifyFailure as failure:
        return emit_failure(failure)
    except (RuntimeError, KeyError, json.JSONDecodeError, subprocess.SubprocessError):
        return emit_failure(VerifyFailure("VERIFY_EVIDENCE_INVALID"))
    except (OSError, ValueError, TypeError):
        return emit_failure(VerifyFailure("VERIFY_WORKDIR_UNSAFE"))
    except Exception:
        return emit_failure(VerifyFailure("VERIFY_EVIDENCE_INVALID"))


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
