#!/usr/bin/env python3
"""Task-local regression gates for the Phase 43 renderer contract and model."""

from __future__ import annotations

import argparse
import ast
import contextlib
import copy
import fcntl
import hashlib
import io
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
import zlib
from dataclasses import replace
from pathlib import Path, PurePosixPath
from xml.etree import ElementTree as ET
from zipfile import BadZipFile, ZIP_DEFLATED, ZipFile

import pptx_paginate
import template_report
from markdown_contract import load_manifest, parse_document


SKILL_DIR = Path(__file__).resolve().parent.parent
SCRIPTS_DIR = SKILL_DIR / "scripts"
MANIFEST_PATH = SKILL_DIR / "templates" / "standard-school.manifest.yaml"
TEMPLATE_PATH = SKILL_DIR / "templates" / "standard-school.pptx"
PUBLIC_CLI = SCRIPTS_DIR / "school-pptx.sh"
FIXTURE_PATH = SKILL_DIR / "fixtures" / "school-pptx-full.md"
DELIVERY_FIXTURE = SKILL_DIR / "fixtures" / "clean-delivery" / "confirmed-assets.md"
DELIVERY_FAULTS = (
    "after_candidate_validation",
    "after_history_reservation",
    "after_archive_snapshot",
    "after_publish_file_1",
    "after_publish_middle_file",
    "before_post_publish_verify",
    "before_work_cleanup",
)
CANONICAL_PHYSICAL_SLIDES = 26
MAX_ZIP_ENTRIES = 256
MAX_ZIP_ENTRY_BYTES = 4 * 1024 * 1024
MAX_ZIP_TOTAL_BYTES = 32 * 1024 * 1024
FORBIDDEN_XML_MARKERS = (b"<!DOCTYPE", b"<!ENTITY")
PRESENTATION_NS = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}


class GateFailure(RuntimeError):
    pass


def require(condition: bool, message: str) -> None:
    if not condition:
        raise GateFailure(message)


def safe_package_entries(path: Path) -> dict[str, bytes]:
    try:
        with ZipFile(path) as package:
            infos = package.infolist()
            require(len(infos) <= MAX_ZIP_ENTRIES, "ZIP_ENTRY_LIMIT")
            names: set[str] = set()
            total = 0
            result: dict[str, bytes] = {}
            for info in infos:
                name = info.filename
                pure = PurePosixPath(name)
                require(not pure.is_absolute() and ".." not in pure.parts, "ZIP_PATH_INVALID")
                require(name not in names, "ZIP_DUPLICATE_ENTRY")
                names.add(name)
                require(info.file_size <= MAX_ZIP_ENTRY_BYTES, "ZIP_ENTRY_SIZE_LIMIT")
                total += info.file_size
                require(total <= MAX_ZIP_TOTAL_BYTES, "ZIP_TOTAL_SIZE_LIMIT")
                if name.endswith(".xml") or name.endswith(".rels"):
                    payload = package.read(info)
                    upper = payload.upper()
                    require(not any(marker in upper for marker in FORBIDDEN_XML_MARKERS), "XML_EXTERNAL_ENTITY_FORBIDDEN")
                    try:
                        ET.fromstring(payload)
                    except ET.ParseError as exc:
                        raise GateFailure("XML_MALFORMED") from exc
                    result[name] = payload
            return result
    except (BadZipFile, OSError) as exc:
        raise GateFailure("ZIP_MALFORMED") from exc


def run_template_report(workdir: Path) -> dict[str, object]:
    report_json = workdir / "template-report.json"
    report_md = workdir / "template-report.md"
    completed = subprocess.run(
        [
            str(SCRIPTS_DIR / "school-pptx.sh"),
            "template-report",
            "--theme",
            "standard-school",
            "--out-md",
            str(report_md),
            "--out-json",
            str(report_json),
        ],
        cwd=SKILL_DIR.parent.parent,
        text=True,
        capture_output=True,
        timeout=30,
        check=False,
    )
    require(completed.returncode == 0, "controlled template-report failed")
    require("Traceback" not in completed.stdout + completed.stderr, "template-report leaked traceback")
    data = json.loads(report_json.read_text(encoding="utf-8"))
    require(not data.get("failures"), "template-report contains failures")
    require(report_md.is_file(), "template-report markdown evidence missing")
    return data


def manifest_renderer_contract_gate(workdir: Path) -> dict[str, object]:
    report = run_template_report(workdir)
    manifest = template_report.read_yaml(MANIFEST_PATH)
    layouts = manifest["layouts"]
    require(manifest["template"]["mapping_strategy"] == "renderer-owned-native-shapes",
            "renderer-owned mapping strategy changed")
    require(report["renderer_contract"]["mapping_strategy"] == "renderer-owned-native-shapes",
            "template-report mapping strategy drifted")
    cover_subtitle = next(slot for slot in layouts["cover"]["slots"] if slot["id"] == "subtitle")
    require(cover_subtitle["text_budget"]["max_chars"] == 96, "cover subtitle max_chars must be 96")
    require(cover_subtitle["text_budget"]["max_lines"] == 3, "cover subtitle max_lines must be 3")
    require(report["renderer_contract"]["cover_subtitle_budget"] == cover_subtitle["text_budget"],
            "template-report cover subtitle evidence drifted")
    table = next(slot for slot in layouts["table"]["slots"] if slot["id"] == "table")
    gallery_items = next(slot for slot in layouts["gallery"]["slots"] if slot["id"] == "gallery_items")
    gallery_caption = next(slot for slot in layouts["gallery"]["slots"] if slot["id"] == "caption")
    timeline = next(slot for slot in layouts["timeline"]["slots"] if slot["id"] == "timeline_items")
    table_name = table["subregions"]["table_name"]
    require(table_name["empty_slot"] == "preserve", "D-14 table_name empty slot not preserved")
    require(gallery_caption["empty_slot"] == "preserve", "D-14 gallery caption empty slot not preserved")
    require(template_report.validate_text_budget(table_name["text_budget"]), "table_name budget invalid")
    require(set(gallery_items["item_presets"]) == {1, 2, 3, 4}, "gallery presets incomplete")
    require(set(timeline["subregions"]) == {"axis", "node_band"}, "timeline subdivisions incomplete")
    require(
        set(timeline["node_template"]["subregions"]) == {"marker", "time", "title", "description"},
        "timeline node subdivisions incomplete",
    )
    require(manifest["inline_styles"]["highlight"]["scheme_color"] == "accent1",
            "sixth-version highlight must use theme-blue accent1")
    require(layouts["closing"]["pptx_layout"] == "ppt/slideLayouts/slideLayout8.xml", "closing part path invalid")
    require(len(report["layouts"]) == 11 and len(report["slots"]) >= 1, "layout regression")

    invalid = copy.deepcopy(manifest)
    invalid_gallery = next(slot for slot in invalid["layouts"]["gallery"]["slots"] if slot["id"] == "gallery_items")
    invalid_gallery["item_presets"][1][0]["card"]["x"] = invalid_gallery["geometry"]["x"] - 1
    bounded = template_report.validate_manifest(invalid, MANIFEST_PATH, TEMPLATE_PATH, 2000)
    require(any("越界" in item for item in bounded["failures"]), "out-of-bounds preset did not fail")
    return {
        "layout_count": len(report["layouts"]),
        "mapping_strategy": report["renderer_contract"]["mapping_strategy"],
        "closing_part_path": layouts["closing"]["pptx_layout"],
        "highlight_scheme_color": manifest["inline_styles"]["highlight"]["scheme_color"],
        "gallery_presets": sorted(gallery_items["item_presets"]),
        "table_name_empty_slot": table_name["empty_slot"],
    }


def local_import_graph(root: Path) -> set[Path]:
    pending = [root]
    visited: set[Path] = set()
    while pending:
        path = pending.pop()
        if path in visited:
            continue
        visited.add(path)
        tree = ast.parse(path.read_text(encoding="utf-8"), filename=str(path))
        for node in ast.walk(tree):
            names: list[str] = []
            if isinstance(node, ast.Import):
                names.extend(alias.name for alias in node.names)
            elif isinstance(node, ast.ImportFrom) and node.module:
                names.append(node.module)
            for name in names:
                require(name != "pptx" and not name.startswith("pptx."), "pptx import boundary violated")
                candidate = SCRIPTS_DIR / f"{name.split('.')[0]}.py"
                if candidate.is_file():
                    pending.append(candidate)
    return visited


def model_determinism_gate(workdir: Path) -> dict[str, object]:
    del workdir
    first = pptx_paginate.run_contract_model_self_check()
    second = pptx_paginate.run_contract_model_self_check()
    require(first == second, "model self-check evidence is not deterministic")
    require(first["model_equal"] is True and first["model_immutable"] is True, "frozen model contract failed")
    graph = local_import_graph(SCRIPTS_DIR / "pptx_paginate.py")
    return {"model_equal": True, "model_immutable": True, "import_graph": sorted(path.name for path in graph)}


def expect_package_failure(path: Path, expected: str) -> None:
    try:
        safe_package_entries(path)
    except GateFailure as exc:
        require(str(exc) == expected, f"expected {expected}, got {exc}")
        return
    raise GateFailure(f"negative package did not fail: {expected}")


def measurement_gate(workdir: Path) -> dict[str, object]:
    manifest = load_manifest(SKILL_DIR)
    evidence = pptx_paginate.run_contract_model_self_check()
    vectors = evidence["vectors"]
    require(vectors["explicit_newline"]["display_lines"] == 3, "explicit newline vector failed")
    require(vectors["font_clamp"]["font_size"] == 18.0, "font clamp vector failed")
    try:
        pptx_paginate.TextMeasure().measure(
            "测" * (pptx_paginate.MAX_TEXT_BYTES + 1),
            width_emu=1_270_000,
            font_size=16,
            font_size_min=12,
            font_size_max=18,
        )
    except ValueError as exc:
        require(str(exc) == "text exceeds bounded measurement input", "unbounded text failure")
    else:
        raise GateFailure("oversized text did not fail")

    controlled = safe_package_entries(TEMPLATE_PATH)
    presentation = ET.fromstring(controlled["ppt/presentation.xml"])
    seed_count = len(presentation.findall(".//p:sldId", PRESENTATION_NS))
    require(seed_count == manifest["template"]["seed_slides"],
            "controlled template seed count differs from manifest")

    malformed = workdir / "malformed.pptx"
    malformed.write_bytes(b"not-a-zip")
    expect_package_failure(malformed, "ZIP_MALFORMED")

    xxe = workdir / "xxe.pptx"
    with ZipFile(xxe, "w", ZIP_DEFLATED) as package:
        package.writestr("ppt/presentation.xml", '<!DOCTYPE x [<!ENTITY e SYSTEM "file:///etc/passwd">]><x>&e;</x>')
    expect_package_failure(xxe, "XML_EXTERNAL_ENTITY_FORBIDDEN")

    malformed_xml = workdir / "malformed-xml.pptx"
    with ZipFile(malformed_xml, "w", ZIP_DEFLATED) as package:
        package.writestr("ppt/presentation.xml", "<x>")
    expect_package_failure(malformed_xml, "XML_MALFORMED")

    bomb = workdir / "entry-bomb.pptx"
    with ZipFile(bomb, "w", ZIP_DEFLATED) as package:
        package.writestr("ppt/presentation.xml", b" " * (MAX_ZIP_ENTRY_BYTES + 1))
    expect_package_failure(bomb, "ZIP_ENTRY_SIZE_LIMIT")
    return {"seed_slide_count": seed_count, "vectors": sorted(vectors), "security_negative_cases": 4}


GATES = {
    "manifest_renderer_contract_gate": manifest_renderer_contract_gate,
    "model_determinism_gate": model_determinism_gate,
    "measurement_gate": measurement_gate,
}


def frozen_slot_content_gate(workdir: Path) -> dict[str, object]:
    manifest = load_manifest(SKILL_DIR)
    document = parse_document(FIXTURE_PATH, manifest)
    first = pptx_paginate.build_deck_plan(document, manifest)
    second = pptx_paginate.build_deck_plan(document, manifest)
    require(first.to_projection() == second.to_projection(), "frozen projection is not deterministic")
    two_column = [slide for slide in first.slides if slide.layout == "two-column"]
    fragments = [fragment for slide in two_column for fragment in slide.fragments]
    require([fragment.target_slot for fragment in fragments] ==
            ["left_body", "right_body", "left_body", "right_body", "left_body"],
            "D-06 target slot sequence changed")
    require([fragment.block_index for fragment in fragments] == list(range(5)), "D-06 block inventory changed")
    cover = next(slide for slide in first.slides if slide.layout == "cover")
    values = dict(cover.slot_values)
    require(values["title"] and len(values["subtitle"]) == 64, "canonical cover descriptor length changed")
    require(values["subtitle"].count("\n") == 2, "canonical cover information bar is not three lines")
    require(all(value in values["subtitle"] for value in
                ("示例职业技术学院", "智能制造学院", "智能产线安装与调试", "课程建设团队", "2026-07-13")),
            "canonical cover descriptor lost metadata")
    require("机电一体化技术" not in values["subtitle"] and "张老师" not in values["subtitle"],
            "program/presenter leaked into cover subtitle")
    overflow = copy.deepcopy(document)
    overflow["metadata"]["school"] = "超长学校" * 30
    overflow_plan = pptx_paginate.build_deck_plan(overflow, manifest)
    require(any(item.code == "COVER_METADATA_OVERFLOW" and item.severity == "error"
                for item in overflow_plan.diagnostics), "cover overflow did not fail")
    overflow_source = workdir / "cover-overflow.md"
    overflow_source.write_text(
        "---\n"
        'title: "独立封面负例"\n'
        'subtitle: "只验证非 canonical 超长封面"\n'
        f'school: "{"超长学校" * 30}"\n'
        'department: "智能制造学院"\n'
        'course: "智能产线安装与调试"\n'
        'author: "课程建设团队"\n'
        'date: "2026-07-13"\n'
        'theme: "standard-school"\n'
        "---\n\n"
        '::: slide {layout="cover"}\n:::\n',
        encoding="utf-8",
    )
    overflow_output = workdir / "cover-overflow-output"
    overflow_output.mkdir()
    overflow_result = run_public_render(overflow_source, overflow_output, "overflow")
    overflow_text = overflow_result.stdout + overflow_result.stderr
    require(overflow_result.returncode != 0 and "COVER_METADATA_OVERFLOW" in overflow_text,
            "independent public cover overflow did not fail bounded")
    require("Traceback" not in overflow_text and len(overflow_text.encode("utf-8")) < 8 * 1024,
            "public cover overflow leaked traceback or unbounded output")
    return {
        "target_slots": [fragment.target_slot for fragment in fragments],
        "cover_length": len(values["subtitle"]),
        "public_overflow_code": "COVER_METADATA_OVERFLOW",
    }


GATES["frozen-slot-content"] = frozen_slot_content_gate


def frozen_numbering_row_heights_gate(workdir: Path) -> dict[str, object]:
    del workdir
    manifest = load_manifest(SKILL_DIR)
    document = parse_document(FIXTURE_PATH, manifest)
    plan = pptx_paginate.build_deck_plan(document, manifest)
    contents = [slide for slide in plan.slides if slide.layout == "contents"]
    visible = [item for slide in contents for fragment in slide.fragments for item in fragment.items]
    require(visible == document["contents_entries"], "contents section entries changed")
    require(all(not re.match(r"^\d+[.、]", item) for item in visible), "contents numbering was not removed")
    require(all(slide.title == "目录" for slide in contents), "contents title changed")
    table_slot = next(slot for slot in manifest["layouts"]["table"]["slots"] if slot["id"] == "table")
    budget = table_slot["subregions"]["table_body"]["geometry"]["height"]
    tables = [slide.fragments[0] for slide in plan.slides if slide.layout == "table"]
    require(budget > 0 and tables, "derived table budget is invalid")
    require(all(len(fragment.row_heights_emu) == len(fragment.rows) for fragment in tables),
            "table row-height vector length changed")
    require(all(all(isinstance(value, int) and value > 0 for value in fragment.row_heights_emu)
                and sum(fragment.row_heights_emu) <= budget for fragment in tables),
            "table row-height vector exceeds derived budget")
    require(all(fragment.row_heights_emu[0] == tables[0].row_heights_emu[0] for fragment in tables),
            "repeated table header height changed")
    require(pptx_paginate._allocate_row_heights((4, 6), 10) == (4, 6),
            "equal-budget allocation changed measured heights")
    require(pptx_paginate._allocate_row_heights((4, 6), 15) == (7, 8),
            "positive slack allocation is unstable")
    projection = plan.to_projection()
    require(any(fragment["row_heights_emu"] for slide in projection["slides"]
                for fragment in slide["fragments"] if fragment["kind"] == "table"),
            "row heights missing from projection")
    return {"contents_items": len(visible), "content_budget_emu": budget, "table_pages": len(tables)}


GATES["frozen-numbering-row-heights"] = frozen_numbering_row_heights_gate


def frozen_plan_emission_gate(workdir: Path) -> dict[str, object]:
    import pptx_emit
    from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN

    manifest = load_manifest(SKILL_DIR)
    document = parse_document(FIXTURE_PATH, manifest)
    plan = pptx_paginate.build_deck_plan(document, manifest)
    output = workdir / "frozen-plan.pptx"
    pptx_emit.emit_deck(plan, manifest, TEMPLATE_PATH, output, media_root=FIXTURE_PATH.parent)
    presentation = pptx_emit.require_dependencies()["pptx"].Presentation(output)
    require(len(presentation.slides) == len(plan.slides), "emitted slide count differs from frozen plan")

    for physical, slide in zip(plan.slides, presentation.slides):
        names = {shape.name: shape for shape in slide.shapes}
        title_slot = next((
            slot for slot in manifest["layouts"][physical.layout].get("slots", ())
            if slot["id"] == "title"
        ), None)
        if title_slot is not None and physical.layout != "cover":
            title = names["school-pptx:slide-title"]
            expected_alignment = {
                "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT,
                "justify": PP_ALIGN.JUSTIFY,
            }[title_slot["paragraph_alignment"]]
            expected_anchor = {
                "top": MSO_VERTICAL_ANCHOR.TOP,
                "middle": MSO_VERTICAL_ANCHOR.MIDDLE,
                "bottom": MSO_VERTICAL_ANCHOR.BOTTOM,
            }[title_slot["vertical_anchor"]]
            require(title.text_frame.paragraphs[0].alignment == expected_alignment,
                    f"{physical.layout} title horizontal alignment changed")
            require(title.text_frame.vertical_anchor == expected_anchor,
                    f"{physical.layout} title vertical alignment changed")
            expected_size = float(title_slot["text_budget"]["font_size_max"])
            require(all(run.font.size and run.font.size.pt == expected_size
                        for run in title.text_frame.paragraphs[0].runs),
                    f"{physical.layout} title font size changed")
        if physical.layout == "cover":
            require(names["school-pptx:cover-title"].text == dict(physical.slot_values)["title"],
                    "cover title differs from frozen slot value")
            require(names["school-pptx:cover-subtitle"].text == dict(physical.slot_values)["subtitle"],
                    "cover subtitle differs from frozen slot value")
        elif physical.layout == "two-column":
            for slot_id in ("left_body", "right_body"):
                require(f"school-pptx:{slot_id}" in names, f"missing native {slot_id} shape")
                require(all(
                    paragraph.alignment == PP_ALIGN.JUSTIFY
                    for paragraph in names[f"school-pptx:{slot_id}"].text_frame.paragraphs
                    if paragraph.text.strip()
                ), f"{slot_id} is not justified")
                expected: list[str] = []
                for fragment in physical.fragments:
                    if fragment.target_slot != slot_id:
                        continue
                    if fragment.heading:
                        expected.append(fragment.heading)
                    if fragment.text is not None:
                        expected.append(fragment.text)
                    expected.extend(fragment.items)
                require(names[f"school-pptx:{slot_id}"].text == "\n".join(expected),
                        f"{slot_id} text differs from frozen fragments")
        elif physical.layout == "contents":
            require(names["school-pptx:body"].text == "\n".join(physical.fragments[0].items),
                    "contents visible entries differ from frozen items")
            require(names["school-pptx:body"].text_frame.vertical_anchor == MSO_VERTICAL_ANCHOR.MIDDLE,
                    "contents body is not vertically centered")
            require(all(paragraph.alignment == PP_ALIGN.JUSTIFY
                        for paragraph in names["school-pptx:body"].text_frame.paragraphs),
                    "contents body is not justified")
        elif physical.layout == "table":
            table_shape = names["school-pptx:native-table"]
            require(tuple(row.height for row in table_shape.table.rows) == physical.fragments[0].row_heights_emu,
                    "native table row heights differ from frozen EMU vector")
            for row_index, row in enumerate(table_shape.table.rows):
                for cell in row.cells:
                    require(cell.vertical_anchor == MSO_VERTICAL_ANCHOR.MIDDLE,
                            "native table cell vertical alignment changed")
                    require(all(p.alignment == PP_ALIGN.CENTER for p in cell.text_frame.paragraphs),
                            "native table cell horizontal alignment changed")
                    expected_fill = "4472C4" if row_index == 0 else "D9EAF7"
                    require(str(cell.fill.fore_color.rgb) == expected_fill,
                            "native table blue fill changed")
        elif physical.layout == "timeline":
            from pptx.oxml.ns import qn

            axis = names["school-pptx:timeline-axis"]
            gradient = axis._element.spPr.find(qn("a:gradFill"))
            require(gradient is not None, "timeline axis gradient missing")
            stops = gradient.findall(".//" + qn("a:gs"))
            require([stop.get("pos") for stop in stops] == ["0", "100000"],
                    "timeline axis gradient stops changed")
            require([stop[0].get("val") for stop in stops] == ["accent6", "accent1"],
                    "timeline axis theme gradient changed")
            require(gradient.find(qn("a:lin")).get("ang") == "10800000",
                    "timeline axis gradient direction changed")
            groups = [shape for shape in slide.shapes if shape.name.startswith("school-pptx:timeline-node:")]
            require(groups, "timeline node groups missing")
            marker_colors: list[tuple[int, int, int]] = []
            for group in groups:
                xfrm = group._element.grpSpPr.xfrm
                require(xfrm.ext.cx == xfrm.chExt.cx and xfrm.ext.cy == xfrm.chExt.cy,
                        "timeline group non-uniformly scales its circular marker")
                marker = next(shape for shape in group.shapes if "timeline-marker" in shape.name)
                require(marker.width == marker.height, "timeline marker is not square")
                marker_colors.append(tuple(marker.fill.fore_color.rgb))
                for role, expected_size in (("timeline-time", 20.0), ("timeline-title", 20.0),
                                            ("timeline-description", 16.0)):
                    text_shape = next(shape for shape in group.shapes if role in shape.name)
                    require(all(run.font.size and run.font.size.pt == expected_size
                                for paragraph in text_shape.text_frame.paragraphs for run in paragraph.runs),
                            f"{role} font size changed")
                    expected_alignment = PP_ALIGN.LEFT if role == "timeline-description" else PP_ALIGN.CENTER
                    require(all(paragraph.alignment == expected_alignment
                                for paragraph in text_shape.text_frame.paragraphs),
                            f"{role} alignment changed")
                    require(text_shape.text_frame.vertical_anchor == MSO_VERTICAL_ANCHOR.MIDDLE,
                            f"{role} vertical anchor changed")
            require(marker_colors[0] != marker_colors[-1], "timeline marker gradient colors collapsed")
            require(marker_colors[0][1] > marker_colors[-1][1]
                    and marker_colors[0][2] < marker_colors[-1][2],
                    "timeline marker colors do not follow green-to-blue axis")

        require(physical.notes_intent and slide.has_notes_slide,
                f"{physical.layout} speaker notes missing")
        require(slide.notes_slide.notes_text_frame.text.strip() == physical.notes_intent.strip(),
                f"{physical.layout} speaker notes changed")

    source = (SCRIPTS_DIR / "pptx_emit.py").read_text(encoding="utf-8")
    require("pptx_paginate" not in source and "pptx_measure" not in source,
            "emitter imports pagination or measurement")
    require("table_height / len(rows)" not in (SCRIPTS_DIR / "pptx_objects.py").read_text(encoding="utf-8"),
            "average row-height formula remains")
    return {"slides": len(plan.slides), "two_column_pages": sum(s.layout == "two-column" for s in plan.slides)}


GATES["frozen-plan-emission"] = frozen_plan_emission_gate


def _minimal_document(slides: list[dict[str, object]]) -> dict[str, object]:
    return {
        "document_title": "测试文稿",
        "logical_slides": slides,
        "implicit_slides": [{"layout": "closing", "position": "end_of_deck", "fixed_template_page": True}],
        "contents_entries": [],
    }


def pagination_text_code_gate(workdir: Path) -> dict[str, object]:
    del workdir
    manifest = load_manifest(SKILL_DIR)
    sentence = "第一句完整结束。第二句仍然很长！第三句结束？"
    sentence_parts = pptx_paginate.split_semantic_text(sentence, 12)
    require("".join(sentence_parts) == sentence and sentence_parts[0].endswith("。"), "D-02 sentence split")
    weak = "前半部分，后半部分仍然很长"
    weak_parts = pptx_paginate.split_semantic_text(weak, 6)
    require("".join(weak_parts) == weak and weak_parts[0].endswith("，"), "D-02 weak punctuation split")
    unicode_text = "A\u0301👩\u200d🏭️B"
    unicode_parts = pptx_paginate.split_semantic_text(unicode_text, 1)
    require("".join(unicode_parts) == unicode_text and unicode_parts[0] == "A\u0301", "D-02 grapheme split")

    long_code_line = "x = '" + "中" * 2000 + "'"
    code_text = "line_1\n" + "\n".join(f"line_{index}" for index in range(2, 28)) + "\n" + long_code_line
    document = _minimal_document([
        {"layout": "title-content", "title": "正文", "source_line": 1, "notes": None, "blocks": [
            {"kind": "paragraph", "source_line": 2, "heading": "同一标题", "text": "短块完整保留。"},
            {"kind": "paragraph", "source_line": 3, "heading": "同一标题", "text": "这是一个很长的句子。" * 80},
        ]},
        {"layout": "code", "title": "代码", "source_line": 10, "notes": None, "blocks": [
            {"kind": "code", "source_line": 11, "heading": "源码", "language": "python", "text": code_text},
        ]},
    ])
    first = pptx_paginate.build_deck_plan(document, manifest)
    second = pptx_paginate.build_deck_plan(document, manifest)
    require(first == second, "pagination determinism")
    text_slides = [slide for slide in first.slides if slide.layout == "title-content"]
    require(text_slides[0].fragments[0].text == "短块完整保留。", "D-01 complete block split")
    require(all(slide.title == "正文" for slide in text_slides), "D-04 repeated title")
    require(all("续" not in slide.title for slide in first.slides), "D-05 visible continuation marker")
    headings = [fragment.heading for slide in text_slides for fragment in slide.fragments]
    require(all(heading == "同一标题" for heading in headings), "D-04 block heading changed")
    code_fragments = [fragment for slide in first.slides for fragment in slide.fragments if fragment.kind == "code"]
    require("\n".join(fragment.text or "" for fragment in code_fragments) == code_text, "D-03 code round trip")
    require(all(fragment.text is not None for fragment in code_fragments), "D-03 code fragment")
    require(any(d.code == "CODE_LINE_OVERFLOW" for d in first.diagnostics), "CODE_LINE_OVERFLOW missing")
    return {
        "sentence_parts": len(sentence_parts), "weak_parts": len(weak_parts),
        "unicode_parts": len(unicode_parts), "physical_slides": len(first.slides),
        "code_fragments": len(code_fragments),
    }


def pagination_structured_gate(workdir: Path) -> dict[str, object]:
    del workdir
    manifest = load_manifest(SKILL_DIR)
    rows = [[f"模块{index}", "说明" * (index % 4 + 3), "成果完整"] for index in range(1, 24)]
    timeline_rows = [[f"第{index}月", f"节点{index}", "阶段说明" * 14] for index in range(1, 16)]
    entries = [f"目录项目 {index} " + ("长标题" * (index % 3)) for index in range(1, 15)]
    document = _minimal_document([
        {"layout": "table", "title": "表格页", "source_line": 1, "notes": None, "blocks": [{
            "kind": "table", "source_line": 2, "heading": "原表名", "table_title": "原表名",
            "headers": ["模块", "说明", "成果"], "rows": rows,
        }]},
        {"layout": "table", "title": "无名表", "source_line": 30, "notes": None, "blocks": [{
            "kind": "table", "source_line": 31, "heading": None, "table_title": None,
            "headers": ["模块", "说明", "成果"], "rows": rows,
        }]},
        {"layout": "timeline", "title": "时间线", "source_line": 60, "notes": None, "blocks": [{
            "kind": "timeline", "source_line": 61, "heading": None,
            "headers": ["时间", "标题", "说明"], "rows": timeline_rows,
        }]},
        {"layout": "contents", "title": None, "source_line": 90, "notes": None, "blocks": []},
    ])
    document["contents_entries"] = entries
    plan = pptx_paginate.build_deck_plan(document, manifest)
    table_slides = [slide for slide in plan.slides if slide.logical_index == 0]
    require(len(table_slides) > 1, "D-07 table did not paginate")
    require(all(slide.fragments[0].rows[0] == ("模块", "说明", "成果") for slide in table_slides), "D-07 header")
    require(sum(len(slide.fragments[0].rows) - 1 for slide in table_slides) == len(rows), "D-08 row loss")
    names = [dict(slide.fragments[0].metadata)["table_name"] for slide in table_slides]
    require(names[0] == "原表名" and all(name == "原表名（续）" for name in names[1:]), "D-06 table name")
    unnamed = [slide for slide in plan.slides if slide.logical_index == 1]
    require(all(dict(slide.fragments[0].metadata)["table_name"] == "" for slide in unnamed), "D-06 empty table name")

    # Demonstrate that a larger size can remove a visible wrapped orphan.
    measure = pptx_paginate.TextMeasure()
    orphan_vector = None
    for length in range(8, 160):
        value = "测" * length
        small = pptx_paginate._wrapped_cell(value, 900_000, 12.0, measure)[1]
        large = pptx_paginate._wrapped_cell(value, 900_000, 18.0, measure)[1]
        if small and not large:
            orphan_vector = length
            break
    require(orphan_vector is not None, "D-08 larger-font orphan vector missing")

    timeline = [slide for slide in plan.slides if slide.logical_index == 2]
    counts = [len(slide.fragments[0].rows) for slide in timeline]
    require(len(timeline) >= 3 and min(counts) >= 3 and max(counts) - min(counts) <= 1, "D-09 global timeline balance")
    flattened = [row[1] for slide in timeline for row in slide.fragments[0].rows]
    require(flattened == [row[1] for row in timeline_rows], "D-09 timeline order")
    infeasible_document = _minimal_document([{
        "layout": "timeline", "title": "不可行", "source_line": 1, "notes": None, "blocks": [{
            "kind": "timeline", "source_line": 2, "heading": None,
            "headers": ["时间", "标题", "说明"],
            "rows": [[str(index), f"节点{index}", "极高说明" * 400] for index in range(5)],
        }],
    }])
    infeasible = pptx_paginate.build_deck_plan(infeasible_document, manifest)
    require(any(item.code == "TIMELINE_BALANCE_INFEASIBLE" for item in infeasible.diagnostics),
            "D-09 infeasible diagnostic")

    contents = [slide for slide in plan.slides if slide.logical_index == 3]
    content_counts = [len(slide.fragments[0].items) for slide in contents]
    require(all(slide.title == "目录" for slide in contents), "D-10 contents title")
    require(max(content_counts) - min(content_counts) <= 1, "D-10 contents count balance")
    require([item for slide in contents for item in slide.fragments[0].items] == entries,
            "D-10 unnumbered contents order")
    try:
        pptx_paginate.ordered_contiguous_partition([1.0] * 300, 10.0)
    except ValueError as exc:
        require(str(exc) == "partition state budget exceeded", "bounded partition diagnostic")
    else:
        raise GateFailure("partition resource limit missing")
    return {
        "table_pages": len(table_slides), "table_font": dict(table_slides[0].selected_font_sizes)["table"],
        "orphan_vector": orphan_vector, "timeline_counts": counts, "contents_counts": content_counts,
    }


def _image(index: int, caption: str | None = None) -> dict[str, object]:
    return {
        "kind": "image", "source_line": index + 2, "heading": None,
        "authored_path": f"media/{index}.png", "caption": caption if caption is not None else f"图 {index}",
    }


def pagination_full_gate(workdir: Path) -> dict[str, object]:
    del workdir
    manifest = load_manifest(SKILL_DIR)
    gallery_vectors: dict[int, list[int]] = {}
    for count in (1, 2, 3, 4, 5, 6, 8, 9):
        images = [_image(index, "" if index == 0 else None) for index in range(count)]
        document = _minimal_document([{
            "layout": "gallery", "title": "图集", "source_line": 1, "notes": None, "blocks": images,
        }])
        plan = pptx_paginate.build_deck_plan(document, manifest)
        gallery = [slide for slide in plan.slides if slide.layout == "gallery"]
        sizes = [len(slide.fragments) for slide in gallery]
        gallery_vectors[count] = sizes
        require(all(size <= 4 for size in sizes), "PPTX-06 gallery capacity")
        require(sum(sizes) == count, "PPTX-06 gallery item loss")
        for slide in gallery:
            require(all(int(dict(fragment.metadata)["gallery_preset"]) == len(slide.fragments)
                        for fragment in slide.fragments), "PPTX-06 gallery preset")
            require(all(dict(fragment.metadata)["caption_placeholder"] == "true"
                        for fragment in slide.fragments), "PPTX-06 caption placeholder")

    fixture = SKILL_DIR / "fixtures" / "school-pptx-full.md"
    document = parse_document(fixture, manifest)
    require(not document["errors"], "full fixture parser errors")
    first = pptx_paginate.build_deck_plan(document, manifest)
    second = pptx_paginate.build_deck_plan(document, manifest)
    require(first == second, "full fixture plan is not deterministic")
    require(first.slides[-1].layout == "closing", "closing is not last")
    require(sum(slide.layout == "closing" for slide in first.slides) == 1, "closing count")
    require(dict(first.slides[-1].fragments[0].metadata)["pptx_layout"] ==
            "ppt/slideLayouts/slideLayout8.xml", "closing part path")

    for logical_index, indices in first.logical_to_physical:
        require(indices == tuple(range(indices[0], indices[-1] + 1)), "logical-to-physical mapping not contiguous")
        logical = document["logical_slides"][logical_index]
        expected_notes = (
            logical["notes"]["markdown"] if logical.get("notes")
            else pptx_paginate._default_notes(logical, logical_index, document["logical_slides"], document)
        )
        derived = [first.slides[index] for index in indices]
        require(all(slide.notes_intent == expected_notes for slide in derived), "D-16 notes propagation")
        require(all(slide.to_projection()["notes_intent"] for slide in derived), "D-16 notes missing")

    section_notes = [slide.notes_intent for slide in first.slides if slide.layout == "section"]
    require(all(note and "本章内容概括" in note for note in section_notes), "section notes lack summary")
    require(first.slides[-1].notes_intent, "closing notes missing")

    layouts = {slide.layout for slide in first.slides}
    require({"title-content", "two-column", "image-text", "table", "timeline", "gallery", "code"} <= layouts,
            "full fixture overflow layout coverage")
    image_text = [slide for slide in first.slides if slide.layout == "image-text"]
    require(len(image_text) >= 3 and all(any(fragment.kind == "image" for fragment in slide.fragments)
                                         for slide in image_text), "image-text cartesian planning")
    contents = [slide for slide in first.slides if slide.layout == "contents"]
    require([item for slide in contents for item in slide.fragments[0].items] == document["contents_entries"],
            "full fixture contents order")
    return {
        "gallery_vectors": gallery_vectors,
        "logical_slides": len(document["logical_slides"]),
        "physical_slides": len(first.slides),
        "diagnostics": [item.code for item in first.diagnostics],
    }


GATES["pagination-text-code"] = pagination_text_code_gate
GATES["pagination-structured"] = pagination_structured_gate
GATES["pagination-full"] = pagination_full_gate


def ooxml_bootstrap_gate(workdir: Path) -> dict[str, object]:
    try:
        import pptx_emit
        from pptx_ooxml import set_run_highlight
        from pptx.util import Inches
    except ImportError as exc:
        raise GateFailure("PPTX_DEPENDENCY_MISSING") from exc

    manifest = load_manifest(SKILL_DIR)
    presentation, layouts, evidence = pptx_emit.bootstrap_template(TEMPLATE_PATH, manifest)
    require(evidence == {
        "seed_slides": manifest["template"]["seed_slides"],
        "removed_seed_slides": manifest["template"]["seed_slides"],
        "removed_seed_notes_relationships": manifest["template"]["seed_notes_relationships"],
    }, "third-version seed removal evidence mismatch")
    require(len(presentation.slides) == 0, "seed slides remain after bootstrap")
    require(set(layouts) == set(manifest["layouts"]), "layout part map incomplete")
    require(str(layouts["closing"].part.partname).lstrip("/") == "ppt/slideLayouts/slideLayout8.xml",
            "closing layout did not resolve by part path")

    invalid = copy.deepcopy(manifest)
    invalid["layouts"]["closing"]["pptx_layout"] = "ppt/slideLayouts/missing.xml"
    try:
        pptx_emit.bootstrap_template(TEMPLATE_PATH, invalid)
    except pptx_emit.PptxEmitError as exc:
        require(exc.code == "PPTX_TEMPLATE_LAYOUT_INVALID", "layout failure code mismatch")
    else:
        raise GateFailure("missing layout part did not fail")

    slide = presentation.slides.add_slide(layouts["closing"])
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run = text_box.text_frame.paragraphs[0].add_run()
    run.text = "保留粗体与文本"
    run.font.bold = True
    scheme = manifest["inline_styles"]["highlight"]["scheme_color"]
    set_run_highlight(run, scheme)
    set_run_highlight(run, scheme)
    staged = workdir / "bootstrap.pptx"
    presentation.save(staged)
    reopened = pptx_emit.require_dependencies()["pptx"].Presentation(staged)
    require(len(reopened.slides) == 1, "bootstrap package did not reopen")
    require(reopened.slides[0].shapes[-1].text == "保留粗体与文本", "highlight changed run text")
    require(reopened.slides[0].shapes[-1].text_frame.paragraphs[0].runs[0].font.bold is True,
            "highlight changed bold")
    entries = safe_package_entries(staged)
    slide_xml = entries["ppt/slides/slide1.xml"]
    require(slide_xml.count(b"<a:highlight>") == 1 and f'val="{scheme}"'.encode() in slide_xml,
            "highlight is not idempotent theme OOXML")

    def missing_import(name: str) -> object:
        raise ModuleNotFoundError(name)

    try:
        pptx_emit.require_dependencies(missing_import)
    except pptx_emit.PptxEmitError as exc:
        require(exc.code == "PPTX_DEPENDENCY_MISSING" and "Traceback" not in str(exc),
                "dependency failure is not bounded")
    else:
        raise GateFailure("missing dependency did not fail")
    return {**evidence, "layouts": len(layouts), "highlight_scheme": scheme, "reopened_slides": 1}


GATES["ooxml-bootstrap"] = ooxml_bootstrap_gate


def code_literal_roundtrip_gate(workdir: Path) -> dict[str, object]:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.util import Inches
        import pptx_emit
        import pptx_objects
    except ImportError as exc:
        raise GateFailure("PPTX_DEPENDENCY_MISSING") from exc

    vector = "if a == b == c\nreturn **value**\nx = ***raw***\n  keep  spaces  "
    source = (SCRIPTS_DIR / "pptx_objects.py").read_text(encoding="utf-8")
    tree = ast.parse(source, filename="pptx_objects.py")
    helper = next(
        node for node in tree.body
        if isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef)) and node.name == "add_literal_text"
    )
    forbidden = {"normalize_rich_text", "inline_spans", "add_rich_text", "set_run_highlight"}
    referenced = {node.id for node in ast.walk(helper) if isinstance(node, ast.Name)}
    require(not forbidden & referenced, "C-02 literal helper reaches delimiter parser or highlight helper")

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    pptx_objects.add_literal_text(
        slide,
        {"x": Inches(1), "y": Inches(1), "width": Inches(8), "height": Inches(4)},
        vector,
        font_size=18,
        name="school-pptx:code",
    )
    output = workdir / "code-literal-roundtrip.pptx"
    presentation.save(output)
    reopened = Presentation(output)
    code_shapes = [shape for shape in reopened.slides[0].shapes if shape.name == "school-pptx:code"]
    require(len(code_shapes) == 1, "literal vector is not in exactly one school-pptx:code shape")
    code_shape = code_shapes[0]
    require(code_shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX, "literal code is not a native textbox")
    require(code_shape.text == vector, "literal code changed after PPTX reopen")
    runs = [run for paragraph in code_shape.text_frame.paragraphs for run in paragraph.runs]
    require(len(runs) == 1 and runs[0].text == vector, "literal code is not stored in one run")
    require(all(run.font.name == "Consolas" for run in runs), "literal code run is not monospace")

    wide_line = "wide = '" + "中" * 48 + "'"
    title_code = "if a == b == c\nreturn **value**\n" + wide_line
    column_code = "column = ==literal== and ***raw***"
    markdown = workdir / "code-roundtrip.md"
    markdown.write_text(
        "---\n"
        'title: "代码保真"\n'
        'theme: "standard-school"\n'
        "---\n"
        '::: slide {layout="title-content"}\n'
        "## 混排代码页\n"
        "普通 **bold** 与 ==highlight==。\n"
        "```text\n"
        + title_code
        + "\n```\n"
        ":::\n"
        '::: slide {layout="two-column"}\n'
        "## 双栏代码页\n"
        "```text\n"
        + column_code
        + "\n```\n"
        "右栏 **bold** 与 ==highlight==。\n"
        ":::\n",
        encoding="utf-8",
    )
    manifest = load_manifest(SKILL_DIR)
    document = parse_document(markdown, manifest)
    require(not document["errors"], f"code round-trip parser errors: {document['errors']}")
    parser_title_code = document["logical_slides"][0]["blocks"][1]["text"]
    parser_column_code = document["logical_slides"][1]["blocks"][0]["text"]
    require(parser_title_code == title_code and parser_column_code == column_code, "parser changed authored code")
    plan = pptx_paginate.build_deck_plan(document, manifest)
    content_slides = [slide for slide in plan.slides if slide.layout == "title-content"]
    column_slides = [slide for slide in plan.slides if slide.layout == "two-column"]
    content_slide = content_slides[0]
    column_slide = column_slides[0]
    frozen_title_fragments = [
        fragment.text or "" for slide in content_slides for fragment in slide.fragments if fragment.kind == "code"
    ]
    frozen_column_fragments = [
        fragment.text or "" for slide in column_slides for fragment in slide.fragments if fragment.kind == "code"
    ]
    frozen_title_code = "\n".join(frozen_title_fragments)
    frozen_column_code = "\n".join(frozen_column_fragments)
    require(frozen_title_code == parser_title_code and frozen_column_code == parser_column_code,
            "frozen plan changed authored code")

    emitter_tree = ast.parse((SCRIPTS_DIR / "pptx_emit.py").read_text(encoding="utf-8"), filename="pptx_emit.py")
    emitter_imports = {
        alias.name
        for node in ast.walk(emitter_tree)
        if isinstance(node, ast.Import)
        for alias in node.names
    } | {
        node.module for node in ast.walk(emitter_tree) if isinstance(node, ast.ImportFrom) and node.module
    }
    require("pptx_paginate" not in emitter_imports, "emitter imports paginator")
    mixed_helper = next(
        node for node in ast.parse(source, filename="pptx_objects.py").body
        if isinstance(node, ast.FunctionDef) and node.name == "add_fragment_text_frame"
    )
    helper_calls = {
        node.func.id for node in ast.walk(mixed_helper)
        if isinstance(node, ast.Call) and isinstance(node.func, ast.Name)
    }
    require(not {"add_rich_text", "add_plain_lines", "inline_spans"} & helper_calls,
            "mixed-fragment helper delegates code through rich-text line normalization")
    require("fragment_paragraph_sequence" in helper_calls,
            "mixed-fragment helper does not consume shared paragraph sequence")
    code_kind_branches = [
        node for node in ast.walk(mixed_helper)
        if isinstance(node, ast.If) and "role == 'code'" in ast.unparse(node.test)
    ]
    require(len(code_kind_branches) == 1, "mixed-fragment helper code branch missing")
    code_branch_source = ast.unparse(code_kind_branches[0].body)
    require("run.text = text" in code_branch_source,
            "mixed-fragment helper does not assign projected code text directly")
    require(not {"normalize_rich_text", "inline_spans", "add_rich_text", "add_plain_lines"} & {
        node.id for statement in code_kind_branches[0].body for node in ast.walk(statement)
        if isinstance(node, ast.Name)
    }, "mixed-fragment code branch reaches delimiter normalization")
    code_branches = [
        node for node in ast.walk(emitter_tree)
        if isinstance(node, ast.If)
        and any(
            isinstance(candidate, ast.Compare)
            and isinstance(candidate.left, ast.Attribute)
            and isinstance(candidate.left.value, ast.Name)
            and candidate.left.value.id == "physical"
            and candidate.left.attr == "layout"
            and len(candidate.ops) == 1
            and isinstance(candidate.ops[0], ast.Eq)
            and len(candidate.comparators) == 1
            and isinstance(candidate.comparators[0], ast.Constant)
            and candidate.comparators[0].value == "code"
            for candidate in ast.walk(node.test)
        )
    ]
    require(len(code_branches) == 1, "dedicated code emitter branch missing")
    code_branch_nodes = [candidate for statement in code_branches[0].body for candidate in ast.walk(statement)]
    branch_calls = {
        node.func.id for node in code_branch_nodes
        if isinstance(node, ast.Call) and isinstance(node.func, ast.Name)
    }
    require("add_literal_text" in branch_calls, "code branch does not call literal helper")
    require("add_plain_lines" not in branch_calls, "code body reaches line-joining rich-text helper")
    literal_call = next(
        node for node in code_branch_nodes
        if isinstance(node, ast.Call) and isinstance(node.func, ast.Name) and node.func.id == "add_literal_text"
    )
    require(
        len(literal_call.args) >= 3 and ast.unparse(literal_call.args[2]) == "code_text",
        "literal helper does not consume the ordered frozen code text",
    )
    code_text_assignments = [
        node for node in code_branch_nodes
        if isinstance(node, ast.Assign)
        and any(isinstance(target, ast.Name) and target.id == "code_text" for target in node.targets)
    ]
    require(
        len(code_text_assignments) == 1
        and "item.text" in ast.unparse(code_text_assignments[0].value)
        and "code_fragments" in ast.unparse(code_text_assignments[0].value),
        "code text is not reconstructed from frozen fragments in plan order",
    )

    generic_branch = next(
        node for node in ast.walk(emitter_tree)
        if isinstance(node, ast.For)
        and isinstance(node.target, ast.Name)
        and node.target.id == "body_slot_id"
    )
    generic_calls = [
        node for statement in generic_branch.body for node in ast.walk(statement)
        if isinstance(node, ast.Call) and isinstance(node.func, ast.Name)
        and node.func.id == "add_fragment_text_frame"
    ]
    require(len(generic_calls) == 1, "emitter must call mixed-fragment helper once per target slot")

    output_root = workdir / "public-code-roundtrip"
    completed = run_public_render(markdown, output_root, "code-full-chain")
    require(completed.returncode == 0, f"public mixed-code render failed: {completed.stdout + completed.stderr}")
    emitted = output_root / "code-full-chain.pptx"
    require({path.name for path in output_root.iterdir()} == {"code-full-chain.md", "code-full-chain.pptx"},
            "public mixed-code render leaked sidecar artifacts")
    public_output = completed.stdout + completed.stderr
    require("Traceback" not in public_output and str(SKILL_DIR) not in public_output
            and str(Path.home() / ".cache" / "uv") not in public_output,
            "public mixed-code success leaked internal paths or traceback")
    reopened_chain = Presentation(emitted)
    content_body_shapes = [
        shape for physical in content_slides for shape in reopened_chain.slides[physical.physical_index].shapes
        if shape.name == "school-pptx:body"
    ]
    left_body_shapes = [
        shape for physical in column_slides for shape in reopened_chain.slides[physical.physical_index].shapes
        if shape.name == "school-pptx:left_body"
    ]
    right_body_shapes = [
        shape for physical in column_slides for shape in reopened_chain.slides[physical.physical_index].shapes
        if shape.name == "school-pptx:right_body"
    ]
    require(content_body_shapes and left_body_shapes and right_body_shapes,
            "non-code layout did not create required body textboxes")
    content_body = content_body_shapes[0]
    title_code_runs = [
        run for shape in content_body_shapes for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs if run.font.name == "Consolas"
    ]
    column_code_runs = [
        run for shape in left_body_shapes for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs if run.font.name == "Consolas"
    ]
    require(content_body.text_frame.paragraphs[0].text == "普通 bold 与 highlight。",
            "title-content rich paragraph order differs from frozen fragment order")
    require("\n".join(run.text for run in column_code_runs) == frozen_column_code == parser_column_code,
            "two-column parser-plan-PPTX code text differs")
    require(sum(shape.text.count(wide_line) for shape in content_body_shapes) == 1,
            "soft wrap inserted a newline into the wide source line")
    code_runs = [*title_code_runs, *column_code_runs]
    require("\n".join(run.text for run in title_code_runs) == frozen_title_code,
            "title-content reopened code differs from frozen fragments")
    require(len(title_code_runs) == len(frozen_title_fragments)
            and len(column_code_runs) == len(frozen_column_fragments),
            "each frozen code fragment must occupy one paragraph and one run")
    require(all(run.font.name == "Consolas" and run.font.bold is not True for run in code_runs),
            "code run gained non-literal styling")
    require(all("<a:highlight>" not in run._r.xml and 'b="1"' not in run._r.xml for run in code_runs),
            "code delimiters generated rich-text OOXML")
    rich_runs = content_body.text_frame.paragraphs[0].runs
    require(any(run.font.bold is True and run.text == "bold" for run in rich_runs),
            "ordinary bold rich text regressed")
    highlight_runs = [run for run in rich_runs if run.text == "highlight"]
    require(len(highlight_runs) == 1, "ordinary highlight run count changed")
    highlight_run = highlight_runs[0]
    highlight_xml = highlight_run._r.xml
    highlight_scheme = manifest["inline_styles"]["highlight"]["scheme_color"]
    require(highlight_scheme == "accent1", "ordinary highlight theme token is not blue accent1")
    require(
        "<a:highlight>" in highlight_xml
        and f'<a:schemeClr val="{highlight_scheme}"' in highlight_xml,
        "ordinary highlight run lost theme-blue highlight OOXML",
    )
    require(highlight_run.font.bold is True, "ordinary highlight run is not bold")
    require(str(highlight_run.font.color.rgb) == "FFFFFF",
            "ordinary highlight run font color is not FFFFFF")
    return {
        "characters": len(parser_title_code) + len(parser_column_code),
        "runs": len(code_runs), "font": code_runs[0].font.name,
        "parser_plan_pptx_equal": True, "wide_line_newlines": 0, "rich_text_regression": "PASS",
        "body_shapes_per_slot": 1, "public_artifacts": 2,
    }


GATES["code-literal-roundtrip"] = code_literal_roundtrip_gate


def mixed_fragment_capacity_gate(workdir: Path) -> dict[str, object]:
    try:
        from pptx import Presentation
        import pptx_emit
        from pptx_model import PhysicalDeckPlan
    except ImportError as exc:
        raise GateFailure("PPTX_DEPENDENCY_MISSING") from exc

    code_text = "\n".join("中" * 40 for _ in range(8))
    markdown = workdir / "mixed-fragment-capacity.md"
    markdown.write_text(
        "---\n"
        'title: "混排容量"\n'
        'theme: "standard-school"\n'
        "---\n"
        '::: slide {layout="title-content"}\n'
        "## 标题正文容量\n"
        "普通 **bold** 与 ==highlight==。\n\n"
        "```text\n" + code_text + "\n```\n"
        ":::\n"
        '::: slide {layout="two-column"}\n'
        "## 双栏容量\n"
        "```text\n" + code_text + "\n```\n"
        "右栏 **bold** 与 ==highlight==。\n"
        ":::\n",
        encoding="utf-8",
    )
    manifest = load_manifest(SKILL_DIR)
    document = parse_document(markdown, manifest)
    require(not document["errors"], f"mixed capacity parser errors: {document['errors']}")
    plan = pptx_paginate.build_deck_plan(document, manifest)
    parser_codes = {
        "title-content": document["logical_slides"][0]["blocks"][1]["text"],
        "two-column": document["logical_slides"][1]["blocks"][0]["text"],
    }
    output_root = workdir / "delivery"
    completed = run_public_render(markdown, output_root, "mixed-fragment-capacity")
    public_output = completed.stdout + completed.stderr
    require(completed.returncode == 0, f"mixed capacity render failed: {public_output[:1024]}")
    pptx_path = output_root / "mixed-fragment-capacity.pptx"
    markdown_path = output_root / "mixed-fragment-capacity.md"
    require({path.name for path in output_root.iterdir()} == {markdown_path.name, pptx_path.name},
            "mixed capacity public artifacts changed")
    reopened = Presentation(pptx_path)

    vectors: dict[str, object] = {}
    for layout, slot_id, expected_font in (
        ("title-content", "body", 28.0),
        ("two-column", "left_body", 22.0),
    ):
        physical_slides = [slide for slide in plan.slides if slide.layout == layout]
        code_slides = [slide for slide in physical_slides if any(fragment.kind == "code" for fragment in slide.fragments)]
        require(len(physical_slides) > 1 and len(code_slides) > 1,
                f"{layout} over-capacity code did not add physical pages")
        frozen_fragments: list[str] = []
        reopened_runs: list[str] = []
        page_equalities: list[bool] = []
        display_heights: list[float] = []
        effective_heights: list[float] = []
        effective_widths: list[float] = []
        frozen_typographies: list[dict[str, float]] = []
        reopened_typographies: list[dict[str, float]] = []
        selected_fonts: list[float] = []
        slot = next(item for item in manifest["layouts"][layout]["slots"] if item["id"] == slot_id)
        for physical in code_slides:
            selected = dict(physical.selected_font_sizes)[slot_id]
            selected_fonts.append(selected)
            require(selected == expected_font, f"{layout} frozen font differs from target slot")
            values = dict(physical.slot_values)
            typography = {
                key: float(values[f"typography.{slot_id}.{key}"])
                for key in (
                    "margin_left_points", "margin_right_points", "margin_top_points", "margin_bottom_points",
                    "line_spacing", "paragraph_spacing_points",
                )
            }
            frozen_typographies.append(typography)
            shape = next(
                shape for shape in reopened.slides[physical.physical_index].shapes
                if shape.name == f"school-pptx:{slot_id}"
            )
            code_paragraph, code_run = next(
                (paragraph, run) for paragraph in shape.text_frame.paragraphs for run in paragraph.runs
                if run.font.name == "Consolas"
            )
            fragment = next(fragment for fragment in physical.fragments if fragment.kind == "code")
            frozen_fragments.append(fragment.text or "")
            reopened_runs.append(code_run.text)
            page_equalities.append(fragment.text == code_run.text)
            require(code_run.font.size is not None and code_run.font.size.pt == expected_font,
                    f"{layout} reopened code font differs from frozen font")
            reopened_typography = {
                "margin_left_points": shape.text_frame.margin_left.pt,
                "margin_right_points": shape.text_frame.margin_right.pt,
                "margin_top_points": shape.text_frame.margin_top.pt,
                "margin_bottom_points": shape.text_frame.margin_bottom.pt,
                "line_spacing": float(code_paragraph.line_spacing),
                "paragraph_spacing_points": code_paragraph.space_after.pt,
            }
            reopened_typographies.append(reopened_typography)
            require(reopened_typography == typography, f"{layout} reopened typography differs from frozen plan")
            effective_width = slot["geometry"]["width"] / pptx_paginate.EMU_PER_POINT
            effective_width -= typography["margin_left_points"] + typography["margin_right_points"]
            effective_height = slot["geometry"]["height"] / pptx_paginate.EMU_PER_POINT
            effective_height -= typography["margin_top_points"] + typography["margin_bottom_points"]
            measurement = pptx_paginate.TextMeasure(
                line_spacing=typography["line_spacing"],
                paragraph_spacing_points=typography["paragraph_spacing_points"],
            ).measure(
                fragment.text or "",
                width_emu=round(effective_width * pptx_paginate.EMU_PER_POINT),
                font_size=selected,
                font_size_min=selected,
                font_size_max=selected,
                kind="code",
            )
            require(measurement.height_points <= effective_height,
                    f"{layout} reopened display height exceeds effective content height")
            display_heights.append(measurement.height_points)
            effective_heights.append(effective_height)
            effective_widths.append(effective_width)
        parser_text = parser_codes[layout]
        frozen_text = "\n".join(frozen_fragments)
        reopened_text = "\n".join(reopened_runs)
        require(parser_text == frozen_text == reopened_text,
                f"{layout} parser-plan-PPTX joined code changed")
        require(all(page_equalities), f"{layout} per-page code text changed")
        vectors[layout] = {
            "public_exit": completed.returncode,
            "physical_pages": len(physical_slides),
            "code_pages": len(code_slides),
            "selected_font_sizes": selected_fonts,
            "frozen_typography": frozen_typographies,
            "reopened_typography": reopened_typographies,
            "effective_content_width_points": effective_widths,
            "display_height_points": display_heights,
            "effective_content_height_points": effective_heights,
            "page_equalities": page_equalities,
            "parser_text_hash": hashlib.sha256(parser_text.encode()).hexdigest(),
            "frozen_text_hash": hashlib.sha256(frozen_text.encode()).hexdigest(),
            "reopened_text_hash": hashlib.sha256(reopened_text.encode()).hexdigest(),
            "text_lengths": [len(parser_text), len(frozen_text), len(reopened_text)],
            "joined_equality": parser_text == frozen_text == reopened_text,
            "output_bytes": len(public_output.encode("utf-8")),
            "artifact_paths": [str(markdown_path), str(pptx_path)],
        }

    body_slide = next(slide for slide in plan.slides if slide.layout == "title-content")
    missing_font_slide = replace(body_slide, selected_font_sizes=())
    missing_font_plan = PhysicalDeckPlan(
        tuple(missing_font_slide if slide is body_slide else slide for slide in plan.slides),
        plan.diagnostics,
        plan.logical_to_physical,
    )
    try:
        pptx_emit.emit_deck(missing_font_plan, manifest, TEMPLATE_PATH, workdir / "missing-font.pptx")
    except pptx_emit.PptxEmitError as exc:
        require(exc.code == "PPTX_PLAN_FONT_MISSING", "missing body font did not fail closed")
    else:
        raise GateFailure("missing body font emitted successfully")
    missing_typography_slide = replace(body_slide, slot_values=())
    missing_typography_plan = PhysicalDeckPlan(
        tuple(missing_typography_slide if slide is body_slide else slide for slide in plan.slides),
        plan.diagnostics,
        plan.logical_to_physical,
    )
    try:
        pptx_emit.emit_deck(missing_typography_plan, manifest, TEMPLATE_PATH, workdir / "missing-typography.pptx")
    except pptx_emit.PptxEmitError as exc:
        require(exc.code == "PPTX_PLAN_TYPOGRAPHY_MISSING", "missing typography did not fail closed")
    else:
        raise GateFailure("missing body typography emitted successfully")
    return {
        "vectors": vectors,
        "fail_closed_codes": ["PPTX_PLAN_FONT_MISSING", "PPTX_PLAN_TYPOGRAPHY_MISSING"],
    }


GATES["mixed-fragment-capacity"] = mixed_fragment_capacity_gate


def frame_capacity_consistency_gate(workdir: Path) -> dict[str, object]:
    try:
        from pptx import Presentation
        import pptx_emit
        from pptx_model import PhysicalDeckPlan, fragment_paragraph_sequence
    except ImportError as exc:
        raise GateFailure("PPTX_DEPENDENCY_MISSING") from exc

    manifest = load_manifest(SKILL_DIR)
    code_lines = "\n".join("中" * 20 for _ in range(30))
    long_code_heading = "**多行代码标题**" + "确定性可见文本" * 12 + "==高亮结尾=="
    markdown = workdir / "frame-capacity-consistency.md"
    markdown.write_text(
        "---\n"
        'title: "整帧容量"\n'
        'theme: "standard-school"\n'
        "---\n"
        '::: slide {layout="contents"}\n:::\n'
        + "".join(
            f'::: slide {{layout="section"}}\n## 第{index}项\n:::\n'
            for index in range(1, 11)
        )
        + "".join(
            f'::: slide {{layout="title-content"}}\n## 第{index}项\n### 小标题\n正文内容。\n:::\n'
            for index in range(1, 6)
        )
        + '::: slide {layout="title-content"}\n## 段落向量\n### **重复标题** ==高亮==\n段落 **粗体** 与 ==高亮==。\n:::\n'
        + '::: slide {layout="title-content"}\n## 列表向量\n### **重复标题** ==高亮==\n- 第一项 **粗体**\n- 第二项 ==高亮==\n:::\n'
        + '::: slide {layout="title-content"}\n## 混排代码向量\n### 源码\n```text\n'
        + "\n".join("中" * 40 for _ in range(4)) + "\n```\n:::\n"
        + '::: slide {layout="code"}\n## 专用代码\n```text\n' + code_lines + "\n```\n:::\n"
        + '::: slide {layout="code"}\n## 带标题专用代码\n### ' + long_code_heading
        + '\n```text\n' + code_lines + "\n```\n:::\n",
        encoding="utf-8",
    )
    document = parse_document(markdown, manifest)
    require(not document["errors"], f"frame capacity parser errors: {document['errors']}")
    plan = pptx_paginate.build_deck_plan(document, manifest)
    require(not [item for item in plan.diagnostics if item.severity == "error"],
            f"frame capacity plan diagnostics: {plan.diagnostics}")
    delivery = workdir / "delivery"
    completed = run_public_render(markdown, delivery, "frame-capacity-consistency")
    public_output = completed.stdout + completed.stderr
    require(completed.returncode == 0, f"frame capacity render failed: {public_output[:1024]}")
    pptx_path = delivery / "frame-capacity-consistency.pptx"
    reopened = Presentation(pptx_path)
    vectors: dict[str, object] = {}

    for layout, slot_id in (("contents", "body"), ("title-content", "body"), ("code", "code")):
        physical_slides = [slide for slide in plan.slides if slide.layout == layout]
        page_evidence: list[dict[str, object]] = []
        for physical in physical_slides:
            plan_sequence = [projection for fragment in physical.fragments
                             for projection in fragment_paragraph_sequence(fragment)]
            plan_paragraphs = [{
                "role": projection.role,
                "authored_text": projection.authored_text,
                "visible_text": projection.visible_text,
            } for projection in plan_sequence]
            names = [f"school-pptx:{slot_id}"]
            if layout == "code" and any(item["role"] == "heading" for item in plan_paragraphs):
                names.insert(0, "school-pptx:code-heading")
            frames: list[dict[str, object]] = []
            paragraph_offset = 0
            for name in names:
                shape = next(shape for shape in reopened.slides[physical.physical_index].shapes
                             if shape.name == name)
                paragraphs: list[dict[str, object]] = []
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs)
                    if not text:
                        continue
                    role = "heading" if name.endswith("code-heading") else (
                        "code" if name.endswith(":code") else plan_paragraphs[paragraph_offset]["role"]
                    )
                    run_font = next(run.font.size.pt for run in paragraph.runs if run.font.size is not None)
                    spacing = float(paragraph.line_spacing)
                    space_after = paragraph.space_after.pt if paragraph.space_after is not None else 0.0
                    width_emu = shape.width - round(
                        (shape.text_frame.margin_left.pt + shape.text_frame.margin_right.pt)
                        * pptx_paginate.EMU_PER_POINT
                    )
                    measurement = pptx_paginate.TextMeasure(
                        line_spacing=spacing, paragraph_spacing_points=space_after
                    ).measure(
                        text, width_emu=width_emu,
                        font_size=run_font, font_size_min=run_font, font_size_max=run_font,
                        kind="code" if role == "code" else "paragraph",
                    )
                    paragraphs.append({
                        "role": role,
                        "text": text,
                        "font_size": run_font,
                        "display_lines": measurement.display_lines,
                        "line_spacing": spacing,
                        "space_after_points": space_after,
                        "runs": [{
                            "text": run.text,
                            "bold": run.font.bold is True,
                            "highlight": "<a:highlight>" in run._r.xml,
                        } for run in paragraph.runs],
                    })
                    paragraph_offset += 1
                frames.append({
                    "shape_name": name,
                    "geometry_emu": {
                        "x": shape.left, "y": shape.top,
                        "width": shape.width, "height": shape.height,
                    },
                    "margins_points": {
                        "left": shape.text_frame.margin_left.pt,
                        "right": shape.text_frame.margin_right.pt,
                        "top": shape.text_frame.margin_top.pt,
                        "bottom": shape.text_frame.margin_bottom.pt,
                    },
                    "paragraphs": paragraphs,
                })
            reopened_sequence = [
                (paragraph["role"], paragraph["text"])
                for frame in frames for paragraph in frame["paragraphs"]
            ]
            require([(item["role"], item["visible_text"]) for item in plan_paragraphs] == reopened_sequence,
                    f"{layout} plan/reopen paragraph sequence changed")
            for frame in frames:
                full_height = sum(
                    paragraph["display_lines"] * paragraph["font_size"] * paragraph["line_spacing"]
                    + paragraph["space_after_points"]
                    for paragraph in frame["paragraphs"]
                )
                geometry = frame["geometry_emu"]
                margins = frame["margins_points"]
                effective_height = geometry["height"] / pptx_paginate.EMU_PER_POINT
                effective_height -= margins["top"] + margins["bottom"]
                require(full_height <= effective_height + 0.001,
                        f"{layout} reopened full frame exceeds effective geometry")
            if layout == "code" and len(frames) == 2:
                heading_geometry, code_geometry = frames[0]["geometry_emu"], frames[1]["geometry_emu"]
                frozen_height = int(dict(physical.slot_values)["geometry.code.heading_height_emu"])
                require(heading_geometry["height"] == frozen_height,
                        "dedicated code heading shape height differs from frozen plan")
                require(heading_geometry["y"] + heading_geometry["height"] == code_geometry["y"]
                        and code_geometry["height"] > 0,
                        "dedicated code remaining geometry is invalid")
            page_evidence.append({
                "physical_index": physical.physical_index,
                "plan_paragraphs": plan_paragraphs,
                "frames": frames,
                "slot_geometry_emu": dict(next(
                    item for item in manifest["layouts"][layout]["slots"] if item["id"] == slot_id
                )["geometry"]),
            })
        vectors[layout] = {
            "physical_pages": len(physical_slides),
            "pages": page_evidence,
        }

    contents_pages = [slide for slide in plan.slides if slide.layout == "contents"]
    require(len(contents_pages) > 1, "ten contents entries did not paginate beyond reviewed six-entry layout")
    contents_items = [item for slide in contents_pages for fragment in slide.fragments for item in fragment.items]
    require(contents_items[:5] == [f"第{index}项" for index in range(1, 6)],
            "contents unnumbered entries changed")
    require(all(dict(slide.selected_font_sizes).get("body") == 32.0 for slide in contents_pages),
            "contents frozen font is not 32pt")
    code_pages = [slide for slide in plan.slides if slide.layout == "code"]
    require(code_pages and all(dict(slide.selected_font_sizes).get("code") == 14.0 for slide in code_pages),
            "dedicated code frozen font is not 14pt")
    require(all(
        all(paragraph["font_size"] == 14.0 for frame in page["frames"] for paragraph in frame["paragraphs"])
        for page in vectors["code"]["pages"]
    ), "dedicated code reopen font is not 14pt")
    styled = [
        (paragraph["role"], paragraph["runs"])
        for page in vectors["title-content"]["pages"]
        for frame in page["frames"] for paragraph in frame["paragraphs"]
    ]
    require({role for role, runs in styled if any(run["bold"] for run in runs)} >= {"heading", "paragraph", "list"}
            and {role for role, runs in styled if any(run["highlight"] for run in runs)} >= {"heading", "paragraph", "list"},
            "bold/highlight heading, paragraph, and list vectors did not reopen with rich runs")

    overflow_source = workdir / "code-heading-overflow.md"
    overflow_source.write_text(
        "---\n"
        'title: "代码标题溢出"\n'
        'theme: "standard-school"\n'
        "---\n"
        '::: slide {layout="code"}\n'
        "## 代码标题溢出\n"
        "### " + "超" * 1400 + "\n"
        "```text\nprint('bounded')\n```\n"
        ":::\n",
        encoding="utf-8",
    )
    overflow_document = parse_document(overflow_source, manifest)
    require(not overflow_document["errors"], "overflow heading parser errors")
    overflow_plan = pptx_paginate.build_deck_plan(overflow_document, manifest)
    overflow_diagnostics = [item.to_projection() for item in overflow_plan.diagnostics]
    require(any(item["code"] == "TEXT_BLOCK_OVERFLOW" and item["severity"] == "error"
                for item in overflow_diagnostics),
            "overflow heading did not return the bounded semantic overflow diagnostic")
    overflow_code_pages = [slide for slide in overflow_plan.slides if slide.layout == "code"]
    code_slot_height = next(
        item for item in manifest["layouts"]["code"]["slots"] if item["id"] == "code"
    )["geometry"]["height"]
    frozen_heading_heights = [
        int(dict(slide.slot_values)["geometry.code.heading_height_emu"])
        for slide in overflow_code_pages
    ]
    require(overflow_code_pages and all(0 < height < code_slot_height for height in frozen_heading_heights),
            "overflow heading did not freeze bounded positive code geometry")

    overflow_delivery = workdir / "overflow-delivery"
    overflow_completed = run_public_render(overflow_source, overflow_delivery, "code-heading-overflow")
    overflow_output = overflow_completed.stdout + overflow_completed.stderr
    require(overflow_completed.returncode != 0 and "TEXT_BLOCK_OVERFLOW" in overflow_output
            and "PPTX_RENDER_FAILED" not in overflow_output,
            "overflow heading bypassed the structured best-effort diagnostic path")
    require("Traceback" not in overflow_output and len(overflow_output.encode("utf-8")) < 8192,
            "overflow heading public failure was unbounded")
    require(command_owned_files(overflow_delivery) == set(),
            "overflow heading published best-effort current artifacts")
    overflow_geometries: list[dict[str, int]] = []
    for heading_height in frozen_heading_heights:
        overflow_geometries.append({
            "heading_height": heading_height,
            "code_height": code_slot_height - heading_height,
            "slot_height": code_slot_height,
        })

    code_slide = code_pages[0]
    fail_closed_codes: list[str] = []
    for mutation, expected_code in (
        (replace(code_slide, selected_font_sizes=()), "PPTX_PLAN_FONT_MISSING"),
        (replace(code_slide, slot_values=()), "PPTX_PLAN_TYPOGRAPHY_MISSING"),
    ):
        mutated = PhysicalDeckPlan(tuple(mutation if slide is code_slide else slide for slide in plan.slides),
                                   plan.diagnostics, plan.logical_to_physical)
        try:
            pptx_emit.emit_deck(mutated, manifest, TEMPLATE_PATH, workdir / f"{expected_code}.pptx")
        except pptx_emit.PptxEmitError as exc:
            require(exc.code == expected_code, f"dedicated code fail-closed mismatch: {exc.code}")
            fail_closed_codes.append(exc.code)
        else:
            raise GateFailure(f"dedicated code mutation emitted: {expected_code}")
    return {
        "public_exit": completed.returncode,
        "output_bytes": len(public_output.encode("utf-8")),
        "artifact_paths": [str(delivery / markdown.name), str(pptx_path)],
        "vectors": vectors,
        "overflow_heading": {
            "public_exit": overflow_completed.returncode,
            "output_bytes": len(overflow_output.encode("utf-8")),
            "diagnostics": overflow_diagnostics,
            "frozen_heading_heights": frozen_heading_heights,
            "geometries": overflow_geometries,
            "artifact_names": sorted(path.name for path in overflow_delivery.iterdir()),
        },
        "fail_closed_codes": fail_closed_codes,
    }


GATES["frame-capacity-consistency"] = frame_capacity_consistency_gate


def editable_objects_gate(workdir: Path) -> dict[str, object]:
    try:
        import pptx_emit
        from pptx_model import BlockFragment, PhysicalDeckPlan, PhysicalSlide
        from pptx_objects import PptxObjectError, normalize_rich_text
    except ImportError as exc:
        raise GateFailure("PPTX_DEPENDENCY_MISSING") from exc

    require(normalize_rich_text("甲**粗体**==高亮==乙") == (
        ("plain", "甲"), ("bold", "粗体"), ("highlight", "高亮"), ("plain", "乙")
    ), "D-15 adjacent span normalization")
    try:
        normalize_rich_text("**重叠 ==样式** 失败==")
    except PptxObjectError as exc:
        require(str(exc) == "PPTX_INLINE_SPAN_OVERLAP", "overlap diagnostic changed")
    else:
        raise GateFailure("nested inline span did not fail")

    manifest = load_manifest(SKILL_DIR)
    table_slot = next(slot for slot in manifest["layouts"]["table"]["slots"] if slot["id"] == "table")
    table_budget = table_slot["subregions"]["table_body"]["geometry"]["height"]
    table_row_heights = (table_budget // 2, table_budget - table_budget // 2)

    rich = BlockFragment(
        kind="paragraph", logical_index=0, block_index=0, fragment_index=0, source_line=1,
        text="甲**粗体**==高亮==乙",
    )
    table = BlockFragment(
        kind="table", logical_index=1, block_index=0, fragment_index=0, source_line=2,
        rows=(("列一", "列二"), ("数据", "内容")),
        metadata=(("table_name", ""), ("table_name_placeholder", "true"), ("repeat_header", "true")),
        row_heights_emu=table_row_heights,
    )
    timeline = BlockFragment(
        kind="timeline", logical_index=2, block_index=0, fragment_index=0, source_line=3,
        rows=(("2026", "启动", "准备"), ("2027", "交付", "完成")),
    )
    gallery = BlockFragment(
        kind="image", logical_index=3, block_index=0, fragment_index=0, source_line=4,
        metadata=(("authored_path", "media/plc-line.png"), ("caption", ""),
                  ("caption_placeholder", "true"), ("gallery_preset", "1"), ("gallery_item_index", "0")),
    )
    code = BlockFragment(
        kind="code", logical_index=4, block_index=0, fragment_index=0, source_line=5,
        text="print('原始换行')\nreturn True",
    )
    body_typography = (
        ("typography.body.margin_left_points", "0.0"),
        ("typography.body.margin_right_points", "0.0"),
        ("typography.body.margin_top_points", "0.0"),
        ("typography.body.margin_bottom_points", "0.0"),
        ("typography.body.line_spacing", "1.2"),
        ("typography.body.paragraph_spacing_points", "2.0"),
    )
    code_typography = tuple((key.replace("body", "code"), value) for key, value in body_typography)
    slides = (
        PhysicalSlide(
            0, 0, "title-content", "富文本", 0, 1, (rich,), "逐字说明",
            selected_font_sizes=(("body", 28.0),), affected_pages=(0,), slot_values=body_typography,
        ),
        PhysicalSlide(1, 1, "table", "原生表格", 0, 2, (table,), None,
                      selected_font_sizes=(("table", 14.0),), affected_pages=(1,)),
        PhysicalSlide(2, 2, "timeline", "原生时间线", 0, 3, (timeline,), None, affected_pages=(2,)),
        PhysicalSlide(3, 3, "gallery", "原生图集", 0, 4, (gallery,), "图集备注", affected_pages=(3,)),
        PhysicalSlide(4, 4, "code", "代码", 0, 5, (code,), None,
                      selected_font_sizes=(("code", 14.0),), affected_pages=(4,),
                      slot_values=code_typography),
        PhysicalSlide(5, 5, "closing", "", 0, 6, (
            BlockFragment(kind="closing", logical_index=5, block_index=0, fragment_index=0, source_line=6),
        ), None, affected_pages=(5,)),
    )
    plan = PhysicalDeckPlan(slides)
    output = workdir / "editable-objects.pptx"
    pptx_emit.emit_deck(plan, manifest, TEMPLATE_PATH, output, media_root=SKILL_DIR / "fixtures")
    reopened = pptx_emit.require_dependencies()["pptx"].Presentation(output)
    require(len(reopened.slides) == len(slides), "editable deck slide count")
    require([slide.has_notes_slide for slide in reopened.slides] == [True, False, False, True, False, False],
            "D-16 notes relationships do not match plan")
    require(reopened.slides[0].notes_slide.notes_text_frame.text == "逐字说明", "notes text mismatch")
    require(reopened.slides[3].notes_slide.notes_text_frame.text == "图集备注", "derived notes text mismatch")

    entries = safe_package_entries(output)
    all_slides = b"\n".join(payload for name, payload in entries.items() if name.startswith("ppt/slides/slide"))
    require(all_slides.count(b"<a:tbl>") == 1, "D-11 native table count")
    require(all_slides.count(b"<p:grpSp>") >= 3, "D-12/D-13 editable group count")
    require(b"<a:highlight>" in all_slides and b"b=\"1\"" in all_slides, "D-15 run OOXML missing")
    require(b"**" not in all_slides and b"==" not in all_slides, "D-15 delimiters visible")
    require(b"school-pptx:table-name" in all_slides and b"school-pptx:gallery-caption:0" in all_slides,
            "D-14 empty editable placeholder missing")
    require(all_slides.count(b"<a:t></a:t>") >= 1, "D-14 empty placeholder has no empty text node")
    require(b"school-pptx:timeline-axis" in all_slides, "timeline axis missing")
    require(b"school-pptx:timeline-axis" not in entries["ppt/slides/slide3.xml"].split(b"<p:grpSp>", 1)[-1],
            "timeline axis entered node group")
    require(b"print('original" not in all_slides, "unexpected code rewrite")
    require("print('原始换行')\nreturn True" in reopened.slides[4].shapes[-1].text, "code text changed")

    pictures = [shape for shape in reopened.slides[3].shapes if getattr(shape, "shape_type", None) == 6]
    require(pictures and len(pictures[0].shapes) == 2, "gallery card did not reopen as native group")
    gallery_picture = next(shape for shape in pictures[0].shapes if getattr(shape, "shape_type", None) == 13)
    require(all(getattr(gallery_picture, crop) == 0 for crop in ("crop_left", "crop_top", "crop_right", "crop_bottom")),
            "contain picture crop changed")
    require(output.stat().st_size > 0, "editable deck is empty")
    return {
        "slides": len(slides), "native_tables": 1, "groups": all_slides.count(b"<p:grpSp>"),
        "notes_relationships": 2, "picture_crop": [0, 0, 0, 0],
    }


GATES["editable-objects"] = editable_objects_gate


def emit_structure_gate(workdir: Path) -> dict[str, object]:
    try:
        import pptx_emit
    except ImportError as exc:
        raise GateFailure("PPTX_DEPENDENCY_MISSING") from exc

    fixture = SKILL_DIR / "fixtures" / "school-pptx-full.md"
    manifest = load_manifest(SKILL_DIR)
    document = parse_document(fixture, manifest)
    require(not document["errors"], "canonical fixture parser errors")
    plan = pptx_paginate.build_deck_plan(document, manifest)
    require(plan.slides and plan.slides[-1].layout == "closing", "physical plan closing invalid")
    output = workdir / "full-fixture.pptx"
    evidence = pptx_emit.emit_deck(plan, manifest, TEMPLATE_PATH, output, media_root=fixture.parent)
    emitter_source = (SCRIPTS_DIR / "pptx_emit.py").read_text(encoding="utf-8")
    require("build_deck_plan" not in emitter_source and "pptx_paginate" not in emitter_source,
            "emitter attempted to repaginate frozen plan")
    require(evidence["physical_slides"] == len(plan.slides), "PPTX-01 emitted slide count evidence")
    require(evidence["transition_mode"] == "none", "PPTX-12 transition mode")
    validated = pptx_emit.validate_staged_package(output, expected_slides=len(plan.slides))
    require(validated["slides"] == len(plan.slides), "staged validator slide count")

    package_entries = safe_package_entries(output)
    slide_xml = {
        index: package_entries[f"ppt/slides/slide{index + 1}.xml"] for index in range(len(plan.slides))
    }
    all_slide_xml = b"\n".join(slide_xml.values())
    require(b"\xe5\x88\x9d\xe8\xaf\x86\xe9\x9a\x94\xe7\xa6\xbb\xe5\xbc\x80\xe5\x85\xb3" not in all_slide_xml,
            "template seed text leaked")
    require(b"<p:transition" not in all_slide_xml, "PPTX-12 unexpected transition XML")

    rel_namespace = {"r": "http://schemas.openxmlformats.org/package/2006/relationships"}
    layout_relationship = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
    notes_relationship = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide"
    image_relationship = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
    notes_indices: list[int] = []
    image_relationships = 0
    for index, physical in enumerate(plan.slides, start=1):
        rel_name = f"ppt/slides/_rels/slide{index}.xml.rels"
        relationships = ET.fromstring(package_entries[rel_name]).findall("r:Relationship", rel_namespace)
        layout_targets = [item.get("Target", "") for item in relationships if item.get("Type") == layout_relationship]
        require(len(layout_targets) == 1, f"slide {index} layout relationship count")
        expected_layout = manifest["layouts"][physical.layout]["pptx_layout"].split("/")[-1]
        require(layout_targets[0].endswith(expected_layout), f"slide {index} layout part mismatch")
        if any(item.get("Type") == notes_relationship for item in relationships):
            notes_indices.append(index - 1)
        image_relationships += sum(item.get("Type") == image_relationship for item in relationships)
    expected_notes = [slide.physical_index for slide in plan.slides if slide.notes_intent is not None]
    require(notes_indices == expected_notes, "PPTX-04 notes relationship mismatch")

    table_slides = [slide for slide in plan.slides if slide.layout == "table"]
    for physical in table_slides:
        xml = slide_xml[physical.physical_index]
        require(xml.count(b"<a:tbl>") == 1, "PPTX-04 native table per page")
        require(b"school-pptx:table-name" in xml, "empty/non-empty table-name shape missing")
        require(all(cell.encode("utf-8") in xml for cell in physical.fragments[0].rows[0]),
                "table repeated header missing")
    timeline_slides = [slide for slide in plan.slides if slide.layout == "timeline"]
    for physical in timeline_slides:
        xml = slide_xml[physical.physical_index]
        require(xml.count(b"<p:grpSp>") == len(physical.fragments[0].rows), "timeline group count")
        require(b"school-pptx:timeline-axis" in xml, "timeline axis missing")
    gallery_slides = [slide for slide in plan.slides if slide.layout == "gallery"]
    for physical in gallery_slides:
        require(len(physical.fragments) <= 4, "gallery capacity")
        xml = slide_xml[physical.physical_index]
        require(xml.count(b"<p:grpSp>") == len(physical.fragments), "gallery group count")
        require(xml.count(b"gallery-caption") == len(physical.fragments), "gallery caption count")
    code_slides = [slide for slide in plan.slides if slide.layout == "code"]
    for physical in code_slides:
        expected = "\n".join(fragment.text or "" for fragment in physical.fragments)
        reopened_text = "\n".join(
            shape.text for shape in pptx_emit.require_dependencies()["pptx"].Presentation(output)
            .slides[physical.physical_index].shapes if getattr(shape, "has_text_frame", False)
        )
        require(expected in reopened_text, "editable code round trip")

    reopened = pptx_emit.require_dependencies()["pptx"].Presentation(output)
    require(len(reopened.slides) == len(plan.slides), "PPTX-01 reopen slide count")
    require(sum(slide.layout == "closing" for slide in plan.slides) == 1, "closing plan count")
    last_relationships = ET.fromstring(
        package_entries[f"ppt/slides/_rels/slide{len(plan.slides)}.xml.rels"]
    ).findall("r:Relationship", rel_namespace)
    require(any(item.get("Target", "").endswith("slideLayout8.xml") for item in last_relationships),
            "closing did not use manifest part")
    for slide in reopened.slides:
        for shape in slide.shapes:
            if getattr(shape, "shape_type", None) == 13:
                require(not (shape.width == reopened.slide_width and shape.height == reopened.slide_height),
                        "whole-slide screenshot shortcut detected")
    require(image_relationships > 0, "native picture relationships missing")
    picture_namespace = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    }
    renderer_picture_names: list[str] = []
    for xml in slide_xml.values():
        root = ET.fromstring(xml)
        for picture in root.findall(".//p:pic", picture_namespace):
            properties = picture.find("./p:nvPicPr/p:cNvPr", picture_namespace)
            name = properties.get("name", "") if properties is not None else ""
            renderer_picture_names.append(name)
            require(
                picture.find("./p:spPr/a:effectLst/a:outerShdw", picture_namespace) is not None,
                f'renderer picture "{name}" lost a:effectLst/a:outerShdw',
            )
    require(renderer_picture_names and len(renderer_picture_names) == image_relationships,
            "renderer picture shadow inventory does not match slide image relationships")
    require(all(b"<a:srcRect" not in xml or b"<a:srcRect/>" in xml or b'l="0"' in xml
                for xml in slide_xml.values()),
            "picture crop is not zero")

    external = workdir / "external-relationship.pptx"
    with ZipFile(output) as source, ZipFile(external, "w", ZIP_DEFLATED) as target:
        for info in source.infolist():
            payload = source.read(info)
            if info.filename == "ppt/_rels/presentation.xml.rels":
                root = ET.fromstring(payload)
                relationship = root.findall("r:Relationship", rel_namespace)[0]
                relationship.set("Target", "https://invalid.example/template")
                relationship.set("TargetMode", "External")
                payload = ET.tostring(root, encoding="utf-8", xml_declaration=True)
            target.writestr(info.filename, payload)
    try:
        pptx_emit.validate_staged_package(external)
    except pptx_emit.PptxEmitError as exc:
        require(exc.code == "PPTX_PACKAGE_INVALID" and "Traceback" not in str(exc),
                "external relationship failure is not bounded")
    else:
        raise GateFailure("external relationship did not fail staged validation")
    return {
        "logical_slides": len(document["logical_slides"]),
        "physical_slides": len(plan.slides),
        "native_table_slides": len(table_slides),
        "timeline_slides": len(timeline_slides),
        "gallery_slides": len(gallery_slides),
        "notes_slides": len(notes_indices),
        "image_relationships": image_relationships,
        "renderer_picture_shadows": len(renderer_picture_names),
        "transition_mode": "none",
        "package": validated,
    }


GATES["emit-structure"] = emit_structure_gate


def run_public_render(input_path: Path, output_root: Path, stem: str) -> subprocess.CompletedProcess[str]:
    environment = os.environ.copy()
    environment["SCHOOL_PPTX_PYTHON"] = sys.executable
    completed = subprocess.run(
        [str(PUBLIC_CLI), "render", "--input", str(input_path), "--out-dir", str(output_root), "--stem", stem],
        cwd=SKILL_DIR.parent.parent,
        text=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        timeout=60,
        check=False,
        env=environment,
    )
    require("Traceback" not in completed.stdout + completed.stderr, "public render leaked traceback")
    require(len(completed.stdout + completed.stderr) < 200_000, "public render output is unbounded")
    return completed


def command_owned_files(root: Path) -> set[str]:
    return {
        path.relative_to(root).as_posix()
        for path in root.rglob("*")
        if path.is_file() or path.is_symlink()
    }


def delivery_tree_snapshot(root: Path) -> dict[str, tuple[str, int]]:
    if not root.exists():
        return {}
    return {
        path.relative_to(root).as_posix(): (hashlib.sha256(path.read_bytes()).hexdigest(), path.stat().st_size)
        for path in root.rglob("*")
        if path.is_file() and not path.is_symlink()
    }


def delivery_transaction_gate(workdir: Path) -> dict[str, object]:
    import pptx_render

    require(DELIVERY_FIXTURE.is_file(), "confirmed-assets fixture Markdown missing")
    fixture_root = DELIVERY_FIXTURE.parent
    input_root = workdir / "input"
    shutil.copytree(fixture_root, input_root)
    source = input_root / "confirmed-assets.md"
    unreferenced = input_root / "assets" / "unreferenced-input.png"
    unreferenced.write_bytes((SKILL_DIR / "fixtures" / "media" / "equipment-cell.png").read_bytes())
    unreferenced_before = hashlib.sha256(unreferenced.read_bytes()).hexdigest()

    signal_delivery = workdir / "first-signal-delivery"
    for signal_name in ("SIGINT", "SIGTERM"):
        environment = os.environ.copy()
        environment["SCHOOL_PPTX_PYTHON"] = sys.executable
        environment["PRESTO_CLEAN_DELIVERY_SIGNAL_BEFORE_RECORD"] = signal_name
        interrupted = subprocess.run(
            [str(PUBLIC_CLI), "render", "--input", str(source), "--out-dir", str(signal_delivery),
             "--stem", "confirmed-assets"],
            cwd=SKILL_DIR.parent.parent, text=True, capture_output=True, timeout=60, check=False, env=environment,
        )
        require(interrupted.returncode != 0 and not delivery_tree_snapshot(signal_delivery),
                f"first publish {signal_name} left partial current")
        require(not (signal_delivery / ".work").exists(), f"first publish {signal_name} leaked work")

    delivery = workdir / "delivery"
    sources = delivery / "sources"
    sources.mkdir(parents=True)
    source_sentinel = sources / "teacher-source.bin"
    source_sentinel.write_bytes(b"teacher-owned-source\x00unchanged")
    sources_before = delivery_tree_snapshot(sources)

    first = run_public_render(source, delivery, "confirmed-assets")
    require(first.returncode == 0 and "发布结果：first" in first.stdout,
            f"confirmed asset first publish failed: {first.stdout + first.stderr}")
    expected_current = {
        "confirmed-assets.md", "confirmed-assets.pptx", "assets/robot-arm.png", "sources/teacher-source.bin"
    }
    require(command_owned_files(delivery) == expected_current, "first publish leaked or omitted root files")
    require(not (delivery / "history").exists() and not (delivery / ".work").exists(),
            "first publish created history or leaked work")
    require((delivery / "assets" / "robot-arm.png").read_bytes()
            == (input_root / "assets" / "robot-arm.png").read_bytes(),
            "confirmed asset bytes changed")
    require(not (delivery / "assets" / unreferenced.name).exists(), "unreferenced input was copied")
    require(delivery_tree_snapshot(sources) == sources_before
            and hashlib.sha256(unreferenced.read_bytes()).hexdigest() == unreferenced_before,
            "sources or unreferenced input changed on first publish")

    current_paths = [delivery / "confirmed-assets.md", delivery / "confirmed-assets.pptx",
                     delivery / "assets" / "robot-arm.png"]
    before_noop = [(path.stat().st_ino, path.stat().st_mtime_ns, hashlib.sha256(path.read_bytes()).hexdigest())
                   for path in current_paths]
    identical = run_public_render(source, delivery, "confirmed-assets")
    after_noop = [(path.stat().st_ino, path.stat().st_mtime_ns, hashlib.sha256(path.read_bytes()).hexdigest())
                  for path in current_paths]
    require(identical.returncode == 0 and "发布结果：identical" in identical.stdout,
            "identical render did not report no-op")
    require(before_noop == after_noop and not (delivery / "history").exists(),
            "no-op touched current inode/mtime or created history")

    history = delivery / "history"
    current_markdown = delivery / "confirmed-assets.md"
    current_refs = pptx_render.managed_asset_references(current_markdown.read_bytes())
    for sequence in ("001", "003"):
        entry = history / sequence
        entry.mkdir(parents=True)
        shutil.copy2(current_markdown, entry / "confirmed-assets.md")
        shutil.copy2(delivery / "confirmed-assets.pptx", entry / "confirmed-assets.pptx")
        for relative in current_refs:
            archived_asset = entry / relative
            archived_asset.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(delivery / relative, archived_asset)
    existing_history = {
        sequence: delivery_tree_snapshot(history / sequence) for sequence in ("001", "003")
    }
    source.write_text(source.read_text(encoding="utf-8").replace(
        "DeliverySpec fixture", "DeliverySpec fixture revision 2"
    ), encoding="utf-8")
    replacement_asset = (SKILL_DIR / "fixtures" / "media" / "plc-line.png").read_bytes()
    (input_root / "assets" / "robot-arm.png").write_bytes(replacement_asset)
    changed = run_public_render(source, delivery, "confirmed-assets")
    require(changed.returncode == 0 and "发布结果：changed" in changed.stdout,
            f"changed publish failed: {changed.stdout + changed.stderr}")
    require((history / "004").is_dir() and not (history / "002").exists(),
            "history did not use max+1 sequence")
    require(all(delivery_tree_snapshot(history / sequence) == existing_history[sequence]
                for sequence in ("001", "003")), "pre-existing history changed")
    archived_markdown = history / "004" / "confirmed-assets.md"
    archived_refs = pptx_render.managed_asset_references(archived_markdown.read_bytes())
    require(archived_refs == ("assets/robot-arm.png",), "archived Markdown reference set changed")
    require(all((history / "004" / relative).is_file() for relative in archived_refs),
            "archived Markdown asset reference does not resolve in the same sequence")
    require((history / "004" / "assets" / "robot-arm.png").read_bytes()
            == fixture_root.joinpath("assets", "robot-arm.png").read_bytes(),
            "archived asset bytes differ from prior current")
    require(delivery_tree_snapshot(sources) == sources_before
            and hashlib.sha256(unreferenced.read_bytes()).hexdigest() == unreferenced_before,
            "sources or unreferenced input changed on revision")

    fault_results: dict[str, int] = {}
    for index, fault in enumerate(DELIVERY_FAULTS, start=1):
        source.write_text(re.sub(
            r'DeliverySpec fixture(?: revision [^\"]+)?',
            f"DeliverySpec fixture revision fault-{index}",
            source.read_text(encoding="utf-8"),
            count=1,
        ), encoding="utf-8")
        before_root = delivery_tree_snapshot(delivery)
        before_external = hashlib.sha256(unreferenced.read_bytes()).hexdigest()
        environment = os.environ.copy()
        environment["SCHOOL_PPTX_PYTHON"] = sys.executable
        environment["PRESTO_CLEAN_DELIVERY_FAULT"] = fault
        completed = subprocess.run(
            [str(PUBLIC_CLI), "render", "--input", str(source), "--out-dir", str(delivery),
             "--stem", "confirmed-assets"],
            cwd=SKILL_DIR.parent.parent, text=True, capture_output=True, timeout=60, check=False, env=environment,
        )
        require(completed.returncode != 0 and "Traceback" not in completed.stdout + completed.stderr,
                f"fault did not fail bounded: {fault}")
        require(delivery_tree_snapshot(delivery) == before_root,
                f"fault changed current/assets/sources/history: {fault}")
        require(hashlib.sha256(unreferenced.read_bytes()).hexdigest() == before_external,
                f"fault changed unreferenced input: {fault}")
        require(not (delivery / ".work").exists(), f"fault leaked work: {fault}")
        fault_results[fault] = completed.returncode

    before_failure = delivery_tree_snapshot(delivery)
    missing_asset = input_root / "assets" / "robot-arm.png"
    saved_asset = missing_asset.read_bytes()
    missing_asset.unlink()
    generation_failure = run_public_render(source, delivery, "confirmed-assets")
    require(generation_failure.returncode != 0 and delivery_tree_snapshot(delivery) == before_failure,
            "missing managed asset changed current/history")
    missing_asset.write_bytes(saved_asset)
    invalid_source = input_root / "invalid.md"
    invalid_source.write_text('---\ntitle: invalid\ntheme: standard-school\n---\n::: slide {layout="unknown"}\n:::\n',
                              encoding="utf-8")
    validation_failure = run_public_render(invalid_source, delivery, "confirmed-assets")
    require(validation_failure.returncode != 0 and delivery_tree_snapshot(delivery) == before_failure,
            "parser/plan failure published best-effort current")

    unknown = delivery / "unknown.log"
    unknown.write_bytes(b"unknown")
    unknown_before = delivery_tree_snapshot(delivery)
    rejected = run_public_render(source, delivery, "confirmed-assets")
    require(rejected.returncode != 0 and "OUTPUT_UNKNOWN" in rejected.stdout + rejected.stderr
            and delivery_tree_snapshot(delivery) == unknown_before, "unknown root entry was not fail-closed")
    unknown.unlink()
    legacy_media = delivery / "media"
    legacy_media.mkdir()
    (legacy_media / "legacy.png").write_bytes(saved_asset)
    legacy_before = delivery_tree_snapshot(delivery)
    legacy_rejected = run_public_render(source, delivery, "confirmed-assets")
    require(legacy_rejected.returncode != 0 and "OUTPUT_UNKNOWN" in legacy_rejected.stdout + legacy_rejected.stderr
            and delivery_tree_snapshot(delivery) == legacy_before, "legacy media was silently migrated")
    (legacy_media / "legacy.png").unlink()
    legacy_media.rmdir()
    outside = workdir / "outside-symlink-target.bin"
    outside.write_bytes(b"outside-unchanged")
    root_symlink = delivery / "rogue-link"
    root_symlink.symlink_to(outside)
    symlink_before = delivery_tree_snapshot(delivery)
    symlink_rejected = run_public_render(source, delivery, "confirmed-assets")
    require(symlink_rejected.returncode != 0 and "OUTPUT_UNKNOWN" in symlink_rejected.stdout + symlink_rejected.stderr
            and delivery_tree_snapshot(delivery) == symlink_before and outside.read_bytes() == b"outside-unchanged",
            "root symlink was followed or mutated")
    root_symlink.unlink()
    traversal = run_public_render(source, delivery, "../escape")
    require(traversal.returncode != 0 and "OUTPUT_STEM_INVALID" in traversal.stdout + traversal.stderr,
            "unsafe stem traversal was accepted")
    (delivery / ".work").mkdir()
    (delivery / ".work" / ".lock").write_bytes(b"held")
    locked_before = delivery_tree_snapshot(delivery)
    locked = run_public_render(source, delivery, "confirmed-assets")
    require(locked.returncode != 0 and delivery_tree_snapshot(delivery) == locked_before,
            "existing lock/work state was mutated")
    (delivery / ".work" / ".lock").unlink()
    (delivery / ".work").rmdir()

    return {
        "registry": {"required": DELIVERY_FAULTS, "called": tuple(fault_results), "dynamic_skips": 0},
        "first": first.returncode,
        "identical": identical.returncode,
        "changed": changed.returncode,
        "history_sequence": "004",
        "archived_assets": list(archived_refs),
        "faults": fault_results,
        "sources_unchanged": delivery_tree_snapshot(sources) == sources_before,
        "unreferenced_not_copied": not (delivery / "assets" / unreferenced.name).exists(),
        "best_effort_published": False,
        "manual_viewer_uat_claimed": False,
    }


def semantic_package_inventory(path: Path) -> dict[str, object]:
    import pptx_emit

    entries = safe_package_entries(path)
    presentation = pptx_emit.require_dependencies()["pptx"].Presentation(path)
    slide_payloads = [entries[f"ppt/slides/slide{index}.xml"] for index in range(1, len(presentation.slides) + 1)]
    relationship_payloads = [
        entries[f"ppt/slides/_rels/slide{index}.xml.rels"] for index in range(1, len(presentation.slides) + 1)
    ]
    relationship_namespace = {"r": "http://schemas.openxmlformats.org/package/2006/relationships"}
    layouts: list[str] = []
    notes = 0
    images = 0
    for payload in relationship_payloads:
        relationships = ET.fromstring(payload).findall("r:Relationship", relationship_namespace)
        layouts.extend(
            item.get("Target", "").split("/")[-1]
            for item in relationships
            if item.get("Type", "").endswith("/slideLayout")
        )
        notes += sum(item.get("Type", "").endswith("/notesSlide") for item in relationships)
        images += sum(item.get("Type", "").endswith("/image") for item in relationships)
    visible_text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    return {
        "slides": len(presentation.slides),
        "layouts": layouts,
        "native_tables": sum(payload.count(b"<a:tbl>") for payload in slide_payloads),
        "groups": sum(payload.count(b"<p:grpSp>") for payload in slide_payloads),
        "pictures": images,
        "notes": notes,
        "visible_text_sha256": hashlib.sha256(visible_text.encode("utf-8")).hexdigest(),
    }


GATES["delivery-transaction"] = delivery_transaction_gate


def cli_publication_gate(workdir: Path) -> dict[str, object]:
    output = workdir / "delivery"
    output.mkdir()
    completed = run_public_render(FIXTURE_PATH, output, "course-deck")
    require(completed.returncode == 0, f"public render failed: {completed.stdout}{completed.stderr}")
    lines = completed.stdout.splitlines()
    require(lines and lines[0] == "渲染成功", "success heading order changed")
    require(f"逻辑页：13；物理页：{CANONICAL_PHYSICAL_SLIDES}" in completed.stdout,
            "success pagination summary missing")
    require(completed.stdout.rstrip().endswith("校验结果：PASS"), "success validation footer missing")
    markdown = output / "course-deck.md"
    pptx = output / "course-deck.pptx"
    require(markdown.read_bytes() == FIXTURE_PATH.read_bytes(), "published Markdown bytes changed")
    require(pptx.stat().st_size > 0, "published PPTX is empty")
    require(command_owned_files(output) == {"course-deck.md", "course-deck.pptx"},
            "success public root contains non-pair artifacts")
    inventory = semantic_package_inventory(pptx)
    require(inventory["slides"] == CANONICAL_PHYSICAL_SLIDES and inventory["native_tables"] > 0,
            "published PPTX did not reopen structurally")
    help_result = subprocess.run([str(PUBLIC_CLI), "--help"], text=True, capture_output=True, check=False)
    require(help_result.returncode == 0 and
            "render --input <reviewed.md> --out-dir <delivery-dir> [--stem <name>]" in help_result.stdout,
            "literal render usage missing")
    return {
        "files": sorted(command_owned_files(output)),
        "inventory": inventory,
        "pagination_explanation": "43-07 freezes three ordered two-column pairs as three physical pages",
    }


def best_effort_gate(workdir: Path) -> dict[str, object]:
    missing_source = workdir / "missing-media.md"
    missing_source.write_bytes(FIXTURE_PATH.read_bytes())
    missing_output = workdir / "missing-output"
    missing_output.mkdir()
    completed = run_public_render(missing_source, missing_output, "missing")
    require(completed.returncode != 0, "missing-media render returned success")
    require("渲染完成但输入存在异常" in completed.stdout and "本次渲染不成功" in completed.stdout,
            "best-effort non-success copy missing")
    require("受影响逻辑页" in completed.stdout and "MEDIA_MISSING" in completed.stdout,
            "best-effort affected slides missing")
    require("渲染成功" not in completed.stdout, "best-effort printed success copy")
    require(command_owned_files(missing_output) == set() and not (missing_output / ".work").exists(),
            "best-effort escaped owned work evidence or leaked work")

    invalid_source = workdir / "invalid.md"
    invalid_source.write_text(
        '---\ntitle: "异常"\ntheme: "standard-school"\n---\n'
        '::: slide {layout="unknown-layout"}\n## 受影响页\n可恢复正文。\n:::\n',
        encoding="utf-8",
    )
    invalid_output = workdir / "invalid-output"
    invalid_output.mkdir()
    invalid = run_public_render(invalid_source, invalid_output, "invalid")
    require(invalid.returncode != 0 and "LAYOUT_UNKNOWN" in invalid.stdout and "受影响页" in invalid.stdout,
            "invalid Markdown did not publish bounded best-effort output")
    require(command_owned_files(invalid_output) == set() and not (invalid_output / ".work").exists(),
            "invalid Markdown published best-effort current or leaked work")
    return {
        "missing_media_current": sorted(command_owned_files(missing_output)),
        "invalid_current": sorted(command_owned_files(invalid_output)),
        "invalid_exit": invalid.returncode,
        "best_effort_published": False,
    }


def invoke_render_module(input_path: Path, output_root: Path, stem: str) -> tuple[int, str]:
    import pptx_render

    stdout = io.StringIO()
    arguments = argparse.Namespace(
        skill_dir=str(SKILL_DIR), input=str(input_path), out_dir=str(output_root), stem=stem,
    )
    with contextlib.redirect_stdout(stdout), contextlib.redirect_stderr(stdout):
        result = pptx_render.render_command(arguments)
    output = stdout.getvalue()
    require("Traceback" not in output and len(output) < 200_000, "in-process render failure is unbounded")
    return result, output


def assert_no_renderer_debris(root: Path) -> None:
    forbidden = [
        path for path in root.rglob("*")
        if path.name.endswith(".bak") or ".tmp-" in path.name or path.suffix in {".json", ".log"}
        or "manifest" in path.name.lower() or "debug" in path.name.lower() or "evidence" in path.name.lower()
    ]
    require(not forbidden, f"renderer debris remains: {[path.name for path in forbidden]}")


def publication_safety_gate(workdir: Path) -> dict[str, object]:
    import pptx_render

    old_pptx = b"old-pptx-sentinel"
    corruption_root = workdir / "corruption"
    corruption_root.mkdir()
    (corruption_root / "deck.md").write_bytes(b"old-markdown")
    (corruption_root / "deck.pptx").write_bytes(old_pptx)
    def corrupt(stream: object) -> None:
        stream.seek(0)
        stream.truncate(0)
        stream.write(b"corrupt")
        stream.flush()

    pptx_render.RENDER_PRE_VALIDATE_HOOK = corrupt
    try:
        result, output = invoke_render_module(FIXTURE_PATH, corruption_root, "deck")
    finally:
        pptx_render.RENDER_PRE_VALIDATE_HOOK = None
    require(result != 0 and "PPTX_PACKAGE_INVALID" in output, "staged corruption did not fail bounded")
    require((corruption_root / "deck.pptx").read_bytes() == old_pptx, "staged corruption replaced old PPTX")
    require((corruption_root / "deck.md").read_bytes() == b"old-markdown", "staged corruption replaced Markdown")
    assert_no_renderer_debris(corruption_root)

    crash_root = workdir / "crash-window"
    crash_root.mkdir()
    (crash_root / "deck.md").write_bytes(b"old-markdown")
    (crash_root / "deck.pptx").write_bytes(old_pptx)
    pptx_render.RENDER_BETWEEN_REPLACE_HOOK = lambda root: (_ for _ in ()).throw(OSError("injected"))
    try:
        result, output = invoke_render_module(FIXTURE_PATH, crash_root, "deck")
    finally:
        pptx_render.RENDER_BETWEEN_REPLACE_HOOK = None
    require(result != 0 and "OUTPUT_PAIR_INTERRUPTED" in output, "between-replace fault did not fail bounded")
    require((crash_root / "deck.md").read_bytes() == b"old-markdown", "handled fault did not restore old Markdown")
    require((crash_root / "deck.pptx").read_bytes() == old_pptx, "PPTX-last crash window replaced old PPTX")
    assert_no_renderer_debris(crash_root)

    exchange_root = workdir / "exchange"
    exchange_root.mkdir()
    attacker = workdir / "attacker"
    attacker.mkdir()
    sentinel = attacker / "sentinel"
    sentinel.write_bytes(b"caller")
    held = workdir / "held-exchange"

    def exchange(root: Path) -> None:
        root.rename(held)
        root.symlink_to(attacker, target_is_directory=True)

    pptx_render.RENDER_PRE_PUBLISH_HOOK = exchange
    try:
        result, output = invoke_render_module(FIXTURE_PATH, exchange_root, "deck")
    finally:
        pptx_render.RENDER_PRE_PUBLISH_HOOK = None
    require(result != 0 and "OUTPUT_ROOT_CHANGED" in output, "output-root exchange was not blocked")
    require(sentinel.read_bytes() == b"caller" and not (attacker / "deck.pptx").exists(),
            "output-root exchange modified attacker target")
    assert_no_renderer_debris(held)

    collision_root = workdir / "collisions"
    collision_root.mkdir()
    outside = workdir / "outside"
    outside.write_bytes(b"outside")
    (collision_root / "link.pptx").symlink_to(outside)
    result, output = invoke_render_module(FIXTURE_PATH, collision_root, "link")
    require(result != 0 and "OUTPUT_PARTIAL" in output and outside.read_bytes() == b"outside",
            "final symlink collision was unsafe")
    (collision_root / "directory.pptx").mkdir()
    result, output = invoke_render_module(FIXTURE_PATH, collision_root, "directory")
    require(result != 0 and "OUTPUT_UNKNOWN" in output, "directory collision was accepted")
    result, output = invoke_render_module(FIXTURE_PATH, collision_root, "../escape")
    require(result != 0 and "OUTPUT_STEM_INVALID" in output and outside.read_bytes() == b"outside",
            "stem traversal was accepted")
    same = workdir / "same.md"
    same.write_bytes(FIXTURE_PATH.read_bytes())
    result, output = invoke_render_module(same, workdir, "same")
    require(result != 0 and ("OUTPUT_UNKNOWN" in output or "OUTPUT_PARTIAL" in output)
            and same.read_bytes() == FIXTURE_PATH.read_bytes(),
            "input/output identity was accepted")
    assert_no_renderer_debris(collision_root)

    host_python = sys.executable
    require(Path(host_python).is_file(), "host python missing for dependency gate")
    dependency_root = workdir / "dependency"
    dependency_root.mkdir()
    (dependency_root / "deck.md").write_bytes(b"old-markdown")
    (dependency_root / "deck.pptx").write_bytes(old_pptx)
    environment = os.environ.copy()
    yaml_root = next(path.parent.parent for path in Path.home().glob(".cache/uv/archive-v0/*/yaml/__init__.py"))
    environment["PYTHONPATH"] = str(yaml_root)
    isolated_python = workdir / "python-no-site"
    isolated_python.write_text(f'#!/bin/sh\nexec "{host_python}" -S "$@"\n', encoding="utf-8")
    isolated_python.chmod(0o700)
    environment["SCHOOL_PPTX_PYTHON"] = str(isolated_python)
    dependency = subprocess.run(
        [str(PUBLIC_CLI), "render", "--input", str(FIXTURE_PATH), "--out-dir", str(dependency_root), "--stem", "deck"],
        text=True, capture_output=True, check=False, env=environment, timeout=30,
    )
    require(dependency.returncode != 0 and "PPTX_DEPENDENCY_MISSING" in dependency.stdout + dependency.stderr,
            "dependency absence did not fail bounded")
    require("Traceback" not in dependency.stdout + dependency.stderr and
            (dependency_root / "deck.pptx").read_bytes() == old_pptx,
            "dependency failure leaked traceback or replaced old PPTX")
    assert_no_renderer_debris(dependency_root)
    return {"staged_corruption": "blocked", "crash_window": "old-bundle-restored", "exchange": "blocked"}


def publication_descriptor_race_gate(workdir: Path) -> dict[str, object]:
    import pptx_render

    normal_root = workdir / "normal"
    normal_root.mkdir()
    descriptor_evidence: dict[str, object] = {}

    def inspect_descriptors(destination: object, duplicate: object) -> None:
        _, identity, owner_fd = destination.temporary["pptx"]
        owner_flags = fcntl.fcntl(owner_fd, fcntl.F_GETFL) & os.O_ACCMODE
        duplicate_flags = fcntl.fcntl(duplicate.fileno(), fcntl.F_GETFL) & os.O_ACCMODE
        owner_status = os.fstat(owner_fd)
        duplicate_status = os.fstat(duplicate.fileno())
        probe = os.fdopen(os.dup(duplicate.fileno()), "r+b")
        probe.write(b"descriptor-probe")
        probe.seek(0)
        require(probe.read() == b"descriptor-probe", "duplicate is not independently seekable/readable/writable")
        probe.close()
        require(os.fstat(owner_fd).st_ino == owner_status.st_ino, "closing duplicate closed owner")
        require(owner_flags == duplicate_flags == os.O_RDWR, "PPTX descriptors are not read-write")
        require((owner_status.st_dev, owner_status.st_ino) ==
                (duplicate_status.st_dev, duplicate_status.st_ino) == identity,
                "owner and duplicate inode differ")
        descriptor_evidence.update({"owner_fd": owner_fd, "duplicate_fd": duplicate.fileno(), "identity": identity})

    pptx_render.RENDER_PRE_SAVE_HOOK = inspect_descriptors
    try:
        result, output = invoke_render_module(FIXTURE_PATH, normal_root, "deck")
    finally:
        pptx_render.RENDER_PRE_SAVE_HOOK = None
    require(result == 0, f"descriptor-only canonical render failed: {output}")
    safe_package_entries(normal_root / "deck.pptx")
    import pptx_emit
    require(len(pptx_emit.require_dependencies()["pptx"].Presentation(normal_root / "deck.pptx").slides) > 0,
            "published descriptor deck did not reopen")
    for key in ("owner_fd", "duplicate_fd"):
        try:
            os.fstat(int(descriptor_evidence[key]))
        except OSError:
            pass
        else:
            raise GateFailure(f"{key} remained open after render")

    race_root = workdir / "race"
    race_root.mkdir()
    sentinel = workdir / "outside-sentinel.bin"
    sentinel_bytes = b"caller-owned-sentinel\x00bytes"
    sentinel.write_bytes(sentinel_bytes)
    sentinel_hash = hashlib.sha256(sentinel_bytes).hexdigest()
    race_fds: dict[str, int] = {}
    attacker_name: dict[str, Path | str] = {}

    def exchange(destination: object, duplicate: object) -> None:
        name, _, owner_fd = destination.temporary["pptx"]
        race_fds.update({"owner_fd": owner_fd, "duplicate_fd": duplicate.fileno()})
        attacker_name["name"] = name
        attacker_name["root"] = destination.output_root
        os.unlink(name, dir_fd=destination.root_fd)
        os.symlink(str(sentinel), name, dir_fd=destination.root_fd)

    pptx_render.RENDER_PRE_SAVE_HOOK = exchange
    try:
        result, output = invoke_render_module(FIXTURE_PATH, race_root, "deck")
    finally:
        pptx_render.RENDER_PRE_SAVE_HOOK = None
    require(result != 0 and "OUTPUT_TEMP_CHANGED" in output, "staged symlink exchange was not blocked")
    require(sentinel.read_bytes() == sentinel_bytes and sentinel.stat().st_size == len(sentinel_bytes)
            and hashlib.sha256(sentinel.read_bytes()).hexdigest() == sentinel_hash,
            "external sentinel changed during descriptor race")
    retained = Path(attacker_name["root"]) / str(attacker_name["name"])
    require(retained.is_symlink() and retained.resolve() == sentinel.resolve(), "attacker symlink was not retained")
    for key, descriptor in race_fds.items():
        try:
            os.fstat(descriptor)
        except OSError:
            pass
        else:
            raise GateFailure(f"race {key} remained open")
    regular_debris = [path for path in race_root.iterdir() if path.is_file() and not path.is_symlink()]
    require(not regular_debris and not (race_root / "deck.pptx").exists() and not (race_root / "deck.md").exists(),
            "renderer-owned publication debris remains after race")
    return {
        "descriptor_access": "O_RDWR",
        "same_inode": True,
        "zip_and_reopen": "PASS",
        "race_code": "OUTPUT_TEMP_CHANGED",
        "sentinel_sha256": sentinel_hash,
        "attacker_symlink_retained": True,
        "fds_closed": True,
        "renderer_owned_inode_reclaimed": True,
        "renderer_owned_regular_debris": [],
    }


def _media_markdown(media_name: str) -> str:
    return (
        '---\ntitle: "对象错误"\ntheme: "standard-school"\n---\n'
        '::: slide {layout="image-text"}\n## 媒体页\n对象错误验证正文。\n\n'
        f'![本地媒体]({media_name})\n:::\n'
    )


def _png_chunk(kind: bytes, payload: bytes) -> bytes:
    return len(payload).to_bytes(4, "big") + kind + payload + zlib.crc32(kind + payload).to_bytes(4, "big")


def _minimal_png(width: int, height: int, *, invalid_idat: bool = False) -> bytes:
    ihdr = width.to_bytes(4, "big") + height.to_bytes(4, "big") + bytes((8, 2, 0, 0, 0))
    chunks = [_png_chunk(b"IHDR", ihdr)]
    if invalid_idat:
        payload = b"not-a-valid-zlib-stream"
        chunks.append(len(payload).to_bytes(4, "big") + b"IDAT" + payload + b"\x00\x00\x00\x00")
    chunks.append(_png_chunk(b"IEND", b""))
    return b"\x89PNG\r\n\x1a\n" + b"".join(chunks)


def table_header_only_gate(workdir: Path) -> dict[str, object]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise GateFailure("PPTX_DEPENDENCY_MISSING") from exc

    source = workdir / "header-only.md"
    source.write_text(
        '---\ntitle: "仅表头表格"\ntheme: "standard-school"\n---\n'
        '::: slide {layout="table"}\n## 字段定义\n| 字段 | 说明 |\n|---|---|\n:::\n',
        encoding="utf-8",
    )
    manifest = load_manifest(SKILL_DIR)
    document = parse_document(source, manifest)
    require(not document["errors"], f"header-only parser errors: {document['errors']}")
    plan = pptx_paginate.build_deck_plan(document, manifest)
    require(not plan.diagnostics, f"header-only pagination diagnostics: {plan.diagnostics}")
    table_slides = [slide for slide in plan.slides if slide.layout == "table"]
    require(len(table_slides) == 1, "header-only table did not produce exactly one table physical page")
    fragment = table_slides[0].fragments[0]
    require(len(fragment.rows) == len(fragment.row_heights_emu) == 1,
            "header-only frozen rows and heights must each contain exactly one item")
    require(fragment.row_heights_emu[0] > 0, "header-only frozen row height must be positive")

    output_root = workdir / "delivery"
    completed = run_public_render(source, output_root, "header-only")
    require(completed.returncode == 0, f"header-only public render failed: {completed.stdout + completed.stderr}")
    require({path.name for path in output_root.iterdir()} == {"header-only.md", "header-only.pptx"},
            "header-only public render leaked sidecars")
    reopened = Presentation(output_root / "header-only.pptx")
    slide = reopened.slides[table_slides[0].physical_index]
    tables = [shape.table for shape in slide.shapes if getattr(shape, "has_table", False)]
    require(len(tables) == 1 and len(tables[0].rows) == 1,
            "header-only output is not one editable native table row")
    require(tables[0].rows[0].height == fragment.row_heights_emu[0],
            "header-only reopened row height differs from frozen plan")
    placeholders = [shape for shape in slide.shapes if shape.name == "school-pptx:table-name"]
    require(len(placeholders) == 1 and placeholders[0].text == "",
            "header-only empty table-name placeholder is missing or contains prompt text")
    return {
        "table_physical_pages": 1,
        "rows": len(fragment.rows),
        "row_heights_emu": list(fragment.row_heights_emu),
        "native_table_rows": len(tables[0].rows),
        "empty_table_name": True,
        "public_exit": completed.returncode,
    }


def object_error_bounded_gate(workdir: Path) -> dict[str, object]:
    import pptx_emit
    from pptx_objects import PptxObjectError

    unknown = pptx_emit.map_object_error(PptxObjectError("private path /tmp/secret"))
    require(unknown.code == "PPTX_OBJECT_INVALID" and "/tmp/secret" not in unknown.message,
            "unknown object error was not normalized")

    vectors = (
        ("bomb-error", "bomb-error.png", _minimal_png(20_000, 10_000), "PPTX_MEDIA_PIXEL_LIMIT"),
        ("bomb-warning", "bomb-warning.png", _minimal_png(10_000, 10_000), "PPTX_MEDIA_PIXEL_LIMIT"),
        ("unidentified", "unidentified.png", b"not-an-image", "PPTX_MEDIA_FORMAT_INVALID"),
        ("verify-decode", "verify-decode.png", _minimal_png(1, 1, invalid_idat=True),
         "PPTX_MEDIA_FORMAT_INVALID"),
    )

    evidence: dict[str, object] = {}
    for label, media_name, payload, expected_code in vectors:
        vector_root = workdir / label
        vector_root.mkdir()
        media_path = vector_root / media_name
        media_path.write_bytes(payload)
        source = vector_root / "deck.md"
        source.write_text(_media_markdown(media_name), encoding="utf-8")
        output_root = vector_root / "delivery"
        output_root.mkdir()
        old_pptx = b"old-pptx-object-error-sentinel"
        (output_root / "deck.md").write_bytes(b"old-markdown")
        (output_root / "deck.pptx").write_bytes(old_pptx)
        completed = run_public_render(source, output_root, "deck")
        output = completed.stdout + completed.stderr
        require(completed.returncode != 0 and expected_code in output,
                f"{label} did not return {expected_code}: {output[:1024]}")
        require(len(output.encode("utf-8")) < 8 * 1024, f"{expected_code} output exceeded 8 KiB")
        forbidden_paths = (str(media_path), str(vector_root), str(workdir), str(SKILL_DIR),
                           str(Path.home() / ".cache" / "uv"))
        require("Traceback" not in output and not any(path in output for path in forbidden_paths),
                f"{expected_code} leaked traceback or absolute workdir")
        require("cannot identify image file" not in output.lower() and "decompressionbomb" not in output.lower()
                and "image size" not in output.lower() and "pillow" not in output.lower(),
                f"{expected_code} leaked Pillow details")
        require((output_root / "deck.pptx").read_bytes() == old_pptx,
                f"{expected_code} replaced the old PPTX")
        require((output_root / "deck.md").read_bytes() == b"old-markdown",
                f"{expected_code} replaced old Markdown")
        assert_no_renderer_debris(output_root)
        evidence[label] = {
            "code": expected_code,
            "exit": completed.returncode,
            "output_bytes": len(output.encode("utf-8")),
            "old_pptx_preserved": True,
        }
    return evidence


def _valid_png_bytes() -> bytes:
    from PIL import Image

    payload = io.BytesIO()
    Image.new("RGB", (1, 1), (20, 80, 140)).save(payload, format="PNG")
    return payload.getvalue()


def _embedded_media_hashes(path: Path) -> list[str]:
    with ZipFile(path) as package:
        return sorted(
            hashlib.sha256(package.read(name)).hexdigest()
            for name in package.namelist()
            if name.startswith("ppt/media/")
        )


def media_descriptor_binding_gate(workdir: Path) -> dict[str, object]:
    try:
        from pptx import Presentation
        import pptx_objects
        import pptx_render
    except ImportError as exc:
        raise GateFailure("PPTX_DEPENDENCY_MISSING") from exc

    original_payload = _valid_png_bytes()
    original_hash = hashlib.sha256(original_payload).hexdigest()
    replacement_payload = _minimal_png(20_000, 10_000)
    replacement_hash = hashlib.sha256(replacement_payload).hexdigest()

    race_root = workdir / "replace-after-validate"
    race_root.mkdir()
    media_path = race_root / "image.png"
    media_path.write_bytes(original_payload)
    replacement_path = race_root / "replacement.png"
    replacement_path.write_bytes(replacement_payload)
    source = race_root / "deck.md"
    source.write_text(_media_markdown("image.png"), encoding="utf-8")
    delivery = race_root / "delivery"
    hook_calls: list[str] = []

    def replace_after_validation(reference: object) -> None:
        hook_calls.append("called")
        os.replace(replacement_path, media_path)

    pptx_objects.MEDIA_AFTER_VALIDATE_HOOK = replace_after_validation
    try:
        race_exit, race_output = invoke_render_module(source, delivery, "deck")
    finally:
        pptx_objects.MEDIA_AFTER_VALIDATE_HOOK = None
    require(race_exit == 0 and hook_calls == ["called"],
            f"validated-byte replacement render failed: {race_output}")
    embedded_hashes = _embedded_media_hashes(delivery / "deck.pptx")
    require(original_hash in embedded_hashes and replacement_hash not in embedded_hashes,
            "replacement path changed embedded validated media")
    reopened = Presentation(delivery / "deck.pptx")
    pictures = [shape for slide in reopened.slides for shape in slide.shapes if shape.shape_type == 13]
    require(len(pictures) == 1 and all(
        value == 0 for value in (
            pictures[0].crop_left, pictures[0].crop_top, pictures[0].crop_right, pictures[0].crop_bottom
        )
    ), "validated-byte vector lost native zero-crop picture")

    absolute_root = workdir / "absolute"
    absolute_root.mkdir()
    absolute_media = absolute_root / "absolute.png"
    absolute_media.write_bytes(original_payload)
    absolute_source = absolute_root / "deck.md"
    absolute_source.write_text(_media_markdown(str(absolute_media.resolve())), encoding="utf-8")
    absolute_delivery = absolute_root / "delivery"
    absolute_result = run_public_render(absolute_source, absolute_delivery, "deck")
    require(absolute_result.returncode == 0, f"absolute PNG failed: {absolute_result.stdout + absolute_result.stderr}")
    absolute_hashes = _embedded_media_hashes(absolute_delivery / "deck.pptx")
    require(original_hash in absolute_hashes, "absolute PNG embedded hash changed")

    symlink_evidence: dict[str, object] = {}
    for label, authored_builder in (
        ("relative-intermediate", lambda root, target: (root / "linked").symlink_to(target.parent, target_is_directory=True) or "linked/image.png"),
        ("relative-final", lambda root, target: (root / "linked.png").symlink_to(target) or "linked.png"),
        ("absolute-final", lambda root, target: (root / "absolute-linked.png").symlink_to(target) or str(root / "absolute-linked.png")),
    ):
        vector_root = workdir / label
        vector_root.mkdir()
        target_root = vector_root / "target"
        target_root.mkdir()
        target = target_root / "image.png"
        target.write_bytes(original_payload)
        authored = authored_builder(vector_root, target)
        vector_source = vector_root / "deck.md"
        vector_source.write_text(_media_markdown(str(authored)), encoding="utf-8")
        vector_delivery = vector_root / "delivery"
        vector_delivery.mkdir()
        old_target = b"old-pptx-media-symlink"
        (vector_delivery / "deck.md").write_bytes(b"old-markdown")
        (vector_delivery / "deck.pptx").write_bytes(old_target)
        old_target_before_hash = hashlib.sha256(old_target).hexdigest()
        result = run_public_render(vector_source, vector_delivery, "deck")
        output = result.stdout + result.stderr
        require(result.returncode != 0 and "PPTX_MEDIA_PATH_INVALID" in output,
                f"{label} symlink did not fail closed: {output[:1024]}")
        require((vector_delivery / "deck.pptx").read_bytes() == old_target
                and (vector_delivery / "deck.md").read_bytes() == b"old-markdown",
                f"{label} symlink replaced public artifacts")
        require(len(output.encode("utf-8")) < 8192 and "Traceback" not in output
                and str(target) not in output and str(vector_root) not in output,
                f"{label} symlink failure was not bounded")
        old_target_after_hash = hashlib.sha256((vector_delivery / "deck.pptx").read_bytes()).hexdigest()
        symlink_evidence[label] = {
            "public_exit": result.returncode,
            "code": "PPTX_MEDIA_PATH_INVALID",
            "output_bytes": len(output.encode("utf-8")),
            "old_target_before_hash": old_target_before_hash,
            "old_target_after_hash": old_target_after_hash,
            "old_target_preserved": old_target_before_hash == old_target_after_hash,
        }

    missing_root = workdir / "runtime-missing"
    missing_root.mkdir()
    missing_media = missing_root / "image.png"
    missing_media.write_bytes(original_payload)
    missing_source = missing_root / "deck.md"
    missing_source.write_text(_media_markdown("image.png"), encoding="utf-8")
    manifest = load_manifest(SKILL_DIR)
    frozen_plan = pptx_paginate.build_deck_plan(parse_document(missing_source, manifest), manifest)
    before_projection = frozen_plan.to_projection()
    missing_delivery = missing_root / "delivery"

    def delete_after_parse(input_path: Path) -> None:
        del input_path
        missing_media.unlink()

    pptx_render.RENDER_POST_PARSE_HOOK = delete_after_parse
    try:
        missing_exit, missing_output = invoke_render_module(missing_source, missing_delivery, "deck")
    finally:
        pptx_render.RENDER_POST_PARSE_HOOK = None
    require(missing_exit != 0 and "MEDIA_MISSING" in missing_output and "受影响逻辑页" in missing_output,
            f"runtime missing did not report structured diagnostic: {missing_output[:1024]}")
    require(command_owned_files(missing_delivery) == set() and not (missing_delivery / ".work").exists(),
            "runtime missing published best-effort current or leaked work")
    require(frozen_plan.to_projection() == before_projection, "runtime missing mutated frozen plan projection")
    projection_before_hash = hashlib.sha256(
        json.dumps(before_projection, ensure_ascii=False, sort_keys=True).encode()
    ).hexdigest()
    projection_after_hash = hashlib.sha256(
        json.dumps(frozen_plan.to_projection(), ensure_ascii=False, sort_keys=True).encode()
    ).hexdigest()

    fault_root = workdir / "add-picture-fault"
    fault_root.mkdir()
    fault_media = fault_root / "image.png"
    fault_media.write_bytes(original_payload)
    fault_source = fault_root / "deck.md"
    fault_source.write_text(_media_markdown("image.png"), encoding="utf-8")
    fault_delivery = fault_root / "delivery"
    fault_delivery.mkdir()
    fault_old = b"old-pptx-add-picture"
    (fault_delivery / "deck.md").write_bytes(b"old-markdown")
    (fault_delivery / "deck.pptx").write_bytes(fault_old)
    fault_before_hash = hashlib.sha256(fault_old).hexdigest()

    def fail_add_picture(reference: object) -> None:
        del reference
        raise OSError("private Pillow stream failure")

    pptx_objects.MEDIA_AFTER_VALIDATE_HOOK = fail_add_picture
    try:
        fault_exit, fault_output = invoke_render_module(fault_source, fault_delivery, "deck")
    finally:
        pptx_objects.MEDIA_AFTER_VALIDATE_HOOK = None
    require(fault_exit != 0 and "PPTX_MEDIA_FORMAT_INVALID" in fault_output,
            "add_picture fault did not map to stable code")
    require(len(fault_output.encode("utf-8")) < 8192 and "Traceback" not in fault_output
            and "private Pillow stream failure" not in fault_output and str(fault_root) not in fault_output,
            "add_picture fault leaked internal details")
    require((fault_delivery / "deck.pptx").read_bytes() == fault_old
            and (fault_delivery / "deck.md").read_bytes() == b"old-markdown",
            "add_picture fault replaced old target")
    fault_after_hash = hashlib.sha256((fault_delivery / "deck.pptx").read_bytes()).hexdigest()

    return {
        "relative_success": {
            "public_exit": race_exit,
            "validated_hash": original_hash,
            "embedded_hashes": embedded_hashes,
            "replacement_hash": replacement_hash,
            "hook_called": len(hook_calls),
            "crop": [0.0, 0.0, 0.0, 0.0],
            "artifact_paths": [str(delivery / "deck.md"), str(delivery / "deck.pptx")],
        },
        "absolute_success": {
            "public_exit": absolute_result.returncode,
            "embedded_hashes": absolute_hashes,
            "crop": [0.0, 0.0, 0.0, 0.0],
            "artifact_paths": [str(absolute_delivery / "deck.md"), str(absolute_delivery / "deck.pptx")],
        },
        "symlink_failures": symlink_evidence,
        "runtime_missing": {
            "public_exit": missing_exit,
            "code": "MEDIA_MISSING",
            "best_effort_generated": "owned work evidence" in missing_output,
            "artifacts": sorted(path.name for path in missing_delivery.iterdir()),
            "projection_before_hash": projection_before_hash,
            "projection_after_hash": projection_after_hash,
            "frozen_projection_equal": projection_before_hash == projection_after_hash,
            "output_bytes": len(missing_output.encode("utf-8")),
            "artifact_paths": [],
        },
        "add_picture_fault": {
            "public_exit": fault_exit,
            "code": "PPTX_MEDIA_FORMAT_INVALID",
            "output_bytes": len(fault_output.encode("utf-8")),
            "old_target_before_hash": fault_before_hash,
            "old_target_after_hash": fault_after_hash,
            "old_target_preserved": fault_before_hash == fault_after_hash,
        },
    }


GATES["media-descriptor-binding"] = media_descriptor_binding_gate


def determinism_gate(workdir: Path) -> dict[str, object]:
    manifest = load_manifest(SKILL_DIR)
    first_document = parse_document(FIXTURE_PATH, manifest)
    second_document = parse_document(FIXTURE_PATH, manifest)
    first_plan = pptx_paginate.build_deck_plan(first_document, manifest)
    second_plan = pptx_paginate.build_deck_plan(second_document, manifest)
    first_projection = first_plan.to_projection()
    second_projection = second_plan.to_projection()
    require(first_projection == second_projection, "physical-plan projection changed")
    inventories: list[dict[str, object]] = []
    for index in range(2):
        output = workdir / f"run-{index}"
        output.mkdir()
        completed = run_public_render(FIXTURE_PATH, output, "deck")
        require(completed.returncode == 0, f"determinism render {index} failed")
        inventories.append(semantic_package_inventory(output / "deck.pptx"))
    require(inventories[0] == inventories[1], "semantic package inventory changed")
    projection_bytes = json.dumps(first_projection, ensure_ascii=False, sort_keys=True).encode("utf-8")
    return {
        "physical_plan_sha256": hashlib.sha256(projection_bytes).hexdigest(),
        "physical_slides": len(first_plan.slides),
        "semantic_inventory": inventories[0],
    }


def phase_41_42_regression_gate(workdir: Path) -> dict[str, object]:
    template = run_template_report(workdir)
    completed = subprocess.run(
        [sys.executable, str(SCRIPTS_DIR / "verify_markdown_contract.py"), "fixture-example"],
        cwd=SKILL_DIR.parent.parent, text=True, capture_output=True, timeout=120, check=False,
    )
    require(completed.returncode == 0, f"Phase 42 fixture-example failed: {completed.stdout}{completed.stderr}")
    require("Traceback" not in completed.stdout + completed.stderr, "Phase 42 regression leaked traceback")
    runtime_files = sorted(path for path in SCRIPTS_DIR.glob("*.py") if path.name.startswith("pptx_"))
    require(runtime_files and all(path.parent == SCRIPTS_DIR for path in runtime_files), "runtime file escaped skill scripts")
    forbidden_calls: list[str] = []
    for path in runtime_files:
        source = path.read_text(encoding="utf-8")
        if "skills/" in source and "school-pptx" not in source:
            forbidden_calls.append(path.name)
    require(not forbidden_calls, f"sibling skill runtime calls found: {forbidden_calls}")
    shell_source = PUBLIC_CLI.read_text(encoding="utf-8")
    require("pptx_render.py" in shell_source and "verify_school_pptx.py" in shell_source,
            "public render or verify dispatcher is missing")
    return {
        "template_layouts": len(template["layouts"]),
        "fixture_example": "PASS",
        "runtime_files": [p.name for p in runtime_files],
        "public_verify_dispatcher": True,
    }


def template_reader_security_gate(workdir: Path) -> dict[str, object]:
    completed = subprocess.run(
        [sys.executable, str(SCRIPTS_DIR / "verify_template_manifest.py")],
        cwd=SKILL_DIR.parent.parent,
        text=True,
        capture_output=True,
        timeout=180,
        check=False,
    )
    output = completed.stdout + completed.stderr
    require(completed.returncode == 0, f"template reader security regressions failed: {output[:1024]}")
    require("Traceback" not in output and str(workdir) not in output,
            "template reader security gate leaked traceback or workdir")
    require(len(output.encode("utf-8")) < 8 * 1024, "template reader security evidence exceeded 8 KiB")
    evidence = json.loads(completed.stdout)
    require(set(evidence["negative_vectors"]) == {
        "entry_count", "declared_size", "actual_size", "duplicate", "traversal",
        "doctype_entity", "external_relationship",
    }, "template reader public vector set changed")
    require(evidence["canonical"] == {"layouts": 11, "status": "PASS"},
            "canonical template reader regression changed")
    return evidence


GATES["cli-publication"] = cli_publication_gate
GATES["best-effort"] = best_effort_gate
GATES["publication-safety"] = publication_safety_gate
GATES["publication-descriptor-race"] = publication_descriptor_race_gate
GATES["table-header-only"] = table_header_only_gate
GATES["object-error-bounded"] = object_error_bounded_gate
GATES["determinism"] = determinism_gate
GATES["phase_41_42_regression"] = phase_41_42_regression_gate
GATES["template-reader-security"] = template_reader_security_gate


PHASE_43_GATE_ORDER = (
    "contract-model",
    "pagination",
    "frozen-slot-content",
    "frozen-numbering-row-heights",
    "ooxml-bootstrap",
    "editable-objects",
    "code-literal-roundtrip",
    "mixed-fragment-capacity",
    "frame-capacity-consistency",
    "emit-structure",
    "frozen-plan-emission",
    "cli-publication",
    "best-effort",
    "publication-safety",
    "publication-descriptor-race",
    "table-header-only",
    "object-error-bounded",
    "media-descriptor-binding",
    "template-reader-security",
    "determinism",
    "phase_41_42_regression",
)

PHASE_43_REQUIRED_GATES = frozenset({
    "contract-model", "pagination", "frozen-slot-content", "frozen-numbering-row-heights",
    "ooxml-bootstrap", "editable-objects", "code-literal-roundtrip", "emit-structure",
    "mixed-fragment-capacity", "frame-capacity-consistency",
    "frozen-plan-emission", "cli-publication", "best-effort", "publication-safety",
    "publication-descriptor-race", "table-header-only", "object-error-bounded", "template-reader-security",
    "media-descriptor-binding",
    "determinism", "phase_41_42_regression",
})

GAP_COVERAGE = {
    "C-01": ("publication-descriptor-race",),
    "C-02": ("code-literal-roundtrip",),
    "W-01": ("object-error-bounded",),
    "W-02": ("frozen-slot-content", "frozen-plan-emission"),
    "W-03": ("frozen-slot-content", "frozen-plan-emission"),
    "W-04": ("frozen-numbering-row-heights", "frozen-plan-emission"),
    "W-05": ("frozen-numbering-row-heights", "frozen-plan-emission"),
    "W-06": ("template-reader-security",),
    "R43-C01": ("code-literal-roundtrip",),
    "R43-C02": ("table-header-only",),
    "R43-W01": ("object-error-bounded",),
    "R43-C03": ("mixed-fragment-capacity",),
    "R43-W02": ("mixed-fragment-capacity", "media-descriptor-binding"),
    "R43-W03": ("media-descriptor-binding",),
    "R43-C04": ("frame-capacity-consistency",),
    "R43-C05": ("frame-capacity-consistency",),
    "R43-W04": ("frame-capacity-consistency",),
    "R43-W05": ("frame-capacity-consistency", "mixed-fragment-capacity", "media-descriptor-binding"),
}

REQUIREMENT_COVERAGE = {
    "PPTX-03": ("code-literal-roundtrip", "editable-objects", "frozen-slot-content", "frozen-plan-emission"),
    "PPTX-04": ("table-header-only", "editable-objects", "frozen-numbering-row-heights", "frozen-plan-emission"),
    "PPTX-05": ("media-descriptor-binding", "editable-objects", "emit-structure"),
    "PPTX-08": (
        "mixed-fragment-capacity", "frame-capacity-consistency", "pagination", "frozen-slot-content",
        "frozen-numbering-row-heights", "frozen-plan-emission",
    ),
    "PPTX-10": ("code-literal-roundtrip", "frame-capacity-consistency"),
    "VER-03": (
        "cli-publication", "best-effort", "publication-safety", "publication-descriptor-race",
        "table-header-only", "object-error-bounded", "media-descriptor-binding", "template-reader-security",
    ),
}


def _assert_gap_outcome_source_derived() -> None:
    source = Path(__file__).read_text(encoding="utf-8")

    def check(candidate: str) -> None:
        tree = ast.parse(candidate, filename=__file__)
        functions = {
            node.name: node for node in tree.body
            if isinstance(node, ast.FunctionDef)
            and node.name in {
                "mixed_fragment_capacity_gate", "media_descriptor_binding_gate",
                "frame_capacity_consistency_gate", "run_phase_43",
            }
        }
        require(len(functions) == 4, "source-derived producer set changed")
        targets = {
            "joined_equality", "embedded_hash", "embedded_hashes",
            "frozen_projection_equal", "old_target_preserved",
        }
        for function in functions.values():
            for mapping in (node for node in ast.walk(function) if isinstance(node, ast.Dict)):
                for key, value in zip(mapping.keys, mapping.values):
                    if not isinstance(key, ast.Constant) or key.value not in targets:
                        continue
                    require(not (isinstance(value, ast.Constant) and value.value is True),
                            f"{key.value} contains a hard-coded success boolean")
                    if key.value in {"embedded_hash", "embedded_hashes"}:
                        require(not (isinstance(value, ast.Name) and value.id in {"original_hash", "expected_hash"}),
                                "embedded hash is copied from expected input")
        frame_function = functions["frame_capacity_consistency_gate"]
        forbidden_frame_fields = {
            "paragraph_hash", "full_height_points", "effective_height_points",
            "capacity_equal", "projection_equal",
        }
        returned_keys = {
            key.value
            for mapping in ast.walk(frame_function) if isinstance(mapping, ast.Dict)
            for key in mapping.keys if isinstance(key, ast.Constant) and isinstance(key.value, str)
        }
        require(not forbidden_frame_fields & returned_keys,
                "frame producer returned derived capacity or equality evidence")
        fail_closed_values = [
            value for mapping in ast.walk(frame_function) if isinstance(mapping, ast.Dict)
            for key, value in zip(mapping.keys, mapping.values)
            if isinstance(key, ast.Constant) and key.value == "fail_closed_codes"
        ]
        require(len(fail_closed_values) == 1 and isinstance(fail_closed_values[0], ast.Name),
                "frame fail-closed evidence is not the captured exception list")
        run_function = functions["run_phase_43"]
        assignments = [
            node for node in ast.walk(run_function)
            if isinstance(node, ast.Assign)
            and any("gap_outcome_audit" in ast.unparse(target) for target in node.targets)
        ]
        require(len(assignments) == 1, "gap outcome audit assignment changed")
        require(not any(isinstance(node, ast.Constant) and node.value is True
                        for node in ast.walk(assignments[0].value)),
                "gap outcome audit contains a hard-coded success boolean")

    check(source)
    mutations = (
        source.replace('"joined_equality": parser_text == frozen_text == reopened_text',
                       '"joined_equality": True', 1),
        source.replace('"embedded_hashes": embedded_hashes', '"embedded_hashes": original_hash', 1),
        source.replace('"frozen_projection_equal": projection_before_hash == projection_after_hash',
                       '"frozen_projection_equal": True', 1),
        source.replace('"old_target_preserved": old_target_before_hash == old_target_after_hash',
                       '"old_target_preserved": True', 1),
        source.replace('"fail_closed_codes": fail_closed_codes',
                       '"fail_closed_codes": ["PPTX_PLAN_FONT_MISSING", "PPTX_PLAN_TYPOGRAPHY_MISSING"]', 1),
        source.replace('"frames": frames,', '"frames": frames, "full_height_points": 1.0,', 1),
    )
    for mutation in mutations:
        rejected = False
        try:
            check(mutation)
        except GateFailure:
            rejected = True
        require(rejected, "source-derived guard accepted a constant evidence mutation")


def _recompute_frame_page(page: dict[str, object]) -> dict[str, object]:
    plan_paragraphs = page["plan_paragraphs"]
    frames = page["frames"]
    plan_sequence = [(item["role"], item["visible_text"]) for item in plan_paragraphs]
    reopen_sequence: list[tuple[str, str]] = []
    frame_heights: list[tuple[float, float]] = []
    for frame in frames:
        geometry = frame["geometry_emu"]
        margins = frame["margins_points"]
        effective_width_emu = geometry["width"] - round(
            (margins["left"] + margins["right"]) * pptx_paginate.EMU_PER_POINT
        )
        effective_height = geometry["height"] / pptx_paginate.EMU_PER_POINT
        effective_height -= margins["top"] + margins["bottom"]
        require(effective_width_emu > 0 and effective_height > 0,
                "frame raw geometry is not independently measurable")
        full_height = 0.0
        for paragraph in frame["paragraphs"]:
            measurement = pptx_paginate.TextMeasure(
                line_spacing=paragraph["line_spacing"],
                paragraph_spacing_points=paragraph["space_after_points"],
            ).measure(
                paragraph["text"],
                width_emu=effective_width_emu,
                font_size=paragraph["font_size"],
                font_size_min=paragraph["font_size"],
                font_size_max=paragraph["font_size"],
                kind="code" if paragraph["role"] == "code" else "paragraph",
            )
            require(measurement.display_lines == paragraph["display_lines"],
                    "captured display lines differ from raw text remeasurement")
            full_height += (
                measurement.display_lines * paragraph["font_size"] * paragraph["line_spacing"]
                + paragraph["space_after_points"]
            )
            reopen_sequence.append((paragraph["role"], paragraph["text"]))
        frame_heights.append((round(full_height, 3), round(effective_height, 3)))
    raw_plan = json.dumps(plan_sequence, ensure_ascii=False, separators=(",", ":")).encode()
    raw_reopen = json.dumps(reopen_sequence, ensure_ascii=False, separators=(",", ":")).encode()
    return {
        "plan_hash": hashlib.sha256(raw_plan).hexdigest(),
        "reopen_hash": hashlib.sha256(raw_reopen).hexdigest(),
        "frame_heights": frame_heights,
        "sequence_equal": plan_sequence == reopen_sequence,
    }


def run_named_gate(name: str, workdir: Path) -> dict[str, object]:
    if name == "contract-model":
        return run_contract_model()
    if name == "pagination":
        names = [candidate for candidate in GATES if candidate.startswith("pagination-")]
        evidence: dict[str, object] = {}
        for candidate in names:
            evidence[candidate] = GATES[candidate](workdir)
            print(f"PASS {candidate}")
        return evidence
    evidence = GATES[name](workdir)
    print(f"PASS {name}")
    return {name: evidence}


def run_phase_43() -> dict[str, object]:
    require(len(PHASE_43_GATE_ORDER) == len(set(PHASE_43_GATE_ORDER)), "phase-43 registry contains duplicates")
    require(set(PHASE_43_GATE_ORDER) == PHASE_43_REQUIRED_GATES, "phase-43 registry coverage changed")
    evidence: dict[str, object] = {}
    called: list[str] = []
    with tempfile.TemporaryDirectory(prefix="school-pptx-phase-43-") as temporary:
        root = Path(temporary)
        for name in PHASE_43_GATE_ORDER:
            gate_dir = root / name
            gate_dir.mkdir()
            evidence[name] = run_named_gate(name, gate_dir)
            called.append(name)
    require(tuple(called) == PHASE_43_GATE_ORDER, "phase-43 registry skipped or reordered a gate")
    called_set = set(called)
    require(set(GAP_COVERAGE) == {
        "C-01", "C-02", "W-01", "W-02", "W-03", "W-04", "W-05", "W-06",
        "R43-C01", "R43-C02", "R43-W01", "R43-C03", "R43-W02", "R43-W03",
        "R43-C04", "R43-C05", "R43-W04", "R43-W05",
    },
            "phase-43 gap coverage set changed")
    require(all(gates and set(gates) <= called_set for gates in GAP_COVERAGE.values()),
            "phase-43 gap coverage names an uncalled gate")
    gap_calls = {
        gap: tuple(gate for gate in gates if gate in called_set)
        for gap, gates in GAP_COVERAGE.items()
    }
    require(gap_calls == GAP_COVERAGE, "phase-43 gap evidence differs from actual called gates")
    require(set(REQUIREMENT_COVERAGE) == {
        "PPTX-03", "PPTX-04", "PPTX-05", "PPTX-08", "PPTX-10", "VER-03",
    },
            "blocked requirement coverage set changed")
    require(all(gates and set(gates) <= called_set for gates in REQUIREMENT_COVERAGE.values()),
            "blocked requirement coverage names an uncalled gate")
    evidence["decision_coverage"] = {
        "D-01..D-05": "pagination",
        "D-06..D-10": "pagination",
        "D-11..D-16": "editable-objects,emit-structure",
        "D-17..D-18": "best-effort",
        "D-19": "publication-safety",
        "D-20..D-21": "cli-publication,publication-safety",
    }
    evidence["gap_coverage"] = GAP_COVERAGE
    evidence["gap_calls"] = gap_calls
    evidence["requirement_coverage"] = REQUIREMENT_COVERAGE
    mixed = evidence["mixed-fragment-capacity"]["mixed-fragment-capacity"]
    frame = evidence["frame-capacity-consistency"]["frame-capacity-consistency"]
    media = evidence["media-descriptor-binding"]["media-descriptor-binding"]
    mixed_vectors = mixed["vectors"]
    for layout, expected_font in (("title-content", 28.0), ("two-column", 22.0)):
        vector = mixed_vectors[layout]
        require(vector["public_exit"] == 0 and vector["physical_pages"] > 1,
                f"{layout} mixed capacity public result failed")
        require(vector["selected_font_sizes"] and all(
            value == expected_font for value in vector["selected_font_sizes"]
        ), f"{layout} mixed capacity font evidence failed")
        require(vector["frozen_typography"] == vector["reopened_typography"],
                f"{layout} typography evidence failed")
        require(all(
            display <= effective
            for display, effective in zip(
                vector["display_height_points"], vector["effective_content_height_points"], strict=True
            )
        ), f"{layout} display height evidence failed")
        require(all(vector["page_equalities"])
                and vector["parser_text_hash"] == vector["frozen_text_hash"] == vector["reopened_text_hash"]
                and len(set(vector["text_lengths"])) == 1,
                f"{layout} text equality evidence failed")
        require(vector["output_bytes"] < 8192 and len(vector["artifact_paths"]) == 2
                and Path(vector["artifact_paths"][0]).stem == Path(vector["artifact_paths"][1]).stem,
                f"{layout} output or artifact evidence failed")
    relative = media["relative_success"]
    absolute = media["absolute_success"]
    require(relative["public_exit"] == absolute["public_exit"] == 0,
            "relative or absolute media success evidence failed")
    require(relative["validated_hash"] in relative["embedded_hashes"]
            and relative["replacement_hash"] not in relative["embedded_hashes"],
            "validated media hash binding evidence failed")
    require(relative["hook_called"] > 0 and relative["crop"] == absolute["crop"] == [0.0, 0.0, 0.0, 0.0],
            "media hook or crop evidence failed")
    require(len(relative["artifact_paths"]) == len(absolute["artifact_paths"]) == 2,
            "media success artifact paths failed")
    symlinks = media["symlink_failures"]
    require(set(symlinks) == {"relative-intermediate", "relative-final", "absolute-final"},
            "media symlink vector set changed")
    require(all(
        item["public_exit"] != 0 and item["code"] == "PPTX_MEDIA_PATH_INVALID"
        and item["output_bytes"] < 8192
        and item["old_target_before_hash"] == item["old_target_after_hash"]
        for item in symlinks.values()
    ), "media symlink failure evidence failed")
    runtime_missing = media["runtime_missing"]
    require(runtime_missing["public_exit"] != 0 and runtime_missing["code"] == "MEDIA_MISSING"
            and runtime_missing["best_effort_generated"] is True
            and runtime_missing["artifacts"] == []
            and runtime_missing["projection_before_hash"] == runtime_missing["projection_after_hash"]
            and runtime_missing["output_bytes"] < 8192
            and runtime_missing["artifact_paths"] == [],
            "runtime missing evidence failed")
    add_picture_fault = media["add_picture_fault"]
    require(add_picture_fault["public_exit"] != 0
            and add_picture_fault["code"] == "PPTX_MEDIA_FORMAT_INVALID"
            and add_picture_fault["output_bytes"] < 8192
            and add_picture_fault["old_target_before_hash"] == add_picture_fault["old_target_after_hash"],
            "add_picture fault evidence failed")
    require(frame["public_exit"] == 0 and frame["output_bytes"] < 8192,
            "frame capacity public evidence failed")
    require(frame["fail_closed_codes"] == [
        "PPTX_PLAN_FONT_MISSING", "PPTX_PLAN_TYPOGRAPHY_MISSING"
    ], "frame fail-closed exceptions differ from the observed codes")
    for layout in ("contents", "title-content", "code"):
        vector = frame["vectors"][layout]
        recomputed = [_recompute_frame_page(page) for page in vector["pages"]]
        require(all(item["sequence_equal"] and item["plan_hash"] == item["reopen_hash"]
                    for item in recomputed), f"{layout} frame text projection differs")
        require(all(full <= effective + 0.001
                    for item in recomputed for full, effective in item["frame_heights"]),
                f"{layout} full-frame capacity failed")
    code_pages = frame["vectors"]["code"]["pages"]
    heading_pages = [page for page in code_pages if len(page["frames"]) == 2]
    require(heading_pages, "dedicated code heading evidence is missing")
    require(any(
        page["frames"][0]["paragraphs"][0]["display_lines"] >= 2
        for page in heading_pages
    ), "dedicated code heading vector did not exercise multiline geometry")
    for page in code_pages:
        frames = page["frames"]
        slot = page["slot_geometry_emu"]
        require(all(frame_item["geometry_emu"]["x"] == slot["x"]
                    and frame_item["geometry_emu"]["width"] == slot["width"] for frame_item in frames),
                "dedicated code frame width differs from the frozen slot")
        require(frames[0]["geometry_emu"]["y"] == slot["y"],
                "dedicated code frame start differs from the frozen slot")
        require(sum(frame_item["geometry_emu"]["height"] for frame_item in frames) == slot["height"],
                "dedicated code heading/code heights do not fill the frozen slot")
        require(all(paragraph["font_size"] == 14.0 and paragraph["line_spacing"] == 1.2
                    for frame_item in frames for paragraph in frame_item["paragraphs"]),
                "dedicated code heading/body typography drifted")
    overflow_heading = frame["overflow_heading"]
    require(overflow_heading["public_exit"] != 0 and overflow_heading["output_bytes"] < 8192,
            "overflow code heading public evidence failed")
    require(any(item["code"] == "TEXT_BLOCK_OVERFLOW" and item["severity"] == "error"
                for item in overflow_heading["diagnostics"]),
            "overflow code heading diagnostic evidence failed")
    require(overflow_heading["artifact_names"] == [],
            "overflow code heading published best-effort current artifacts")
    require(all(
        geometry["heading_height"] > 0 and geometry["code_height"] > 0
        and geometry["heading_height"] + geometry["code_height"] == geometry["slot_height"]
        for geometry in overflow_heading["geometries"]
    ), "overflow code heading raw geometry is invalid")
    _assert_gap_outcome_source_derived()
    code_outcome = evidence["code-literal-roundtrip"]["code-literal-roundtrip"]
    table_outcome = evidence["table-header-only"]["table-header-only"]
    object_outcome = evidence["object-error-bounded"]["object-error-bounded"]
    evidence["gap_outcome_audit"] = {
        "R43-C01": {"gate": "code-literal-roundtrip", "evidence": code_outcome},
        "R43-C02": {"gate": "table-header-only", "evidence": table_outcome},
        "R43-W01": {"gate": "object-error-bounded", "evidence": object_outcome},
        "R43-C03": {"gate": "mixed-fragment-capacity", "evidence": mixed_vectors},
        "R43-W02": {
            "gates": ("mixed-fragment-capacity", "media-descriptor-binding"),
            "mixed": mixed_vectors,
            "media": media,
        },
        "R43-W03": {"gate": "media-descriptor-binding", "evidence": media},
        "R43-C04": {"gate": "frame-capacity-consistency", "evidence": frame["vectors"]["contents"]},
        "R43-C05": {"gate": "frame-capacity-consistency", "evidence": frame["vectors"]["title-content"]},
        "R43-W04": {"gate": "frame-capacity-consistency", "evidence": frame["vectors"]["code"]},
        "R43-W05": {
            "gates": GAP_COVERAGE["R43-W05"],
            "frame": frame,
            "mixed": mixed_vectors,
            "media": media,
        },
    }
    evidence["registry"] = {
        "required": PHASE_43_GATE_ORDER,
        "called": tuple(called),
        "unique": len(PHASE_43_GATE_ORDER) == len(set(PHASE_43_GATE_ORDER)),
        "dynamic_skips": len(PHASE_43_GATE_ORDER) - len(called),
    }
    evidence["phase_boundary"] = {
        "public_verify_command": False,
        "runtime_adapter_changes": False,
        "manual_viewer_uat_claimed": False,
    }
    return evidence


def run_contract_model() -> dict[str, object]:
    names = ["manifest_renderer_contract_gate", "model_determinism_gate", "measurement_gate"]
    require(len(names) == len(set(names)) == 3, "gate registry must contain three unique contract gates")
    called: list[str] = []
    evidence: dict[str, object] = {}
    with tempfile.TemporaryDirectory(prefix="school-pptx-contract-model-") as temporary:
        root = Path(temporary)
        for name in names:
            gate = GATES[name]
            gate_dir = root / name
            gate_dir.mkdir()
            evidence[name] = gate(gate_dir)
            called.append(name)
            print(f"PASS {name}")
    require(called == names, "gate registry execution was skipped or reordered")
    return evidence


def parse_args(argv: list[str]) -> argparse.Namespace:
    parser = argparse.ArgumentParser(prog="verify_pptx_renderer.py")
    parser.add_argument("gate", nargs="?", choices=["contract-model", "pagination", "phase-43", *GATES])
    parser.add_argument("--self-test", choices=["delivery-transaction"])
    return parser.parse_args(argv)


def main(argv: list[str]) -> int:
    args = parse_args(argv)
    selected = args.self_test or args.gate or "all"
    try:
        if selected == "all":
            evidence = {"phase-43": run_phase_43()}
            with tempfile.TemporaryDirectory(prefix="school-pptx-delivery-transaction-") as temporary:
                evidence["delivery-transaction"] = run_named_gate("delivery-transaction", Path(temporary))
        elif selected == "phase-43":
            evidence = run_phase_43()
        else:
            with tempfile.TemporaryDirectory(prefix=f"school-pptx-{selected}-") as temporary:
                evidence = run_named_gate(selected, Path(temporary))
        print(json.dumps(evidence, ensure_ascii=False, sort_keys=True))
        return 0
    except (GateFailure, OSError, ValueError, KeyError, TypeError, json.JSONDecodeError) as exc:
        print(f"FAIL {selected}: {exc}", file=sys.stderr)
        return 1


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
