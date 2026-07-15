#!/usr/bin/env python3
"""Native editable PPTX object helpers driven only by manifest geometry."""

from __future__ import annotations

import warnings
import errno
import hashlib
import io
import os
import stat
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Callable, Iterable

from markdown_contract import inline_spans
from pptx_ooxml import set_run_highlight


MAX_MEDIA_BYTES = 16 * 1024 * 1024
MAX_MEDIA_PIXELS = 40_000_000
ALLOWED_MEDIA_FORMATS = {"JPEG", "PNG"}
MEDIA_AFTER_VALIDATE_HOOK: Callable[["MediaReference"], None] | None = None


class PptxObjectError(RuntimeError):
    pass


@dataclass(frozen=True, slots=True)
class MediaReference:
    traversal_root: Path
    path_parts: tuple[str, ...]
    absolute: bool


def normalize_rich_text(authored: str) -> tuple[tuple[str, str], ...]:
    """Remove inline delimiters once and preserve authored visible-text order."""
    spans = inline_spans(authored)
    segments: list[tuple[str, str]] = []
    cursor = 0
    for span in spans:
        start, end = int(span["start"]), int(span["end"])
        if start < cursor:
            raise PptxObjectError("PPTX_INLINE_SPAN_OVERLAP")
        if start > cursor:
            segments.append(("plain", authored[cursor:start]))
        segments.append((str(span["kind"]), str(span["text"])))
        cursor = end
    if cursor < len(authored):
        segments.append(("plain", authored[cursor:]))
    return tuple(segment for segment in segments if segment[1]) or (("plain", ""),)


def _geometry(value: dict[str, int]) -> tuple[int, int, int, int]:
    return value["x"], value["y"], value["width"], value["height"]


def _configure_text_frame(
    text_frame: Any,
    *,
    margin_left_points: float = 0.0,
    margin_right_points: float = 0.0,
    margin_top_points: float = 0.0,
    margin_bottom_points: float = 0.0,
) -> None:
    from pptx.enum.text import MSO_AUTO_SIZE
    from pptx.util import Pt

    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.margin_left = Pt(margin_left_points)
    text_frame.margin_right = Pt(margin_right_points)
    text_frame.margin_top = Pt(margin_top_points)
    text_frame.margin_bottom = Pt(margin_bottom_points)


def add_rich_text(
    slide: Any,
    geometry: dict[str, int],
    authored: str,
    *,
    font_size: float,
    highlight_scheme: str,
    name: str,
    monospace: bool = False,
    font_name: str | None = None,
    font_theme_color: str | None = None,
    margin_left_points: float = 0.0,
    margin_right_points: float = 0.0,
    margin_top_points: float = 0.0,
    margin_bottom_points: float = 0.0,
    line_spacing: float | None = None,
    paragraph_spacing_points: float | None = None,
    paragraph_alignment: str | None = None,
    vertical_anchor: str | None = None,
) -> Any:
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_THEME_COLOR
    from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
    from pptx.util import Pt

    shape = slide.shapes.add_textbox(*_geometry(geometry))
    shape.name = name
    text_frame = shape.text_frame
    _configure_text_frame(
        text_frame,
        margin_left_points=margin_left_points,
        margin_right_points=margin_right_points,
        margin_top_points=margin_top_points,
        margin_bottom_points=margin_bottom_points,
    )
    paragraph = text_frame.paragraphs[0]
    alignments = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    anchors = {
        "top": MSO_VERTICAL_ANCHOR.TOP,
        "middle": MSO_VERTICAL_ANCHOR.MIDDLE,
        "bottom": MSO_VERTICAL_ANCHOR.BOTTOM,
    }
    if paragraph_alignment:
        try:
            paragraph.alignment = alignments[paragraph_alignment]
        except KeyError as exc:
            raise PptxObjectError("PPTX_PARAGRAPH_ALIGNMENT_INVALID") from exc
    if vertical_anchor:
        try:
            text_frame.vertical_anchor = anchors[vertical_anchor]
        except KeyError as exc:
            raise PptxObjectError("PPTX_VERTICAL_ANCHOR_INVALID") from exc
    if line_spacing is not None:
        paragraph.line_spacing = line_spacing
    if paragraph_spacing_points is not None:
        paragraph.space_after = Pt(paragraph_spacing_points)
    for kind, text in normalize_rich_text(authored):
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        if font_name or monospace:
            run.font.name = font_name or "Consolas"
        if font_theme_color:
            theme_colors = {
                "background1": MSO_THEME_COLOR.BACKGROUND_1,
                "text1": MSO_THEME_COLOR.TEXT_1,
                "background2": MSO_THEME_COLOR.BACKGROUND_2,
                "text2": MSO_THEME_COLOR.TEXT_2,
            }
            try:
                run.font.color.theme_color = theme_colors[font_theme_color]
            except KeyError as exc:
                raise PptxObjectError("PPTX_FONT_THEME_COLOR_INVALID") from exc
        if kind == "bold":
            run.font.bold = True
        elif kind == "highlight":
            set_run_highlight(run, highlight_scheme)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
    return shape


def add_literal_text(
    slide: Any,
    geometry: dict[str, int],
    text: str,
    *,
    font_size: float,
    name: str,
    font_name: str = "Consolas",
    margin_left_points: float = 0.0,
    margin_right_points: float = 0.0,
    margin_top_points: float = 0.0,
    margin_bottom_points: float = 0.0,
    line_spacing: float = 1.2,
    paragraph_spacing_points: float = 2.0,
) -> Any:
    from pptx.util import Pt

    shape = slide.shapes.add_textbox(*_geometry(geometry))
    shape.name = name
    text_frame = shape.text_frame
    _configure_text_frame(
        text_frame,
        margin_left_points=margin_left_points,
        margin_right_points=margin_right_points,
        margin_top_points=margin_top_points,
        margin_bottom_points=margin_bottom_points,
    )
    paragraph = text_frame.paragraphs[0]
    paragraph.line_spacing = line_spacing
    paragraph.space_after = Pt(paragraph_spacing_points)
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.name = font_name
    return shape


def add_fragment_text_frame(
    slide: Any,
    geometry: dict[str, int],
    fragments: Iterable[Any],
    *,
    font_size: float,
    highlight_scheme: str,
    name: str,
    margin_left_points: float,
    margin_right_points: float,
    margin_top_points: float,
    margin_bottom_points: float,
    line_spacing: float,
    paragraph_spacing_points: float,
    code_font_name: str = "Consolas",
    paragraph_alignment: str | None = None,
    vertical_anchor: str | None = None,
) -> Any:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
    from pptx.util import Pt
    from pptx_model import fragment_paragraph_sequence

    shape = slide.shapes.add_textbox(*_geometry(geometry))
    shape.name = name
    text_frame = shape.text_frame
    _configure_text_frame(
        text_frame,
        margin_left_points=margin_left_points,
        margin_right_points=margin_right_points,
        margin_top_points=margin_top_points,
        margin_bottom_points=margin_bottom_points,
    )
    alignments = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    anchors = {
        "top": MSO_VERTICAL_ANCHOR.TOP,
        "middle": MSO_VERTICAL_ANCHOR.MIDDLE,
        "bottom": MSO_VERTICAL_ANCHOR.BOTTOM,
    }
    if vertical_anchor:
        try:
            text_frame.vertical_anchor = anchors[vertical_anchor]
        except KeyError as exc:
            raise PptxObjectError("PPTX_VERTICAL_ANCHOR_INVALID") from exc
    first_paragraph = True

    def append_rich_paragraph(authored: str) -> None:
        nonlocal first_paragraph
        paragraph = text_frame.paragraphs[0] if first_paragraph else text_frame.add_paragraph()
        first_paragraph = False
        paragraph.line_spacing = line_spacing
        paragraph.space_after = Pt(paragraph_spacing_points)
        if paragraph_alignment:
            try:
                paragraph.alignment = alignments[paragraph_alignment]
            except KeyError as exc:
                raise PptxObjectError("PPTX_PARAGRAPH_ALIGNMENT_INVALID") from exc
        for kind, text in normalize_rich_text(authored):
            run = paragraph.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            if kind == "bold":
                run.font.bold = True
            elif kind == "highlight":
                set_run_highlight(run, highlight_scheme)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)

    for fragment in fragments:
        for projection in fragment_paragraph_sequence(fragment):
            if projection.role == "code":
                text = projection.authored_text
                paragraph = text_frame.paragraphs[0] if first_paragraph else text_frame.add_paragraph()
                first_paragraph = False
                paragraph.line_spacing = line_spacing
                paragraph.space_after = Pt(paragraph_spacing_points)
                paragraph.alignment = PP_ALIGN.LEFT
                run = paragraph.add_run()
                run.text = text
                run.font.size = Pt(font_size)
                run.font.name = code_font_name
            else:
                append_rich_paragraph(projection.authored_text)
    return shape


def add_plain_lines(
    slide: Any, geometry: dict[str, int], lines: Iterable[str], *, font_size: float,
    highlight_scheme: str, name: str, monospace: bool = False
) -> Any:
    return add_rich_text(
        slide, geometry, "\n".join(lines), font_size=font_size, highlight_scheme=highlight_scheme,
        name=name, monospace=monospace,
    )


def add_table(
    slide: Any, slot: dict[str, Any], fragment: Any, *, font_size: float,
    row_heights_emu: tuple[int, ...],
) -> tuple[Any, Any]:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
    from pptx.util import Pt

    geometry = slot["geometry"]
    table_name_geometry = slot["subregions"]["table_name"]["geometry"]
    table_name = dict(fragment.metadata).get("table_name", "")
    name_shape = slide.shapes.add_textbox(*_geometry(table_name_geometry))
    name_shape.name = "school-pptx:table-name"
    _configure_text_frame(name_shape.text_frame)
    name_shape.text_frame.paragraphs[0].text = table_name
    if name_shape.text_frame.paragraphs[0].runs:
        name_shape.text_frame.paragraphs[0].runs[0].font.size = Pt(
            slot["subregions"]["table_name"]["text_budget"]["font_size_max"]
        )

    rows = fragment.rows or (("",),)
    if len(row_heights_emu) != len(rows) or any(not isinstance(value, int) or value <= 0 for value in row_heights_emu):
        raise PptxObjectError("PPTX_TABLE_ROW_HEIGHT_MISMATCH")
    columns = max(1, max(len(row) for row in rows))
    table_body_geometry = slot.get("subregions", {}).get("table_body", {}).get("geometry")
    if table_body_geometry is None:
        table_y = table_name_geometry["y"] + table_name_geometry["height"]
        table_body_geometry = {
            "x": geometry["x"],
            "y": table_y,
            "width": geometry["width"],
            "height": geometry["y"] + geometry["height"] - table_y,
        }
    graphic = slide.shapes.add_table(len(rows), columns, *_geometry(table_body_geometry))
    graphic.name = "school-pptx:native-table"
    table = graphic.table
    style = slot.get("table_style", {})
    header_fill = RGBColor.from_string(style.get("header_fill", "4472C4"))
    body_fill = RGBColor.from_string(style.get("body_fill", "D9EAF7"))
    header_text = RGBColor.from_string(style.get("header_text", "FFFFFF"))
    body_text = RGBColor.from_string(style.get("body_text", "1F4E79"))
    for column in table.columns:
        column.width = int(table_body_geometry["width"] / columns)
    for row, height in zip(table.rows, row_heights_emu):
        row.height = height
    for row_index, values in enumerate(rows):
        for column_index in range(columns):
            cell = table.cell(row_index, column_index)
            cell.text = values[column_index] if column_index < len(values) else ""
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if row_index == 0 else body_fill
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)
                    run.font.bold = row_index == 0
                    run.font.color.rgb = header_text if row_index == 0 else body_text
    return name_shape, graphic


def _load_validated_image(reference: MediaReference) -> tuple[bytes, int, int, str]:
    from PIL import Image, UnidentifiedImageError

    descriptors: list[int] = []
    try:
        root_fd = os.open(reference.traversal_root, os.O_RDONLY | os.O_DIRECTORY | os.O_NOFOLLOW)
        descriptors.append(root_fd)
        directory_fd = root_fd
        for part in reference.path_parts[:-1]:
            next_fd = os.open(part, os.O_RDONLY | os.O_DIRECTORY | os.O_NOFOLLOW, dir_fd=directory_fd)
            descriptors.append(next_fd)
            if not stat.S_ISDIR(os.fstat(next_fd).st_mode):
                raise PptxObjectError("PPTX_MEDIA_PATH_INVALID")
            directory_fd = next_fd
        final_fd = os.open(reference.path_parts[-1], os.O_RDONLY | os.O_NOFOLLOW, dir_fd=directory_fd)
        descriptors.append(final_fd)
        final_status = os.fstat(final_fd)
        if not stat.S_ISREG(final_status.st_mode):
            raise PptxObjectError("PPTX_MEDIA_PATH_INVALID")
        if final_status.st_size > MAX_MEDIA_BYTES:
            raise PptxObjectError("PPTX_MEDIA_SIZE_LIMIT")
        chunks: list[bytes] = []
        remaining = MAX_MEDIA_BYTES + 1
        while remaining:
            chunk = os.read(final_fd, min(1024 * 1024, remaining))
            if not chunk:
                break
            chunks.append(chunk)
            remaining -= len(chunk)
        payload = b"".join(chunks)
        if len(payload) > MAX_MEDIA_BYTES:
            raise PptxObjectError("PPTX_MEDIA_SIZE_LIMIT")
        with warnings.catch_warnings():
            warnings.simplefilter("error", Image.DecompressionBombWarning)
            with Image.open(io.BytesIO(payload)) as image:
                if image.format not in ALLOWED_MEDIA_FORMATS:
                    raise PptxObjectError("PPTX_MEDIA_FORMAT_INVALID")
                width, height = image.size
                if width <= 0 or height <= 0 or width * height > MAX_MEDIA_PIXELS:
                    raise PptxObjectError("PPTX_MEDIA_PIXEL_LIMIT")
                image.verify()
    except PptxObjectError:
        raise
    except FileNotFoundError:
        raise PptxObjectError("MEDIA_MISSING") from None
    except OSError as exc:
        if exc.errno == errno.ENOENT:
            raise PptxObjectError("MEDIA_MISSING") from None
        if exc.errno in {errno.ELOOP, errno.ENOTDIR, errno.EACCES, errno.EPERM}:
            raise PptxObjectError("PPTX_MEDIA_PATH_INVALID") from None
        raise PptxObjectError("PPTX_MEDIA_FORMAT_INVALID") from None
    except (Image.DecompressionBombError, Image.DecompressionBombWarning):
        raise PptxObjectError("PPTX_MEDIA_PIXEL_LIMIT") from None
    except (UnidentifiedImageError, OSError, SyntaxError, ValueError):
        raise PptxObjectError("PPTX_MEDIA_FORMAT_INVALID") from None
    finally:
        for descriptor in reversed(descriptors):
            try:
                os.close(descriptor)
            except OSError:
                pass
    return payload, width, height, hashlib.sha256(payload).hexdigest()


def _missing_media_placeholder(slide: Any, geometry: dict[str, int], name: str) -> Any:
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, *_geometry(geometry))
    shape.name = name.replace("picture", "missing-media")
    shape.text = "媒体缺失"
    return shape


def _apply_picture_shadow(picture: Any) -> None:
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement

    properties = picture._element.spPr
    for existing in tuple(properties):
        if existing.tag in {qn("a:effectLst"), qn("a:effectDag")}:
            properties.remove(existing)
    effects = OxmlElement("a:effectLst")
    shadow = OxmlElement("a:outerShdw")
    shadow.set("blurRad", "63500")
    shadow.set("dist", "38100")
    shadow.set("dir", "2700000")
    shadow.set("algn", "ctr")
    shadow.set("rotWithShape", "0")
    color = OxmlElement("a:srgbClr")
    color.set("val", "000000")
    alpha = OxmlElement("a:alpha")
    alpha.set("val", "25000")
    color.append(alpha)
    shadow.append(color)
    effects.append(shadow)
    properties.insert_element_before(effects, "a:scene3d", "a:sp3d", "a:extLst")


def add_contain_picture(
    slide: Any, reference: MediaReference | None, geometry: dict[str, int], *, name: str
) -> tuple[Any, bool, str | None]:
    from PIL import Image, UnidentifiedImageError

    if reference is None:
        return _missing_media_placeholder(slide, geometry, name), True, None
    try:
        payload, width, height, payload_hash = _load_validated_image(reference)
    except PptxObjectError as exc:
        if str(exc) == "MEDIA_MISSING":
            return _missing_media_placeholder(slide, geometry, name), True, None
        raise
    scale = min(geometry["width"] / width, geometry["height"] / height)
    target_width, target_height = int(width * scale), int(height * scale)
    left = geometry["x"] + int((geometry["width"] - target_width) / 2)
    top = geometry["y"] + int((geometry["height"] - target_height) / 2)
    try:
        if MEDIA_AFTER_VALIDATE_HOOK is not None:
            MEDIA_AFTER_VALIDATE_HOOK(reference)
        picture = slide.shapes.add_picture(io.BytesIO(payload), left, top, target_width, target_height)
    except (Image.DecompressionBombError, Image.DecompressionBombWarning):
        raise PptxObjectError("PPTX_MEDIA_PIXEL_LIMIT") from None
    except (UnidentifiedImageError, OSError, SyntaxError, ValueError):
        raise PptxObjectError("PPTX_MEDIA_FORMAT_INVALID") from None
    except Exception:
        raise PptxObjectError("PPTX_OBJECT_INVALID") from None
    picture.name = name
    picture.crop_left = picture.crop_top = picture.crop_right = picture.crop_bottom = 0
    _apply_picture_shadow(picture)
    return picture, False, payload_hash


def add_gallery_card(
    slide: Any, preset: dict[str, Any], reference: MediaReference | None, caption: str, *, index: int,
    highlight_scheme: str, font_size: float, paragraph_alignment: str | None = None,
) -> tuple[Any, bool, str | None]:
    picture, missing, payload_hash = add_contain_picture(
        slide, reference, preset["picture"], name=f"school-pptx:gallery-picture:{index}"
    )
    caption_shape = add_rich_text(
        slide, preset["caption"], caption, font_size=font_size, highlight_scheme=highlight_scheme,
        name=f"school-pptx:gallery-caption:{index}", paragraph_alignment=paragraph_alignment,
    )
    group = slide.shapes.add_group_shape((picture, caption_shape))
    group.name = f"school-pptx:gallery-card:{index}"
    return group, missing, payload_hash


def _set_gradient_fill(shape: Any, start_scheme: str, end_scheme: str) -> None:
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement

    properties = shape._element.spPr
    fill_tags = {qn(name) for name in (
        "a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill"
    )}
    for child in tuple(properties):
        if child.tag in fill_tags:
            properties.remove(child)
    gradient = OxmlElement("a:gradFill")
    gradient.set("rotWithShape", "1")
    stops = OxmlElement("a:gsLst")
    for position, scheme in (("0", start_scheme), ("100000", end_scheme)):
        stop = OxmlElement("a:gs")
        stop.set("pos", position)
        color = OxmlElement("a:schemeClr")
        color.set("val", scheme)
        stop.append(color)
        stops.append(stop)
    gradient.append(stops)
    direction = OxmlElement("a:lin")
    direction.set("ang", "10800000")
    direction.set("scaled", "1")
    gradient.append(direction)
    insertion_index = next(
        (index for index, child in enumerate(properties) if child.tag in {
            qn("a:ln"), qn("a:effectLst"), qn("a:effectDag"), qn("a:scene3d"), qn("a:sp3d"), qn("a:extLst")
        }),
        len(properties),
    )
    properties.insert(insertion_index, gradient)
    shape.line.fill.background()


def _interpolate_rgb(start: str, end: str, position: float) -> tuple[int, int, int]:
    ratio = min(1.0, max(0.0, position))
    start_channels = tuple(int(start[index:index + 2], 16) for index in (0, 2, 4))
    end_channels = tuple(int(end[index:index + 2], 16) for index in (0, 2, 4))
    return tuple(round(left + (right - left) * ratio) for left, right in zip(start_channels, end_channels))


def add_timeline(slide: Any, slot: dict[str, Any], rows: tuple[tuple[str, ...], ...], *, highlight_scheme: str) -> list[Any]:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    axis_geometry = slot["subregions"]["axis"]["geometry"]
    gradient = slot["subregions"]["axis"].get("gradient", {})
    axis = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, *_geometry(axis_geometry))
    axis.name = "school-pptx:timeline-axis"
    _set_gradient_fill(
        axis,
        str(gradient.get("start_scheme", "accent6")),
        str(gradient.get("end_scheme", "accent1")),
    )
    results: list[Any] = [axis]
    node_band = slot["subregions"]["node_band"]["geometry"]
    template = slot["node_template"]
    count = max(1, len(rows))
    node_width = int(slot.get("node_width_emu") or node_band["width"] / count)
    node_step = (node_band["width"] - node_width) / max(1, count - 1)
    for index, row in enumerate(rows):
        base_x = node_band["x"] + round(index * node_step)
        shapes: list[Any] = []
        for role, value in zip(("time", "title", "description"), (*row, "", "")[:3]):
            local = template["subregions"][role]["geometry"]
            geometry = {
                "x": base_x + int(local["x"] * node_width / template["geometry"]["width"]),
                "y": node_band["y"] + local["y"],
                "width": max(1, int(local["width"] * node_width / template["geometry"]["width"])),
                "height": local["height"],
            }
            budget = template["subregions"][role]["text_budget"]
            region = template["subregions"][role]
            shapes.append(add_rich_text(
                slide, geometry, value, font_size=budget["font_size_min"], highlight_scheme=highlight_scheme,
                name=f"school-pptx:timeline-{role}:{index}",
                paragraph_alignment=region.get("paragraph_alignment"),
                vertical_anchor=region.get("vertical_anchor"),
            ))
        marker_local = template["subregions"]["marker"]["geometry"]
        marker_size = min(int(marker_local["width"]), int(marker_local["height"]))
        marker = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            base_x + int(marker_local["x"] * node_width / template["geometry"]["width"])
            + (int(marker_local["width"]) - marker_size) // 2,
            node_band["y"] + marker_local["y"]
            + (int(marker_local["height"]) - marker_size) // 2,
            marker_size, marker_size,
        )
        marker.name = f"school-pptx:timeline-marker:{index}"
        marker_center = marker.left + marker.width / 2
        marker_position = (marker_center - axis_geometry["x"]) / axis_geometry["width"]
        marker.fill.solid()
        marker.fill.fore_color.rgb = RGBColor(*_interpolate_rgb(
            str(gradient.get("start_rgb", "70AD47")),
            str(gradient.get("end_rgb", "4472C4")),
            marker_position,
        ))
        marker.line.fill.background()
        shapes.append(marker)
        group = slide.shapes.add_group_shape(shapes)
        group.name = f"school-pptx:timeline-node:{index}"
        if slot.get("node_group_width_emu"):
            group.width = int(slot["node_group_width_emu"])
        if slot.get("node_group_height_emu"):
            group.height = int(slot["node_group_height_emu"])
        results.append(group)
    return results


def set_notes(slide: Any, notes: str | None) -> None:
    if notes is None:
        return
    text_frame = slide.notes_slide.notes_text_frame
    if text_frame is None:
        raise PptxObjectError("PPTX_NOTES_FRAME_MISSING")
    text_frame.text = notes
